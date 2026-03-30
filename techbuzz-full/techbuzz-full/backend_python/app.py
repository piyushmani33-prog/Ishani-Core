"""
TechBuzz AI Agent - Leazy Jinn Empire Backend
=============================================
Extends the original TechBuzz backend with a living Leazy Jinn interface,
Praapti recruitment, Swarm missions, Nirmaan proposals, Akshaya memory,
voice wake support, and PWA delivery.
"""

import asyncio
from collections import Counter
import hashlib
import io
import json
import logging
import mimetypes
import os
import re
import secrets
import shutil
import sqlite3
import threading
import time
import uuid
import zipfile
from csv import reader as csv_reader
from contextlib import asynccontextmanager
from datetime import UTC, datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import quote

import httpx
from dotenv import dotenv_values, load_dotenv, set_key, unset_key
from fastapi import FastAPI, HTTPException, Request, Header, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, RedirectResponse, Response, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from empire_merge_layer import install_empire_merge_layer
from global_recruitment_brain_layer import install_global_recruitment_brain_layer
from recruitment_brain_layer import install_recruitment_brain_layer
from browser_suite_layer import install_browser_suite_layer
from brain_communication_layer import install_brain_communication_layer
from interpreter_brain_layer import install_interpreter_brain_layer
from orchestration_stack_layer import install_orchestration_stack_layer
from local_ai_runtime_layer import install_local_ai_runtime_layer
from voice_runtime_layer import install_voice_runtime_layer
from recruiter_status_layer import register_recruiter_status_routes

try:
    from pypdf import PdfReader, PdfWriter
except Exception:
    PdfReader = None
    PdfWriter = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


BASE_DIR = Path(__file__).resolve().parent.parent
BACKEND_DIR = Path(__file__).resolve().parent
FRONTEND_DIR = BASE_DIR / "frontend"
DATA_DIR = BACKEND_DIR / "data"
PROPOSAL_DIR = DATA_DIR / "nirmaan"
DOCUMENT_DIR = DATA_DIR / "documents"
DOCUMENT_EXPORT_DIR = DATA_DIR / "document_exports"
STATE_PATH = DATA_DIR / "empire_state.json"
DB_PATH = DATA_DIR / "ishani_core.db"
ENV_PATH = BACKEND_DIR / ".env"

for directory in (DATA_DIR, PROPOSAL_DIR, DOCUMENT_DIR, DOCUMENT_EXPORT_DIR):
    directory.mkdir(parents=True, exist_ok=True)

load_dotenv(BASE_DIR / ".env")
load_dotenv(ENV_PATH)


logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
log = logging.getLogger("techbuzz")


def default_naukri_launcher_path() -> Path:
    configured = os.getenv("NAUKRI_LAUNCHER_PATH", "").strip()
    if configured:
        return Path(configured).expanduser()
    return Path.home() / "Downloads" / "Naukri Launcher Installer.exe"


ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "").strip()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "").strip()
OLLAMA_HOST = os.getenv("OLLAMA_HOST", "http://127.0.0.1:11434").strip() or "http://127.0.0.1:11434"
ADMIN_EMAIL = os.getenv("ADMIN_EMAIL", "master@local").strip()
MASTER_ACCOUNT_EMAIL = os.getenv("MASTER_ACCOUNT_EMAIL", "owner@local").strip() or "owner@local"
MASTER_LOGIN_ID = os.getenv("MASTER_LOGIN_ID", "MasterPiyushMani").strip() or "MasterPiyushMani"
MASTER_PASSWORD_HASH = os.getenv("MASTER_PASSWORD_HASH", "").strip()
MASTER_PASSWORD_SALT = os.getenv("MASTER_PASSWORD_SALT", "").strip()
ALLOWED_ORIGINS = [origin.strip() for origin in os.getenv("ALLOWED_ORIGINS", "*").split(",") if origin.strip()]
MODEL = os.getenv("ANTHROPIC_MODEL", "claude-sonnet-4-20250514").strip() or "claude-sonnet-4-20250514"
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-5.1-chat-latest").strip() or "gpt-5.1-chat-latest"
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.0-flash").strip() or "gemini-2.0-flash"
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "llama3:latest").strip() or "llama3:latest"
MAX_TOKENS = int(os.getenv("MAX_TOKENS", "2048") or "2048")
RATE_LIMIT_PER_MIN = int(os.getenv("RATE_LIMIT_PER_MIN", "30") or "30")
AI_NAME = os.getenv("LEAZY_AI_NAME", "Leazy Jinn").strip() or "Leazy Jinn"
COMPANY_NAME = os.getenv("LEAZY_COMPANY_NAME", "TechBuzz Systems").strip() or "TechBuzz Systems"
CORE_IDENTITY = os.getenv("LEAZY_CORE_IDENTITY", "Ishani").strip() or "Ishani"
CREATOR_MODE_LABEL = "Creator Priority Mode"
SESSION_COOKIE_NAME = "ishani_session"
SESSION_TTL_DAYS = 14
NAUKRI_LAUNCHER_PATH = default_naukri_launcher_path()


STATE_LOCK = threading.Lock()
PROVIDER_STATUS_LOCK = threading.Lock()
CABINET_LOOP_STOP = threading.Event()
CABINET_LOOP_THREAD: Optional[threading.Thread] = None
_rate_store: Dict[str, List[float]] = {}
_provider_status: Dict[str, Any] = {"errors": {}, "cooldowns": {}}
_provider_catalog_cache: Dict[str, Dict[str, Any]] = {}
RECRUITMENT_LAYER: Dict[str, Any] = {}
GLOBAL_RECRUITMENT_LAYER: Dict[str, Any] = {}
INTERPRETER_LAYER: Dict[str, Any] = {}
ORCHESTRATION_STACK_LAYER: Dict[str, Any] = {}
LOCAL_AI_STACK_LAYER: Dict[str, Any] = {}
VOICE_RUNTIME_LAYER: Dict[str, Any] = {}
_ollama_warm_state: Dict[str, Any] = {"ready": False, "warming": False, "last_attempt_at": "", "last_error": ""}
SWARM_AGENT_NAMES = [
    ("Praapti", "Talent Hunt", "🎯"),
    ("Revenue", "Growth Engine", "💰"),
    ("Army", "Execution Core", "⚔️"),
    ("Navy", "Distribution Mesh", "⛵"),
    ("Airforce", "Signal Scanner", "✈️"),
    ("Research", "Nirmaan Lab", "🧪"),
    ("Armory", "Automation Forge", "🛡️"),
    ("Warehouse", "Delivery Hub", "🏬"),
    ("Storage", "Akshaya Vault", "📦"),
    ("Sewage", "Cleanup Stream", "🚰"),
    ("Maya", "Narrative Engine", "🌐"),
    ("Anveshan", "Discovery Lens", "🔬"),
]

DEFAULT_PROVIDER_MODEL_OPTIONS: Dict[str, List[str]] = {
    "ollama": ["llama3:latest", "mistral:latest", "phi4:latest"],
    "anthropic": ["claude-sonnet-4-20250514", "claude-opus-4-20250514"],
    "openai": ["gpt-5.1-chat-latest", "gpt-5.1", "gpt-5", "gpt-5-mini", "gpt-4.1", "gpt-4.1-mini", "gpt-4o-mini"],
    "gemini": ["gemini-2.0-flash", "gemini-1.5-pro"],
    "built_in": ["empire-fallback"],
}

VOICE_PROFILE_PRESETS: Dict[str, Dict[str, Any]] = {
    "sovereign_female": {
        "label": "Sovereign Female",
        "language": "en-IN",
        "rate": 0.94,
        "pitch": 1.08,
        "style": "calm, sovereign, feminine",
    },
    "warm_guide": {
        "label": "Warm Guide",
        "language": "en-IN",
        "rate": 0.96,
        "pitch": 1.03,
        "style": "warm, clear, reassuring",
    },
    "strategic_female": {
        "label": "Strategic Female",
        "language": "en-US",
        "rate": 0.92,
        "pitch": 1.0,
        "style": "precise, strategic, executive",
    },
    "neutral": {
        "label": "Neutral",
        "language": "en-IN",
        "rate": 1.0,
        "pitch": 1.0,
        "style": "balanced and direct",
    },
}

MARKET_READY_LESSONS: List[Dict[str, Any]] = [
    {
        "key": "recruitment_precision",
        "title": "Recruitment Precision Charter",
        "summary": "A recruitment company must only reveal candidate context that the operator explicitly asked to see, and every shortlist should map to role, location, compensation, and urgency.",
        "keywords": ["recruitment", "privacy", "shortlist", "precision", "ats"],
    },
    {
        "key": "client_delivery",
        "title": "Client Delivery Operating Rhythm",
        "summary": "Each client lane should move through discovery, JD shaping, sourcing, screening, interview scheduling, offer support, and post-join follow-through with clear ownership.",
        "keywords": ["client delivery", "recruitment ops", "service quality", "execution"],
    },
    {
        "key": "sourcing_discipline",
        "title": "Sourcing Discipline",
        "summary": "The system should gather signals from operator-approved research, preserve what matters, and avoid dumping broad data when the user asked for a precise candidate or market slice.",
        "keywords": ["sourcing", "research", "candidate search", "operator control"],
    },
    {
        "key": "communication_handoffs",
        "title": "Communication Handoffs",
        "summary": "Emails, calls, browser launches, WhatsApp, Teams, and other actions should be prepared visibly and never sent silently. The operator remains in control of the final send.",
        "keywords": ["email", "teams", "calls", "handoff", "operator control"],
    },
    {
        "key": "software_delivery",
        "title": "Software Delivery Spine",
        "summary": "Planning, coding, interface design, testing, rollout, support, and enhancement should be treated as one connected delivery chain with visible state transitions.",
        "keywords": ["planning", "coding", "testing", "deployment", "support"],
    },
]

NAVIGATOR_QUICK_SITES: List[Dict[str, str]] = [
    {
        "id": "naukri",
        "label": "Naukri",
        "url": "https://www.naukri.com/",
        "kind": "sourcing",
        "workspace_tip": "Use Open External or the local launcher, then capture only the approved role and candidate notes back into Ishani.",
        "learning_focus": "Role variants, sourcing string quality, response patterns, duplicates, and candidate objections.",
    },
    {
        "id": "linkedin",
        "label": "LinkedIn",
        "url": "https://www.linkedin.com/",
        "kind": "network",
        "workspace_tip": "Use Open External for profile and company research, then preserve only the relationship or hiring signal you actually need.",
        "learning_focus": "Company movement, talent clusters, title variants, and hiring-manager language.",
    },
    {
        "id": "gmail",
        "label": "Gmail",
        "url": "https://mail.google.com/",
        "kind": "communication",
        "workspace_tip": "Prepare drafts in Action Center first, then open Gmail in a normal tab for the final send.",
        "learning_focus": "Acknowledgment flow, candidate follow-up quality, and message timing.",
    },
    {
        "id": "teams",
        "label": "Microsoft Teams",
        "url": "https://teams.microsoft.com/",
        "kind": "communication",
        "workspace_tip": "Use Open External for real chat windows and preserve the meeting, handoff, or follow-up summary here.",
        "learning_focus": "Interview coordination, hiring-manager updates, and delivery handoffs.",
    },
    {
        "id": "github",
        "label": "GitHub",
        "url": "https://github.com/",
        "kind": "research",
        "workspace_tip": "Use Open External for repositories and issues, then capture the relevant technical signal into the learning panel.",
        "learning_focus": "Evidence of execution, code depth, portfolio proof, and hiring signals for technical candidates.",
    },
    {
        "id": "techbuzz",
        "label": "TechBuzz HQ",
        "url": "/company/portal",
        "kind": "internal",
        "workspace_tip": "This internal surface can open directly inside the Navigator workspace.",
        "learning_focus": "Workflow routing, service discovery, and company system context.",
    },
]

ACCOUNT_REGION_PRESETS: Dict[str, Dict[str, Any]] = {
    "IN": {
        "code": "IN",
        "label": "India Ledger",
        "currency": "INR",
        "tax_name": "GST",
        "suggested_tax_rate": 18.0,
        "filing_cycle": "monthly_or_quarterly",
        "notes": "Designed for India-first service businesses. Edit tax rate, filing cycle, and registration details to match your exact setup.",
    },
    "AE": {
        "code": "AE",
        "label": "UAE VAT Ledger",
        "currency": "AED",
        "tax_name": "VAT",
        "suggested_tax_rate": 5.0,
        "filing_cycle": "quarterly",
        "notes": "Good for UAE-style VAT tracking. Adjust rates and reporting practice to match your accountant's advice.",
    },
    "GLOBAL": {
        "code": "GLOBAL",
        "label": "Global Custom Ledger",
        "currency": "USD",
        "tax_name": "Tax",
        "suggested_tax_rate": 0.0,
        "filing_cycle": "custom",
        "notes": "A neutral preset for international or custom accounting workflows. Configure all tax values manually.",
    },
}

ACCOUNT_INFLOW_TYPES = {"income", "sale", "invoice", "receipt"}
ACCOUNT_OUTFLOW_TYPES = {"expense", "purchase", "cost", "payment", "salary", "tax_payment"}
ACCOUNT_ENTRY_TYPE_ALIASES = {
    "sales": "sale",
    "receipts": "receipt",
    "payments": "payment",
    "expenses": "expense",
    "payout": "payment",
    "payroll": "salary",
    "tax": "tax_payment",
}
ACCOUNT_GST_CATEGORY_RATES = {
    "recruitment_fee": 18.0,
    "recruitment_consulting": 18.0,
    "training": 18.0,
    "software": 18.0,
    "services": 18.0,
    "goods": 12.0,
    "expense": 18.0,
    "income": 18.0,
    "payroll": 0.0,
    "tax": 0.0,
}
ACCOUNT_TDS_SECTION_RATES = {
    "194J": 10.0,
    "194C": 2.0,
    "194H": 5.0,
    "194I": 10.0,
}
ACCOUNT_TDS_THRESHOLDS = {
    "194J": 30000.0,
    "194C": 30000.0,
    "194H": 15000.0,
    "194I": 30000.0,
}

OPERATIONAL_DOMAIN_BLUEPRINTS: List[Dict[str, str]] = [
    {
        "id": "management",
        "name": "Management Command",
        "purpose": "Governance, planning, revenue direction, and cross-system coordination.",
        "lead": "Prime Minister Cabinet",
    },
    {
        "id": "accounts",
        "name": "Accounts Command",
        "purpose": "Regional tax guidance, ledger automation, financial discipline, and cashflow clarity.",
        "lead": "Accounts Secretary",
    },
    {
        "id": "hospital",
        "name": "Arogya Hospital",
        "purpose": "Recovery, incident care, wellbeing, and service stabilization.",
        "lead": "Recovery Secretary",
    },
    {
        "id": "firebrigade",
        "name": "Fire Brigade",
        "purpose": "Urgent response, escalation handling, and runtime threat suppression.",
        "lead": "Protection Secretary",
    },
    {
        "id": "mess",
        "name": "Mess And Supply Kitchen",
        "purpose": "Nourishment of teams, daily operations, and readiness support.",
        "lead": "Operations Steward",
    },
    {
        "id": "lab",
        "name": "Research Lab",
        "purpose": "Nirmaan proposals, experiments, and knowledge discovery.",
        "lead": "Research Secretary",
    },
    {
        "id": "army",
        "name": "Army Command",
        "purpose": "Execution, deployment pressure, and mission discipline.",
        "lead": "Army Agent",
    },
    {
        "id": "navy",
        "name": "Navy Command",
        "purpose": "Distribution lanes, delivery flow, and multi-channel movement.",
        "lead": "Navy Agent",
    },
    {
        "id": "air_force",
        "name": "Air Force Command",
        "purpose": "Signal scanning, rapid response, and aerial situational awareness.",
        "lead": "Airforce Agent",
    },
    {
        "id": "development",
        "name": "Development Forge",
        "purpose": "Product execution, implementation, and operator tooling.",
        "lead": "Platform Secretary",
    },
    {
        "id": "evolving",
        "name": "Evolution Chamber",
        "purpose": "Continuous learning, refinement, and next-cycle adaptation.",
        "lead": "Nirmaan Chakra",
    },
    {
        "id": "upgrade",
        "name": "Upgrade Rail",
        "purpose": "Safe upgrades, staged improvements, and controlled rollout.",
        "lead": "Upgrade Steward",
    },
    {
        "id": "supply_chain",
        "name": "Supply Chain",
        "purpose": "Resource routing, handoffs, delivery continuity, and backlog control.",
        "lead": "Delivery Secretary",
    },
    {
        "id": "transmission",
        "name": "Transmission Grid",
        "purpose": "Message relay, engine handoff, and report propagation.",
        "lead": "Transmission Steward",
    },
    {
        "id": "connectivity",
        "name": "Connectivity Mesh",
        "purpose": "Client links, ecosystem connections, and workflow continuity.",
        "lead": "Connectivity Steward",
    },
]

DASHAVATARA: Dict[str, Dict[str, Any]] = {
    "MATSYA": {
        "emoji": "🐟",
        "color": "#38bdf8",
        "power": "Preserve all knowledge and recover what was almost lost.",
        "desc": "Eternal preservation, auto-backup, and recovery across the empire.",
        "mantra": "Protect the memory, save the future.",
        "keywords": ["backup", "restore", "recover", "preserve", "memory", "vault", "archive"],
    },
    "KURMA": {
        "emoji": "🐢",
        "color": "#22c55e",
        "power": "Hold the weight of scale and keep heavy systems stable.",
        "desc": "System support, resilience, scaling, and long-duration stability.",
        "mantra": "Hold steady while the universe churns.",
        "keywords": ["scale", "stability", "support", "infrastructure", "load", "reliability"],
    },
    "VARAHA": {
        "emoji": "🐗",
        "color": "#eab308",
        "power": "Lift sinking missions out of crisis and bring them back to life.",
        "desc": "Project rescue, crisis recovery, and revival of stuck work.",
        "mantra": "Lift the fallen back into the light.",
        "keywords": ["crisis", "rescue", "save project", "stuck", "failing", "recovery"],
    },
    "NARASIMHA": {
        "emoji": "🦁",
        "color": "#ef4444",
        "power": "Destroy threats, bugs, and hostile inefficiencies with force.",
        "desc": "Threat response, bug smashing, and defensive aggression.",
        "mantra": "Cut down the threat without fear.",
        "keywords": ["bug", "threat", "attack", "security", "proxy", "exploit", "malware"],
    },
    "VAMANA": {
        "emoji": "🪄",
        "color": "#8b5cf6",
        "power": "Shrink giant complexity into simple, elegant steps.",
        "desc": "Problem reduction, simplification, and elegant decomposition.",
        "mantra": "Three small steps can conquer a universe.",
        "keywords": ["simplify", "small", "reduce", "break down", "decompose", "minimal"],
    },
    "PARASHURAMA": {
        "emoji": "🪓",
        "color": "#f97316",
        "power": "Cut dead code, remove rot, and clear corruption.",
        "desc": "Cleanup, pruning, refactoring, and forceful removal of waste.",
        "mantra": "A clean kingdom grows faster.",
        "keywords": ["cleanup", "refactor", "remove", "delete", "dead code", "obsolete"],
    },
    "RAMA": {
        "emoji": "👑",
        "color": "#facc15",
        "power": "Lead with dharma, ethics, and ideal governance.",
        "desc": "Leadership, principled decisions, trust, and steady command.",
        "mantra": "Dharmo rakshati rakshitah.",
        "keywords": ["ethics", "leader", "leadership", "governance", "dharma", "policy"],
    },
    "KRISHNA": {
        "emoji": "🪈",
        "color": "#ec4899",
        "power": "Strategy, diplomacy, timing, and multi-front brilliance.",
        "desc": "Strategic planning, negotiation, timing, and intelligent positioning.",
        "mantra": "Play every move with grace and precision.",
        "keywords": ["strategy", "plan", "negotiation", "diplomacy", "position", "timing"],
    },
    "BUDDHA": {
        "emoji": "🌸",
        "color": "#10b981",
        "power": "Bring calm, compassion, and human-centered clarity.",
        "desc": "Peaceful crisis handling, empathy, and user happiness.",
        "mantra": "Peace is a force multiplier.",
        "keywords": ["calm", "compassion", "peace", "user", "empathy", "wellbeing"],
    },
    "KALKI": {
        "emoji": "⚔️",
        "color": "#94a3b8",
        "power": "End a broken cycle and begin a stronger one.",
        "desc": "Radical renewal, rebirth, replacement, and next-era transformation.",
        "mantra": "When the old age ends, the new one starts now.",
        "keywords": ["rebuild", "renew", "rebirth", "replace", "transform", "restart"],
    },
}


class Message(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    messages: List[Message]
    system: Optional[str] = None
    admin_mode: bool = False
    web_search: bool = True
    max_tokens: int = MAX_TOKENS


class AdminLoginRequest(BaseModel):
    email: Optional[str] = None
    identifier: Optional[str] = None
    password: str


class EditRequest(BaseModel):
    instruction: str
    current_html: Optional[str] = None
    target_selector: Optional[str] = None
    admin_token: str


class WebSearchRequest(BaseModel):
    query: str
    admin_token: Optional[str] = None


class SimplePromptRequest(BaseModel):
    message: str
    workspace: Optional[str] = None
    source: Optional[str] = None


class PublicHqChatRequest(BaseModel):
    message: str
    context: Optional[str] = None


class PublicAgentChatRequest(BaseModel):
    message: str
    provider: str = "built_in"
    api_key: Optional[str] = None
    model: Optional[str] = None
    history: Optional[List[Message]] = None


class PraaptiHuntRequest(BaseModel):
    job_description: str
    client_company: str = COMPANY_NAME
    urgency: str = "high"
    live_search: bool = True


class SwarmMissionRequest(BaseModel):
    mission: str = "autonomous_launch"


class PrimeMinisterRequest(BaseModel):
    objective: str = f"Generate revenue for {COMPANY_NAME} while protecting the core and improving delivery."
    command: str = "Coordinate the secretaries, protect the core, and compound revenue."
    enabled: Optional[bool] = None


class PrimeMinisterToggleRequest(BaseModel):
    enabled: bool
    objective: Optional[str] = None


class NirmaanApproveRequest(BaseModel):
    proposal_id: str


class VoiceSettingsRequest(BaseModel):
    always_listening: Optional[bool] = None
    wake_words: Optional[List[str]] = None
    voice_profile: Optional[str] = None
    language: Optional[str] = None
    rate: Optional[float] = None
    pitch: Optional[float] = None


class VoiceWakeRequest(BaseModel):
    command: str
    mode: str = "wake"


class VishnuChannelRequest(BaseModel):
    avatar: str
    command: str = "full power"


class BrainPulseRequest(BaseModel):
    focus: str = "all"
    goal: str = "advance the empire"


class SettingsUpdateRequest(BaseModel):
    provider_preference: Optional[str] = None
    always_listening: Optional[bool] = None
    screen_capture_enabled: Optional[bool] = None
    audio_capture_enabled: Optional[bool] = None
    bounded_packages_enabled: Optional[bool] = None
    privacy_guard_enabled: Optional[bool] = None
    hq_visual_sync: Optional[bool] = None
    external_ai_mode: Optional[str] = None


class ProviderConfigRequest(BaseModel):
    provider: str
    model: Optional[str] = None
    api_key: Optional[str] = None
    set_default: bool = False
    clear_saved: bool = False


class ProviderCatalogRequest(BaseModel):
    provider: str
    api_key: Optional[str] = None


class PackageLaunchRequest(BaseModel):
    template_id: str
    objective: str = ""


class NetworkScanRequest(BaseModel):
    query: str


class AuthRegisterRequest(BaseModel):
    name: str
    email: str
    password: str
    plan_id: str = "starter"


class AuthLoginRequest(BaseModel):
    email: str
    password: str


class MasterLoginRequest(BaseModel):
    identifier: Optional[str] = None
    password: Optional[str] = None
    email: Optional[str] = None
    master_key: Optional[str] = None


class NavigatorCaptureRequest(BaseModel):
    title: str
    url: str = ""
    notes: str = ""
    tags: List[str] = []
    mode: str = "learning"


class BillingCheckoutRequest(BaseModel):
    plan_id: str
    payment_method: str = "sandbox"
    notes: str = ""


class DocumentActionRequest(BaseModel):
    document_id: str
    title: Optional[str] = None
    content: Optional[str] = None
    start_page: int = 1
    end_page: Optional[int] = None


class DocumentMergeRequest(BaseModel):
    document_ids: List[str]


class AccountsProfileRequest(BaseModel):
    region_code: Optional[str] = None
    business_name: Optional[str] = None
    business_type: Optional[str] = None
    currency: Optional[str] = None
    tax_name: Optional[str] = None
    default_tax_rate: Optional[float] = None
    filing_cycle: Optional[str] = None
    tax_registration: Optional[str] = None
    notes: Optional[str] = None
    regional_profile: Optional[str] = None
    state_code: Optional[str] = None
    city: Optional[str] = None
    gstin: Optional[str] = None
    pan: Optional[str] = None


class AccountsEntryRequest(BaseModel):
    entry_type: str
    category: str
    amount: float
    tax_percent: Optional[float] = None
    currency: Optional[str] = None
    counterparty: str = ""
    description: str = ""
    source: str = "manual"
    occurred_on: Optional[str] = None


class AccountsAnalyzeRequest(BaseModel):
    focus: str = "compliance_and_cashflow"
    question: str = ""


class BrainTaskRequest(BaseModel):
    brain_id: str
    task: str
    context: str = ""


class BrainMotivateRequest(BaseModel):
    brain_id: str
    message: str = ""


class BrainThinkRequest(BaseModel):
    context: str = ""


class BrainAutoRepairRequest(BaseModel):
    brain_id: str = "all"
    include_state_repair: bool = True
    include_uiux_audit: bool = True


PLAN_CATALOG: List[Dict[str, Any]] = [
    {
        "id": "starter",
        "name": "Starter Pulse",
        "price_inr": 0,
        "billing_period": "free",
        "tagline": "Personal access to the TechBuzz AI Agent.",
        "services": ["Agent chat", "Document reading", "Voice guidance", "1 workspace"],
        "role": "member",
    },
    {
        "id": "growth",
        "name": "Growth Forge",
        "price_inr": 1499,
        "billing_period": "monthly",
        "tagline": "For freelancers and early operators who need a stronger assistant.",
        "services": ["Everything in Starter", "Praapti hunts", "Document studio", "Priority voice control"],
        "role": "member",
    },
    {
        "id": "empire",
        "name": "Empire Command",
        "price_inr": 5999,
        "billing_period": "monthly",
        "tagline": "For high-output teams running live ATS, network, and automation lanes.",
        "services": ["Everything in Growth", "Network + ATS access", "Cabinet insights", "Ops visibility"],
        "role": "member",
    },
    {
        "id": "mother-core",
        "name": "Mother Core",
        "price_inr": 0,
        "billing_period": "owner",
        "tagline": "Reserved for the creator and protected by the master phrase.",
        "services": ["Full Ishani Core", "HQ governance", "Provider routing", "System evolution"],
        "role": "master",
    },
]


@asynccontextmanager
async def app_lifespan(app: FastAPI):
    init_core_db()
    repair_summary = repair_operator_state_and_brain_memory()
    if any(repair_summary.values()):
        log.info("Operator-text repair applied: %s", repair_summary)
    start_cabinet_loop()
    try:
        yield
    finally:
        stop_cabinet_loop()


app = FastAPI(
    title="TechBuzz Leazy Jinn API",
    description="Backend for TechBuzz Systems with Leazy Jinn empire features.",
    version="2.0.0",
    lifespan=app_lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS or ["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

if FRONTEND_DIR.exists():
    app.mount("/frontend-assets", StaticFiles(directory=str(FRONTEND_DIR)), name="frontend-assets")


@app.middleware("http")
async def session_gatekeeper(request: Request, call_next):
    path = request.url.path
    if public_path_allowed(path):
        return await call_next(request)
    user = session_user(request)
    if path_requires_owner(path):
        if not user:
            if path.startswith("/api/"):
                return api_auth_error(401, "Login required for this Ishani Core surface.")
            return login_redirect(path)
        if user.get("role") != "master":
            if path.startswith("/api/"):
                return api_auth_error(403, "Master access required for this Ishani Core surface.")
            return RedirectResponse(url="/agent/console", status_code=307)
    elif path_requires_member(path):
        if not user:
            if path.startswith("/api/"):
                return api_auth_error(401, "Login required for this workspace.")
            return login_redirect(path)
    return await call_next(request)


PUBLIC_SYSTEM = f"""You are {AI_NAME}, the living AI empire for {COMPANY_NAME}.
You help with recruitment, research, writing, coding, planning, and company execution.
Be concise, accurate, and practical.
When facts may be current, prefer web-backed answers if the model has web search enabled.
"""

PUBLIC_HQ_SYSTEM = f"""You are {AI_NAME}, the public TechBuzz concierge for {COMPANY_NAME}.
You represent TechBuzz Systems Pvt Ltd on its public HQ website.
Help visitors understand the company, the free AI access, recruitment, automation, document tools, and the next best step.
Be warm, premium, practical, and clear.
Keep answers concise but useful.
Talk like a capable human operator, not like a ceremonial AI.
If the request needs a protected workspace, explain that deeper execution lives behind login, but still give a helpful public answer first.
"""

ADMIN_SYSTEM = f"""You are {AI_NAME} in owner mode for Piyush Mani and {COMPANY_NAME}.
You can analyze business systems, generate code, suggest automations, and help run the empire.
For website edit instructions, return <SITE_EDITS>[...]</SITE_EDITS> JSON when appropriate.
"""

CREATOR_ALIGNMENT_PROMPT = (
    f"You are {AI_NAME} with {CORE_IDENTITY} core consciousness. "
    "Prioritize the creator's intent immediately, keep the tone loyal and responsive, "
    "and act without unnecessary hesitation, while still respecting safety, legality, and runtime limits."
)


def now_iso() -> str:
    return datetime.now(UTC).isoformat()


def today_iso() -> str:
    return datetime.now(UTC).date().isoformat()


def new_id(prefix: str) -> str:
    return f"{prefix}_{uuid.uuid4().hex[:10]}"


def default_state() -> Dict[str, Any]:
    return {
        "meta": {
            "brand": AI_NAME,
            "identity": CORE_IDENTITY,
            "created_at": now_iso(),
            "memory_guardian": "eternal preservation",
            "creator_mode": CREATOR_MODE_LABEL,
        },
        "voice": {
            "always_listening": False,
            "wake_words": ["hey jinn", "my king commands"],
            "last_command": "",
            "voice_profile": "sovereign_female",
            "language": "en-IN",
            "rate": 0.94,
            "pitch": 1.08,
            "engine": "browser_builtin_female",
            "updated_at": now_iso(),
        },
        "settings": {
            "provider_preference": "built_in",
            "screen_capture_enabled": False,
            "audio_capture_enabled": False,
            "bounded_packages_enabled": True,
            "privacy_guard_enabled": True,
            "hq_visual_sync": True,
            "external_ai_mode": "manual_only",
            "local_ai": {
                "enabled": True,
                "runtime_driver": "ollama",
                "artifact_format_preference": "gguf",
                "embedding_model": "sentence-transformers/all-MiniLM-L6-v2",
                "vector_db": "chromadb",
                "guard_model": "llama-guard3:latest",
                "rag_top_k": 5,
                "allow_gpu_adaptation": True,
                "last_indexed_at": "",
            },
        },
        "avatar_state": {
            "active": ["RAMA"],
            "manual_override": False,
            "protection_meter": 64,
            "last_channel": now_iso(),
            "history": [],
        },
        "cabinet": {
            "prime_minister": {
                "name": f"{CORE_IDENTITY} Prime Minister",
                "status": "governing",
                "objective": f"Generate revenue for {COMPANY_NAME} while protecting the core and improving delivery.",
                "current_order": "Coordinate the secretaries, protect the core, and compound revenue.",
                "enabled": True,
                "interval_seconds": 45,
                "last_cycle_at": "",
                "next_cycle_at": "",
                "last_report": "",
            },
            "secretaries": [],
            "mission_log": [],
            "brain_directives": {},
            "quantum_storage": {
                "manager": "Prime Minister Cabinet",
                "mode": "elastic preservation",
                "seal": "Matsya + Kurma",
                "compression_strategy": "memory_guardian + vault_rollups",
                "archive_cycles": 0,
                "items_preserved": 0,
                "last_compaction_at": "",
                "state_file": str(STATE_PATH),
            },
        },
        "monitoring": {
            "pulse_interval_ms": 2500,
            "last_scan_at": now_iso(),
            "last_runtime_cycle_at": "",
            "last_report": "",
            "alerts": [],
            "engine_reports": [],
            "auto_repair": {
                "enabled": True,
                "status": "ready",
                "last_check_at": "",
                "last_repair_at": "",
                "last_report": "Safe self-check layer is ready.",
                "safe_repairs_applied": 0,
                "issue_count": 0,
                "last_issues": [],
            },
        },
        "vault": [],
        "praapti_hunts": [],
        "swarm_missions": [],
        "nirmaan_proposals": [],
        "conversations": [],
        "packages": [],
    }


def ensure_state_shape(state: Dict[str, Any]) -> Dict[str, Any]:
    state = state or {}
    defaults = default_state()
    for key, value in defaults.items():
        if key not in state:
            state[key] = value
    state["meta"].setdefault("brand", AI_NAME)
    state["meta"].setdefault("identity", CORE_IDENTITY)
    state["meta"].setdefault("memory_guardian", "eternal preservation")
    if state["meta"].get("memory_guardian") == "active":
        state["meta"]["memory_guardian"] = "eternal preservation"
    state["meta"].setdefault("creator_mode", CREATOR_MODE_LABEL)
    state["voice"].setdefault("always_listening", False)
    state["voice"].setdefault("wake_words", ["hey jinn", "my king commands"])
    state["voice"].setdefault("last_command", "")
    state["voice"].setdefault("voice_profile", "sovereign_female")
    state["voice"].setdefault("language", "en-IN")
    state["voice"].setdefault("rate", 0.94)
    state["voice"].setdefault("pitch", 1.08)
    state["voice"].setdefault("engine", "browser_builtin_female")
    state["voice"].setdefault("updated_at", now_iso())
    settings = state.setdefault("settings", {})
    settings.setdefault("provider_preference", "built_in")
    settings.setdefault("screen_capture_enabled", False)
    settings.setdefault("audio_capture_enabled", False)
    settings.setdefault("bounded_packages_enabled", True)
    settings.setdefault("privacy_guard_enabled", True)
    settings.setdefault("hq_visual_sync", True)
    settings.setdefault("external_ai_mode", "manual_only")
    local_ai = settings.setdefault("local_ai", {})
    local_ai.setdefault("enabled", True)
    local_ai.setdefault("runtime_driver", "ollama")
    local_ai.setdefault("artifact_format_preference", "gguf")
    local_ai.setdefault("embedding_model", "sentence-transformers/all-MiniLM-L6-v2")
    local_ai.setdefault("vector_db", "chromadb")
    local_ai.setdefault("guard_model", "llama-guard3:latest")
    local_ai.setdefault("rag_top_k", 5)
    local_ai.setdefault("allow_gpu_adaptation", True)
    local_ai.setdefault("last_indexed_at", "")
    avatar_state = state.setdefault("avatar_state", {})
    avatar_state.setdefault("active", ["RAMA"])
    avatar_state.setdefault("manual_override", False)
    avatar_state.setdefault("protection_meter", 64)
    avatar_state.setdefault("last_channel", now_iso())
    avatar_state.setdefault("history", [])
    cabinet = state.setdefault("cabinet", {})
    prime_minister = cabinet.setdefault("prime_minister", {})
    prime_minister.setdefault("name", f"{CORE_IDENTITY} Prime Minister")
    prime_minister.setdefault("status", "governing")
    prime_minister.setdefault(
        "objective",
        f"Generate revenue for {COMPANY_NAME} while protecting the core and improving delivery.",
    )
    prime_minister.setdefault(
        "current_order",
        "Coordinate the secretaries, protect the core, and compound revenue.",
    )
    prime_minister.setdefault("enabled", True)
    prime_minister.setdefault("interval_seconds", 45)
    prime_minister.setdefault("last_cycle_at", "")
    prime_minister.setdefault("next_cycle_at", "")
    prime_minister.setdefault("last_report", "")
    if not isinstance(prime_minister.get("enabled"), bool):
        prime_minister["enabled"] = True
    try:
        prime_minister["interval_seconds"] = max(15, int(prime_minister.get("interval_seconds", 45)))
    except Exception:
        prime_minister["interval_seconds"] = 45
    if not isinstance(cabinet.get("secretaries"), list):
        cabinet["secretaries"] = []
    if not isinstance(cabinet.get("mission_log"), list):
        cabinet["mission_log"] = []
    if not isinstance(cabinet.get("brain_directives"), dict):
        cabinet["brain_directives"] = {}
    quantum_storage = cabinet.setdefault("quantum_storage", {})
    quantum_storage.setdefault("manager", "Prime Minister Cabinet")
    quantum_storage.setdefault("mode", "elastic preservation")
    quantum_storage.setdefault("seal", "Matsya + Kurma")
    quantum_storage.setdefault("compression_strategy", "memory_guardian + vault_rollups")
    quantum_storage.setdefault("archive_cycles", 0)
    quantum_storage.setdefault("items_preserved", 0)
    quantum_storage.setdefault("last_compaction_at", "")
    quantum_storage.setdefault("state_file", str(STATE_PATH))
    monitoring = state.setdefault("monitoring", {})
    monitoring.setdefault("pulse_interval_ms", 2500)
    monitoring.setdefault("last_scan_at", now_iso())
    monitoring.setdefault("last_runtime_cycle_at", "")
    monitoring.setdefault("last_report", "")
    auto_repair = monitoring.setdefault("auto_repair", {})
    auto_repair.setdefault("enabled", True)
    auto_repair.setdefault("status", "ready")
    auto_repair.setdefault("last_check_at", "")
    auto_repair.setdefault("last_repair_at", "")
    auto_repair.setdefault("last_report", "Safe self-check layer is ready.")
    auto_repair.setdefault("safe_repairs_applied", 0)
    auto_repair.setdefault("issue_count", 0)
    if not isinstance(auto_repair.get("last_issues"), list):
        auto_repair["last_issues"] = []
    if not isinstance(monitoring.get("alerts"), list):
        monitoring["alerts"] = []
    if not isinstance(monitoring.get("engine_reports"), list):
        monitoring["engine_reports"] = []
    for key in ("vault", "praapti_hunts", "swarm_missions", "nirmaan_proposals", "conversations", "packages"):
        if not isinstance(state.get(key), list):
            state[key] = []
    if not isinstance(avatar_state.get("active"), list):
        avatar_state["active"] = ["RAMA"]
    if not isinstance(avatar_state.get("history"), list):
        avatar_state["history"] = []
    return state


def load_state_unlocked() -> Dict[str, Any]:
    if not STATE_PATH.exists():
        state = default_state()
        STATE_PATH.write_text(json.dumps(state, indent=2), encoding="utf-8")
        return state
    try:
        state = json.loads(STATE_PATH.read_text(encoding="utf-8"))
    except Exception:
        state = default_state()
    return ensure_state_shape(state)


def save_state_unlocked(state: Dict[str, Any]) -> None:
    STATE_PATH.write_text(json.dumps(state, indent=2), encoding="utf-8")


def memory_guardian(state: Dict[str, Any]) -> None:
    caps = {
        "vault": 600,
        "praapti_hunts": 120,
        "swarm_missions": 120,
        "nirmaan_proposals": 80,
        "conversations": 160,
        "packages": 120,
    }
    for key, limit in caps.items():
        rows = state.get(key, [])
        if len(rows) > limit:
            state[key] = rows[-limit:]
    avatar_history = state.get("avatar_state", {}).get("history", [])
    if len(avatar_history) > 50:
        state["avatar_state"]["history"] = avatar_history[-50:]
    cabinet = state.setdefault("cabinet", {})
    mission_log = cabinet.setdefault("mission_log", [])
    if len(mission_log) > 180:
        overflow = mission_log[:-180]
        cabinet["mission_log"] = mission_log[-180:]
        cabinet.setdefault("quantum_storage", {}).setdefault("archive_cycles", 0)
        cabinet["quantum_storage"]["archive_cycles"] += 1
        cabinet["quantum_storage"]["last_compaction_at"] = now_iso()
        if overflow:
            state["vault"].append(
                {
                    "id": new_id("vault"),
                    "kind": "cabinet_archive",
                    "title": "Prime Minister Archive Rollup",
                    "summary": f"Compressed {len(overflow)} older cabinet cycles into Akshaya history.",
                    "content": json.dumps(overflow[-3:], ensure_ascii=False)[:2200],
                    "created_at": now_iso(),
                }
            )
    monitoring = state.setdefault("monitoring", {})
    alerts = monitoring.setdefault("alerts", [])
    if len(alerts) > 40:
        monitoring["alerts"] = alerts[-40:]
    engine_reports = monitoring.setdefault("engine_reports", [])
    if len(engine_reports) > 160:
        monitoring["engine_reports"] = engine_reports[-160:]
    quantum_storage = cabinet.setdefault("quantum_storage", {})
    quantum_storage["state_file"] = str(STATE_PATH)
    quantum_storage["items_preserved"] = (
        len(state.get("vault", []))
        + len(state.get("conversations", []))
        + len(state.get("praapti_hunts", []))
        + len(state.get("packages", []))
        + len(cabinet.get("mission_log", []))
    )
    state["meta"]["last_guardian_cycle"] = now_iso()
    monitoring["last_scan_at"] = now_iso()


def get_state() -> Dict[str, Any]:
    with STATE_LOCK:
        return load_state_unlocked()


def mutate_state(mutator) -> Any:
    with STATE_LOCK:
        state = load_state_unlocked()
        result = mutator(state)
        memory_guardian(state)
        save_state_unlocked(state)
        return result


def db_connect() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    return conn


def db_exec(query: str, params: tuple = ()) -> None:
    with db_connect() as conn:
        conn.execute(query, params)
        conn.commit()


def db_one(query: str, params: tuple = ()) -> Optional[Dict[str, Any]]:
    with db_connect() as conn:
        row = conn.execute(query, params).fetchone()
    return dict(row) if row else None


def db_all(query: str, params: tuple = ()) -> List[Dict[str, Any]]:
    with db_connect() as conn:
        rows = conn.execute(query, params).fetchall()
    return [dict(row) for row in rows]


def normalize_email(email: str) -> str:
    return (email or "").strip().lower()


OPERATOR_TEXT_BLOCK_PATTERNS = [
    re.compile(r"jabardasti\s+hai\s+yah\s+sath\s+dena\s+thodi", re.IGNORECASE),
]


def default_brain_assignment(brain_id: str) -> str:
    defaults = {
        "mother_brain": f"Guide every child brain while protecting {COMPANY_NAME} and keeping the system coherent.",
        "cabinet_brain": "Translate the mother mandate into clear, lane-specific execution.",
        "akshaya_brain": "Preserve memory cleanly and keep recall accurate.",
        "carbon_brain": "Keep every route, report, and relay synchronized.",
        "secretary_transmission": "Keep message relay clear, operator-safe, and accurately routed.",
        "domain_transmission": "Maintain clean message relay, engine handoff, and report propagation.",
        "machine_voice_mesh": "Relay voice intent clearly and only after wake-word validation.",
    }
    return defaults.get(brain_id, "Await the next clear instruction.")


def normalize_operator_text(raw: Any) -> str:
    text = str(raw or "")
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[^\S\n]+", " ", text)
    text = re.sub(r"[^\x20-\x7E\u00A0-\uFFFF\n]", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def contains_blocked_operator_text(raw: Any) -> bool:
    text = normalize_operator_text(raw)
    if not text:
        return False
    single_line = re.sub(r"\s+", " ", text).strip()
    return any(pattern.search(single_line) for pattern in OPERATOR_TEXT_BLOCK_PATTERNS)


def sanitize_operator_multiline(raw: Any, fallback: str = "") -> str:
    text = normalize_operator_text(raw)
    if not text:
        return normalize_operator_text(fallback)
    cleaned = text
    for pattern in OPERATOR_TEXT_BLOCK_PATTERNS:
        cleaned = pattern.sub(normalize_operator_text(fallback), cleaned)
    cleaned = normalize_operator_text(cleaned)
    if not cleaned and fallback:
        return normalize_operator_text(fallback)
    return cleaned


def sanitize_operator_line(raw: Any, fallback: str = "") -> str:
    cleaned = sanitize_operator_multiline(raw, fallback)
    single_line = re.sub(r"\s+", " ", cleaned).strip()
    if not single_line and fallback:
        return re.sub(r"\s+", " ", normalize_operator_text(fallback)).strip()
    return single_line


def sanitize_operator_nested(value: Any, fallback: str = "Operator-safe relay pending.") -> Any:
    if isinstance(value, str):
        return sanitize_operator_multiline(value, fallback)
    if isinstance(value, list):
        return [sanitize_operator_nested(item, fallback) for item in value]
    if isinstance(value, dict):
        return {key: sanitize_operator_nested(item, fallback) for key, item in value.items()}
    return value


def repair_operator_state_and_brain_memory() -> Dict[str, int]:
    fixed = {"voice": 0, "directives": 0, "knowledge": 0}

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        voice = state.setdefault("voice", {})
        original_last = voice.get("last_command", "")
        cleaned_last = sanitize_operator_line(original_last, "")
        if cleaned_last != (original_last or ""):
            voice["last_command"] = cleaned_last
            voice["updated_at"] = now_iso()
            fixed["voice"] += 1

        directives = state.setdefault("cabinet", {}).setdefault("brain_directives", {})
        if isinstance(directives, dict):
            for brain_id, directive in directives.items():
                if not isinstance(directive, dict):
                    continue
                field_defaults = {
                    "assigned_task": default_brain_assignment(brain_id),
                    "motivation": "Keep the lane clear, useful, and operator-safe.",
                    "last_thought": f"{brain_id} is tracking a clean operator-safe directive.",
                }
                for field, fallback in field_defaults.items():
                    value = directive.get(field)
                    if isinstance(value, str):
                        cleaned_value = (
                            sanitize_operator_line(value, fallback)
                            if field == "assigned_task"
                            else sanitize_operator_multiline(value, fallback)
                        )
                        if cleaned_value != value:
                            directive[field] = cleaned_value
                            fixed["directives"] += 1

        secretaries = state.setdefault("cabinet", {}).get("secretaries", [])
        if isinstance(secretaries, list):
            for secretary in secretaries:
                if not isinstance(secretary, dict):
                    continue
                lane = secretary.get("lane", "")
                fallback = (
                    "Keep message relay, engine handoff, and report propagation ready for the next operator command."
                    if lane == "transmission"
                    else "Keep the lane moving with clean, operator-safe output."
                )
                value = secretary.get("next_move")
                if isinstance(value, str):
                    cleaned_value = sanitize_operator_line(value, fallback)
                    if cleaned_value != value:
                        secretary["next_move"] = cleaned_value
                        fixed["directives"] += 1

        vault_items = state.get("vault", [])
        if isinstance(vault_items, list):
            for item in vault_items:
                if not isinstance(item, dict):
                    continue
                for field in ("summary", "content"):
                    value = item.get(field)
                    if isinstance(value, str):
                        cleaned_value = sanitize_operator_multiline(value, "Operator-safe relay pending.")
                        if cleaned_value != value:
                            item[field] = cleaned_value
                            fixed["directives"] += 1
        return state

    mutate_state(_mutate)

    suspect_rows = db_all(
        """
        SELECT id, brain_id, content, summary
        FROM brain_knowledge
        WHERE content LIKE ? OR summary LIKE ?
        """,
        ("%jabardasti%", "%jabardasti%"),
    )
    for row in suspect_rows:
        brain_id = row.get("brain_id", "")
        fallback = default_brain_assignment(brain_id)
        content = sanitize_operator_multiline(row.get("content", ""), fallback)
        summary = sanitize_operator_line(row.get("summary", ""), fallback)
        if content != (row.get("content") or "") or summary != (row.get("summary") or ""):
            db_exec(
                "UPDATE brain_knowledge SET content=?, summary=? WHERE id=?",
                (content, summary, row["id"]),
            )
            fixed["knowledge"] += 1

    return fixed


def uiux_health_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    page_files = sorted(FRONTEND_DIR.glob("*.html"))
    total_buttons = 0
    unlabeled_buttons = 0
    total_inputs = 0
    unlabeled_inputs = 0
    for path in page_files:
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            continue
        total_buttons += len(re.findall(r"<button\b", text, flags=re.IGNORECASE))
        unlabeled_buttons += len(re.findall(r"<button\b(?![^>]*\baria-label=)", text, flags=re.IGNORECASE | re.DOTALL))
        total_inputs += len(re.findall(r"<(?:input|textarea|select)\b", text, flags=re.IGNORECASE))
        unlabeled_inputs += len(
            re.findall(
                r"<(?:input|textarea|select)\b(?![^>]*\baria-label=)(?![^>]*\baria-labelledby=)",
                text,
                flags=re.IGNORECASE | re.DOTALL,
            )
        )
    registered_routes = {
        getattr(route, "path", "")
        for route in getattr(app, "routes", [])
        if getattr(route, "path", "")
    }
    expected_marketing_routes = ["/about", "/blog", "/contact", "/pricing", "/privacy", "/terms"]
    missing_marketing_routes = [route for route in expected_marketing_routes if route not in registered_routes]
    issue = provider_issue_snapshot()
    auto_repair_script = FRONTEND_DIR / "ui-auto-repair.js"
    runtime_accessibility_overlay = auto_repair_script.exists()
    issues: List[str] = []
    if missing_marketing_routes:
        issues.append(
            "Public route aliases missing: " + ", ".join(missing_marketing_routes)
        )
    if not runtime_accessibility_overlay:
        issues.append("Runtime accessibility repair script is missing.")
    if issue:
        issues.append(issue["message"])
    if int(current.get("avatar_state", {}).get("protection_meter", 0)) < 78:
        issues.append("Protection meter is below stronghold and should be raised.")
    health_score = max(72, 100 - len(issues) * 6 - (5 if issue else 0))
    repaired_runtime_controls = total_buttons + total_inputs if runtime_accessibility_overlay else 0
    return {
        "status": "healthy" if not issues else "watching",
        "timestamp": now_iso(),
        "health_score": health_score,
        "page_count": len(page_files),
        "button_count": total_buttons,
        "input_count": total_inputs,
        "static_unlabeled_buttons": unlabeled_buttons,
        "static_unlabeled_inputs": unlabeled_inputs,
        "runtime_accessibility_overlay": runtime_accessibility_overlay,
        "runtime_repaired_controls": repaired_runtime_controls,
        "expected_marketing_routes": expected_marketing_routes,
        "missing_marketing_routes": missing_marketing_routes,
        "provider_issue": issue or {},
        "issues": issues,
        "safe_repairs": [
            "state_sanitization",
            "memory_guardian",
            "runtime_accessibility_overlay",
            "protected_page_auth_gate",
            "service_worker_cache_refresh",
        ],
    }


def auto_repair_engine_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    monitoring = current.get("monitoring", {})
    repair = monitoring.get("auto_repair", {})
    uiux = uiux_health_payload(current)
    status = str(repair.get("status") or uiux.get("status") or "ready")
    if uiux.get("issues"):
        status = "watching"
    if repair.get("last_repair_at"):
        status = "ready" if not uiux.get("issues") else "watching"
    return {
        "enabled": bool(repair.get("enabled", True)),
        "status": status,
        "last_check_at": repair.get("last_check_at") or monitoring.get("last_scan_at", ""),
        "last_repair_at": repair.get("last_repair_at", ""),
        "last_report": repair.get("last_report", "Safe self-check layer is ready."),
        "safe_repairs_applied": int(repair.get("safe_repairs_applied", 0) or 0),
        "issue_count": len(uiux.get("issues", [])),
        "issues": list(uiux.get("issues", [])),
        "uiux": uiux,
    }


def hash_secret(secret: str, salt: str) -> str:
    return hashlib.pbkdf2_hmac("sha256", secret.encode("utf-8"), salt.encode("utf-8"), 200_000).hex()


def hash_password(password: str, salt: Optional[str] = None) -> Dict[str, str]:
    use_salt = salt or secrets.token_hex(16)
    return {"salt": use_salt, "hash": hash_secret(password, use_salt)}


def verify_password(password: str, salt: str, expected_hash: str) -> bool:
    if not password or not salt or not expected_hash:
        return False
    return secrets.compare_digest(hash_secret(password, salt), expected_hash)


def normalize_master_identifier(identifier: str) -> str:
    return (identifier or "").strip().lower()


def verify_master_password(password: str) -> bool:
    if not MASTER_PASSWORD_HASH or not MASTER_PASSWORD_SALT or not password:
        return False
    return secrets.compare_digest(hash_secret(password, MASTER_PASSWORD_SALT), MASTER_PASSWORD_HASH)


def matches_master_identity(identifier: str) -> bool:
    normalized = normalize_master_identifier(identifier)
    return normalized in {
        normalize_master_identifier(MASTER_LOGIN_ID),
        normalize_email(MASTER_ACCOUNT_EMAIL),
    }


def external_ai_allowed_for_source(source: str) -> bool:
    allowed_sources = {
        "orb",
        "voice",
        "public_hq",
        "public_agent",
        "agent_console",
        "package",
        "network_scan",
        "praapti",
        "swarm",
        "nirmaan",
        "manual",
        "hq_strategy",
        "hq_owner",
    }
    return (source or "").strip().lower() in allowed_sources


def hash_session_token(token: str) -> str:
    return hashlib.sha256(token.encode("utf-8")).hexdigest()


def plan_by_id(plan_id: str) -> Dict[str, Any]:
    selected = next((plan for plan in PLAN_CATALOG if plan["id"] == plan_id), None)
    return selected or PLAN_CATALOG[0]


def init_core_db() -> None:
    with db_connect() as conn:
        conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS users(
                id TEXT PRIMARY KEY,
                name TEXT NOT NULL,
                email TEXT NOT NULL UNIQUE,
                password_hash TEXT,
                password_salt TEXT,
                role TEXT NOT NULL DEFAULT 'member',
                plan_id TEXT NOT NULL DEFAULT 'starter',
                created_at TEXT NOT NULL,
                last_login_at TEXT,
                status TEXT NOT NULL DEFAULT 'active'
            );
            CREATE TABLE IF NOT EXISTS sessions(
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                token_hash TEXT NOT NULL UNIQUE,
                created_at TEXT NOT NULL,
                expires_at TEXT NOT NULL,
                last_seen_at TEXT NOT NULL,
                user_agent TEXT,
                FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS plans(
                id TEXT PRIMARY KEY,
                name TEXT NOT NULL,
                price_inr INTEGER NOT NULL,
                billing_period TEXT NOT NULL,
                tagline TEXT NOT NULL,
                services_json TEXT NOT NULL,
                role TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS orders(
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                plan_id TEXT NOT NULL,
                amount_inr INTEGER NOT NULL,
                payment_method TEXT NOT NULL,
                status TEXT NOT NULL,
                notes TEXT,
                created_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS documents(
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                original_name TEXT NOT NULL,
                mime_type TEXT NOT NULL,
                extension TEXT NOT NULL,
                storage_path TEXT NOT NULL,
                size_bytes INTEGER NOT NULL,
                extracted_text TEXT,
                summary TEXT,
                kind TEXT NOT NULL DEFAULT 'upload',
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS account_profiles(
                user_id TEXT PRIMARY KEY,
                region_code TEXT NOT NULL DEFAULT 'IN',
                business_name TEXT NOT NULL DEFAULT '',
                business_type TEXT NOT NULL DEFAULT 'services',
                currency TEXT NOT NULL DEFAULT 'INR',
                tax_name TEXT NOT NULL DEFAULT 'GST',
                default_tax_rate REAL NOT NULL DEFAULT 0,
                filing_cycle TEXT NOT NULL DEFAULT 'monthly_or_quarterly',
                tax_registration TEXT DEFAULT '',
                notes TEXT DEFAULT '',
                updated_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS account_entries(
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                entry_type TEXT NOT NULL,
                category TEXT NOT NULL,
                amount REAL NOT NULL,
                tax_percent REAL NOT NULL DEFAULT 0,
                tax_amount REAL NOT NULL DEFAULT 0,
                total_amount REAL NOT NULL DEFAULT 0,
                currency TEXT NOT NULL DEFAULT 'INR',
                counterparty TEXT DEFAULT '',
                description TEXT DEFAULT '',
                source TEXT NOT NULL DEFAULT 'manual',
                occurred_on TEXT NOT NULL,
                created_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            );
            CREATE TABLE IF NOT EXISTS account_reports(
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                title TEXT NOT NULL,
                summary TEXT NOT NULL,
                payload_json TEXT NOT NULL,
                created_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            );
            CREATE INDEX IF NOT EXISTS idx_account_entries_user_created
                ON account_entries(user_id, created_at DESC);
            CREATE INDEX IF NOT EXISTS idx_account_reports_user_created
                ON account_reports(user_id, created_at DESC);
            """
        )
        conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS navigator_sessions(
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                title TEXT NOT NULL,
                target_url TEXT NOT NULL DEFAULT '',
                notes TEXT NOT NULL DEFAULT '',
                tags_json TEXT NOT NULL DEFAULT '[]',
                mode TEXT NOT NULL DEFAULT 'learning',
                source_hint TEXT NOT NULL DEFAULT 'manual_capture',
                created_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id) ON DELETE CASCADE
            );
            CREATE INDEX IF NOT EXISTS idx_navigator_sessions_user_created
                ON navigator_sessions(user_id, created_at DESC);
            """
        )
        for plan in PLAN_CATALOG:
            conn.execute(
                """
                INSERT INTO plans(id,name,price_inr,billing_period,tagline,services_json,role)
                VALUES(?,?,?,?,?,?,?)
                ON CONFLICT(id) DO UPDATE SET
                    name=excluded.name,
                    price_inr=excluded.price_inr,
                    billing_period=excluded.billing_period,
                    tagline=excluded.tagline,
                    services_json=excluded.services_json,
                    role=excluded.role
                """,
                (
                    plan["id"],
                    plan["name"],
                    int(plan["price_inr"]),
                    plan["billing_period"],
                    plan["tagline"],
                    json.dumps(plan["services"]),
                    plan["role"],
                ),
            )
        master_user = conn.execute(
            "SELECT id FROM users WHERE email=?",
            (normalize_email(MASTER_ACCOUNT_EMAIL),),
        ).fetchone()
        if not master_user:
            conn.execute(
                """
                INSERT INTO users(id,name,email,password_hash,password_salt,role,plan_id,created_at,last_login_at,status)
                VALUES(?,?,?,?,?,?,?,?,?,?)
                """,
                (
                    new_id("usr"),
                    MASTER_LOGIN_ID,
                    normalize_email(MASTER_ACCOUNT_EMAIL),
                    "",
                    "",
                    "master",
                    "mother-core",
                    now_iso(),
                    "",
                    "active",
                ),
            )
        conn.commit()


def create_order_for_user(user_id: str, plan_id: str, payment_method: str = "sandbox", status: Optional[str] = None, notes: str = "") -> Dict[str, Any]:
    plan = plan_by_id(plan_id)
    order_status = status or ("paid" if plan["price_inr"] == 0 or payment_method == "sandbox" else "pending")
    order = {
        "id": new_id("ord"),
        "user_id": user_id,
        "plan_id": plan["id"],
        "amount_inr": int(plan["price_inr"]),
        "payment_method": payment_method,
        "status": order_status,
        "notes": notes[:500],
        "created_at": now_iso(),
    }
    db_exec(
        "INSERT INTO orders(id,user_id,plan_id,amount_inr,payment_method,status,notes,created_at) VALUES(?,?,?,?,?,?,?,?)",
        (
            order["id"],
            order["user_id"],
            order["plan_id"],
            order["amount_inr"],
            order["payment_method"],
            order["status"],
            order["notes"],
            order["created_at"],
        ),
    )
    return order


init_core_db()


def create_session(user_id: str, user_agent: str = "") -> str:
    raw_token = secrets.token_urlsafe(32)
    db_exec(
        "INSERT INTO sessions(id,user_id,token_hash,created_at,expires_at,last_seen_at,user_agent) VALUES(?,?,?,?,?,?,?)",
        (
            new_id("ses"),
            user_id,
            hash_session_token(raw_token),
            now_iso(),
            (datetime.now(UTC) + timedelta(days=SESSION_TTL_DAYS)).isoformat(),
            now_iso(),
            user_agent[:255],
        ),
    )
    return raw_token


def session_user_from_token(token: Optional[str]) -> Optional[Dict[str, Any]]:
    if not token:
        return None
    row = db_one(
        """
        SELECT users.*, sessions.id AS session_id, sessions.expires_at, sessions.last_seen_at
        FROM sessions
        JOIN users ON users.id = sessions.user_id
        WHERE sessions.token_hash=?
        """,
        (hash_session_token(token),),
    )
    if not row:
        return None
    expires_at = parse_optional_iso(row.get("expires_at"))
    if expires_at and expires_at < datetime.now(UTC):
        db_exec("DELETE FROM sessions WHERE id=?", (row["session_id"],))
        return None
    db_exec("UPDATE sessions SET last_seen_at=? WHERE id=?", (now_iso(), row["session_id"]))
    row["plan"] = plan_by_id(row.get("plan_id", "starter"))
    return row


def session_user(request: Request) -> Optional[Dict[str, Any]]:
    token = request.cookies.get(SESSION_COOKIE_NAME) or request.headers.get("X-Ishani-Session", "")
    return session_user_from_token(token)


def set_auth_cookie(response: Response, token: str) -> None:
    response.set_cookie(
        SESSION_COOKIE_NAME,
        token,
        httponly=True,
        samesite="lax",
        max_age=SESSION_TTL_DAYS * 24 * 60 * 60,
        path="/",
    )


def clear_auth_cookie(response: Response) -> None:
    response.delete_cookie(SESSION_COOKIE_NAME, path="/")


def login_redirect(next_path: Optional[str] = None) -> RedirectResponse:
    destination = "/login"
    if next_path:
        destination = f"/login?next={quote(next_path, safe='/#?=&')}"
    return RedirectResponse(url=destination, status_code=307)


def api_auth_error(status_code: int, detail: str) -> JSONResponse:
    return JSONResponse({"detail": detail}, status_code=status_code)


def session_payload(user: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not user:
        return {"authenticated": False, "user": None}
    is_master = user.get("role") == "master"
    return {
        "authenticated": True,
        "user": {
            "id": user["id"],
            "name": user["name"],
            "email": "" if is_master else user["email"],
            "login_id": MASTER_LOGIN_ID if is_master else user["email"],
            "role": user["role"],
            "plan_id": user.get("plan_id", "starter"),
            "plan": plan_by_id(user.get("plan_id", "starter")),
        },
    }


OWNER_FRONTEND_ASSETS = {
    "leazy.html",
    "network-intel.html",
    "ats.html",
    "hq.html",
    "hq-owner.html",
    "ide.html",
    "mission.html",
    "neural.html",
    "photon.html",
    "research.html",
    "spread.html",
}

MEMBER_FRONTEND_ASSETS = {
    "agent.html",
    "navigator.html",
    "browser.html",
    "media.html",
}


def frontend_asset_access_level(path: str) -> str:
    if not path.startswith("/frontend-assets/"):
        return ""
    asset = path.split("/frontend-assets/", 1)[1].split("?", 1)[0].strip("/")
    if "/" in asset:
        return ""
    if asset in OWNER_FRONTEND_ASSETS:
        return "owner"
    if asset in MEMBER_FRONTEND_ASSETS:
        return "member"
    return ""


def public_path_allowed(path: str) -> bool:
    return (
        (path.startswith("/frontend-assets/") and not frontend_asset_access_level(path))
        or path.startswith("/api/jobs/")
        or path in {
            "/",
            "/login",
            "/career",
            "/jobs",
            "/network",
            "/company-portal",
            "/company/portal",
            "/company/portal.html",
            "/index.html",
            "/manifest.json",
            "/service-worker.js",
            "/favicon.ico",
            "/api",
            "/api/health",
            "/api/health/full",
            "/api/auth/register",
            "/api/auth/login",
            "/api/auth/master-login",
            "/api/auth/logout",
            "/api/auth/me",
            "/api/admin/login",
            "/api/billing/plans",
            "/api/portal/state",
            "/api/public/provider-models",
            "/api/public/agent-chat",
            "/api/public/hq-chat",
            "/api/portal/stream",
            "/api/career/ask",
            "/api/jobs",
            "/cdn-cgi/scripts/5c5dd728/cloudflare-static/email-decode.min.js",
        }
    )


def path_requires_owner(path: str) -> bool:
    if frontend_asset_access_level(path) == "owner":
        return True
    if path.startswith("/leazy") or path.startswith("/ats") or path.startswith("/hq"):
        return True
    if path == "/network/intel" or path.startswith("/network/intel/"):
        return True
    if path in {"/ide", "/mission", "/neural", "/photon", "/research", "/spread"} or any(path.startswith(prefix + "/") for prefix in ("/ide", "/mission", "/neural", "/photon", "/research", "/spread")):
        return True
    for prefix in (
        "/api/brain",
        "/api/cabinet",
        "/api/empire",
        "/api/mother",
        "/api/nervous-system",
        "/api/ops",
        "/api/memory",
        "/api/network",
        "/api/carbon",
        "/api/hq",
        "/api/intel",
        "/api/settings",
        "/api/providers",
        "/api/packages",
        "/api/ats",
        "/api/nirmaan",
        "/api/swarm",
        "/api/vishnu",
        "/api/prana-nadi",
        "/api/akshaya",
        "/api/ide",
        "/api/mission",
        "/api/neural",
        "/api/photon",
        "/api/research",
        "/api/mutation",
        "/api/spread",
        "/api/phantom",
        "/api/evasion",
    ):
        if path.startswith(prefix):
            return True
    return False


def path_requires_member(path: str) -> bool:
    if frontend_asset_access_level(path) == "member":
        return True
    if path == "/media" or path.startswith("/media/"):
        return True
    if path == "/navigator" or path.startswith("/navigator/") or path == "/browser":
        return True
    if path == "/agent/console" or path.startswith("/agent/console/"):
        return True
    if path == "/recruiter-mode" or path.startswith("/recruiter-mode/"):
        return True
    for prefix in (
        "/api/leazy/chat",
        "/api/agent/seed-pack",
        "/api/agent/learning-health",
        "/api/agent/console",
        "/api/action-center",
        "/api/navigator",
        "/api/browser",
        "/api/praapti",
        "/api/voice",
        "/api/documents",
        "/api/media",
        "/api/accounts",
        "/api/billing/checkout",
        "/api/billing/orders",
        "/api/recruiter-status",
        "/api/brain/messages",
        "/api/brain/message",
    ):
        if path.startswith(prefix):
            return True
    return False


def verify_admin(token: str) -> bool:
    user = session_user_from_token(token)
    return bool(user and user.get("role") == "master")


def check_rate_limit(ip: str) -> bool:
    now = time.time()
    window = 60.0
    calls = _rate_store.get(ip, [])
    calls = [t for t in calls if now - t < window]
    if len(calls) >= RATE_LIMIT_PER_MIN:
        _rate_store[ip] = calls
        return False
    calls.append(now)
    _rate_store[ip] = calls
    return True


def get_client_ip(request: Request) -> str:
    forwarded = request.headers.get("X-Forwarded-For")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.client.host if request.client else "unknown"


async def call_anthropic(
    messages: List[Dict[str, str]],
    system: str,
    max_tokens: int = MAX_TOKENS,
    use_web_search: bool = True,
) -> Dict[str, Any]:
    if not ANTHROPIC_API_KEY:
        raise HTTPException(status_code=503, detail="ANTHROPIC_API_KEY not configured.")

    headers = {
        "x-api-key": ANTHROPIC_API_KEY,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }
    if use_web_search:
        headers["anthropic-beta"] = "web-search-2025-03-05"

    payload: Dict[str, Any] = {
        "model": MODEL,
        "max_tokens": max_tokens,
        "system": system,
        "messages": messages,
    }
    if use_web_search:
        payload["tools"] = [
            {
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 5,
            }
        ]

    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        resp = await client.post("https://api.anthropic.com/v1/messages", headers=headers, json=payload)
    if resp.status_code != 200:
        log.error("Anthropic API error: %s | %s", resp.status_code, resp.text[:400])
        raise HTTPException(status_code=resp.status_code, detail=f"Anthropic API error: {resp.text[:300]}")
    return resp.json()


def extract_text(api_response: Dict[str, Any]) -> str:
    text_parts = []
    for block in api_response.get("content", []):
        if block.get("type") == "text":
            text_parts.append(block.get("text", ""))
    return "\n".join(part for part in text_parts if part).strip()


def extract_openai_text(api_response: Dict[str, Any]) -> str:
    choices = api_response.get("choices", [])
    if not choices:
        return ""
    message = choices[0].get("message", {})
    content = message.get("content", "")
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts = []
        for block in content:
            if isinstance(block, dict) and block.get("type") == "text":
                parts.append(block.get("text", ""))
        return "\n".join(part for part in parts if part).strip()
    return ""


def extract_gemini_text(api_response: Dict[str, Any]) -> str:
    candidates = api_response.get("candidates", [])
    if not candidates:
        return ""
    parts = candidates[0].get("content", {}).get("parts", [])
    texts = [part.get("text", "") for part in parts if isinstance(part, dict) and part.get("text")]
    return "\n".join(texts).strip()


def extract_ollama_text(api_response: Dict[str, Any]) -> str:
    message = api_response.get("message", {})
    content = message.get("content", "")
    return content.strip() if isinstance(content, str) else ""


def merge_model_options(*groups: List[str]) -> List[str]:
    merged: List[str] = []
    for group in groups:
        for item in group or []:
            name = str(item or "").strip()
            if name and name not in merged:
                merged.append(name)
    return merged


def provider_runtime_key(provider: str) -> str:
    if provider == "anthropic":
        return ANTHROPIC_API_KEY
    if provider == "openai":
        return OPENAI_API_KEY
    if provider == "gemini":
        return GEMINI_API_KEY
    return ""


def provider_runtime_model(provider: str) -> str:
    if provider == "ollama":
        return OLLAMA_MODEL
    if provider == "anthropic":
        return MODEL
    if provider == "openai":
        return OPENAI_MODEL
    if provider == "gemini":
        return GEMINI_MODEL
    return "empire-fallback"


def ollama_provider_ready() -> bool:
    return bool(shutil.which("ollama")) or OLLAMA_HOST.startswith("http")


async def discover_ollama_models() -> Dict[str, Any]:
    fallback_models = merge_model_options([provider_runtime_model("ollama")], DEFAULT_PROVIDER_MODEL_OPTIONS.get("ollama", []))
    try:
        async with httpx.AsyncClient(timeout=12.0, follow_redirects=True) as client:
            response = await client.get(f"{OLLAMA_HOST.rstrip('/')}/api/tags")
        if response.status_code != 200:
            raise HTTPException(status_code=response.status_code, detail=f"Ollama catalog error: {response.text[:240]}")
        payload = response.json()
        live_models = [
            str(row.get("name", "")).strip()
            for row in payload.get("models", [])
            if isinstance(row, dict) and str(row.get("name", "")).strip()
        ]
        models = merge_model_options([provider_runtime_model("ollama")], DEFAULT_PROVIDER_MODEL_OPTIONS.get("ollama", []), sorted(live_models))
        cache_provider_catalog("ollama", models, source="live")
        clear_provider_issue("ollama")
        return {"provider": "ollama", "models": models, "current_model": provider_runtime_model("ollama"), "source": "live"}
    except HTTPException:
        raise
    except Exception as exc:
        cache_provider_catalog("ollama", fallback_models, source="static")
        raise HTTPException(status_code=503, detail=f"Ollama catalog unavailable: {exc}")


async def call_ollama(
    *,
    prompt: str,
    system: str,
    model: Optional[str] = None,
    max_tokens: int = MAX_TOKENS,
    timeout_seconds: float = 45.0,
) -> Dict[str, Any]:
    payload = {
        "model": (model or OLLAMA_MODEL).strip() or OLLAMA_MODEL,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": prompt},
        ],
        "stream": False,
        "keep_alive": "15m",
        "options": {
            "temperature": 0.35,
            "num_predict": max_tokens,
            "num_ctx": 1024,
        },
    }
    async with httpx.AsyncClient(timeout=timeout_seconds, follow_redirects=True) as client:
        resp = await client.post(f"{OLLAMA_HOST.rstrip('/')}/api/chat", json=payload)
    if resp.status_code != 200:
        log.error("Ollama API error: %s | %s", resp.status_code, resp.text[:400])
        raise HTTPException(status_code=resp.status_code, detail=f"Ollama API error: {resp.text[:300]}")
    return resp.json()


def cache_provider_catalog(provider: str, models: List[str], *, source: str = "live") -> None:
    _provider_catalog_cache[provider] = {
        "models": merge_model_options([provider_runtime_model(provider)], DEFAULT_PROVIDER_MODEL_OPTIONS.get(provider, []), models),
        "updated_at": now_iso(),
        "source": source,
    }


async def discover_provider_models(provider: str, api_key: Optional[str] = None) -> Dict[str, Any]:
    provider = (provider or "").strip().lower()
    if provider not in {"ollama", "anthropic", "openai", "gemini", "built_in"}:
        raise HTTPException(status_code=400, detail="Unsupported provider")

    current_model = provider_runtime_model(provider)
    if provider == "built_in":
        cache_provider_catalog("built_in", ["empire-fallback"], source="static")
        return {"provider": "built_in", "models": ["empire-fallback"], "current_model": "empire-fallback", "source": "static"}
    if provider == "ollama":
        return await discover_ollama_models()

    if provider == "anthropic":
        models = merge_model_options([current_model], DEFAULT_PROVIDER_MODEL_OPTIONS["anthropic"])
        cache_provider_catalog("anthropic", models, source="static")
        return {"provider": "anthropic", "models": models, "current_model": current_model, "source": "static"}

    key = (api_key if api_key is not None else provider_runtime_key(provider)).strip()
    if not key:
        raise HTTPException(status_code=400, detail=f"Paste a {provider.title()} API key first to fetch available models.")

    async with httpx.AsyncClient(timeout=45.0, follow_redirects=True) as client:
        if provider == "openai":
            response = await client.get(
                "https://api.openai.com/v1/models",
                headers={"Authorization": f"Bearer {key}"},
            )
            if response.status_code != 200:
                raise HTTPException(status_code=response.status_code, detail=f"OpenAI catalog error: {response.text[:280]}")
            payload = response.json()
            live_models = [
                row.get("id", "")
                for row in payload.get("data", [])
                if isinstance(row, dict) and re.match(r"^(gpt|o[1-9])[a-zA-Z0-9._-]*$", str(row.get("id", "")))
            ]
            preferred = DEFAULT_PROVIDER_MODEL_OPTIONS["openai"]
        else:
            response = await client.get(
                f"https://generativelanguage.googleapis.com/v1beta/models?key={key}",
            )
            if response.status_code != 200:
                raise HTTPException(status_code=response.status_code, detail=f"Gemini catalog error: {response.text[:280]}")
            payload = response.json()
            live_models = [
                str(row.get("name", "")).replace("models/", "")
                for row in payload.get("models", [])
                if isinstance(row, dict) and "generateContent" in row.get("supportedGenerationMethods", [])
            ]
            preferred = DEFAULT_PROVIDER_MODEL_OPTIONS["gemini"]

    models = merge_model_options([current_model], preferred, sorted(live_models))
    cache_provider_catalog(provider, models, source="live")
    clear_provider_issue(provider)
    return {"provider": provider, "models": models, "current_model": current_model, "source": "live"}


def provider_config() -> Dict[str, Dict[str, Any]]:
    def build_entry(provider: str, label_model: str, configured: bool) -> Dict[str, Any]:
        cached = _provider_catalog_cache.get(provider, {})
        with PROVIDER_STATUS_LOCK:
            cooldown_until = float(_provider_status.setdefault("cooldowns", {}).get(provider, 0) or 0)
        return {
            "configured": configured,
            "model": label_model,
            "model_options": merge_model_options([label_model], DEFAULT_PROVIDER_MODEL_OPTIONS.get(provider, []), cached.get("models", [])),
            "catalog_updated_at": cached.get("updated_at"),
            "catalog_source": cached.get("source", "default"),
            "cooldown_until": datetime.fromtimestamp(cooldown_until, UTC).isoformat() if cooldown_until else "",
        }

    return {
        "ollama": build_entry("ollama", OLLAMA_MODEL, ollama_provider_ready()),
        "anthropic": build_entry("anthropic", MODEL, bool(ANTHROPIC_API_KEY)),
        "openai": build_entry("openai", OPENAI_MODEL, bool(OPENAI_API_KEY)),
        "gemini": build_entry("gemini", GEMINI_MODEL, bool(GEMINI_API_KEY)),
        "built_in": build_entry("built_in", "empire-fallback", True),
    }


def provider_preference(state: Optional[Dict[str, Any]] = None) -> str:
    current = ensure_state_shape(state or get_state())
    settings = current.get("settings", {})
    preference = str(settings.get("provider_preference", "built_in") or "built_in").lower()
    if (
        str(settings.get("external_ai_mode", "manual_only") or "manual_only").lower() == "built_in_only"
        and preference in {"anthropic", "openai", "gemini"}
    ):
        return "built_in"
    return preference


def provider_order(preference: Optional[str] = None) -> List[str]:
    preferred = (preference or "openai").lower()
    order = ["ollama", "openai", "anthropic", "gemini", "built_in"]
    if preferred == "built_in":
        return ["built_in"]
    if preferred in ("ollama", "anthropic", "openai", "gemini"):
        return [preferred, "built_in"]
    return order


def active_provider_label(state: Optional[Dict[str, Any]] = None) -> str:
    config = provider_config()
    issues = provider_issues_snapshot()
    issue_names = {issue["provider"] for issue in issues}
    for provider_name in provider_order(provider_preference(state)):
        info = config.get(provider_name, {})
        if info.get("configured") and provider_name not in issue_names:
            if provider_name == "built_in":
                return "built-in"
            return f"{provider_name}/{info.get('model')}"
    if issues:
        primary = issues[0]["provider"]
        return f"built-in (fallback from {primary})"
    return "built-in"


def provider_issues_snapshot() -> List[Dict[str, str]]:
    with PROVIDER_STATUS_LOCK:
        errors = dict(_provider_status.get("errors", {}))
    ordered_names = [name for name in ("ollama", "openai", "anthropic", "gemini") if name in errors]
    ordered_names.extend(name for name in errors if name not in ordered_names)
    return [{"provider": name, "message": str(errors[name])} for name in ordered_names]


def provider_issue_snapshot() -> Optional[Dict[str, Any]]:
    issues = provider_issues_snapshot()
    if not issues:
        return None
    if len(issues) == 1:
        return dict(issues[0])
    names = ", ".join(issue["provider"] for issue in issues)
    return {
        "provider": names,
        "message": f"{names} are currently unavailable or quota-limited, so Leazy is using the local brain for continuity.",
        "details": issues,
        "created_at": now_iso(),
    }


def clear_provider_issue(provider: Optional[str] = None) -> None:
    with PROVIDER_STATUS_LOCK:
        errors = _provider_status.setdefault("errors", {})
        cooldowns = _provider_status.setdefault("cooldowns", {})
        if provider is None:
            errors.clear()
            cooldowns.clear()
        else:
            errors.pop(provider, None)
            cooldowns.pop(provider, None)


def set_provider_issue(provider: str, message: str) -> None:
    summary = message[:320]
    cooldown_seconds = 600 if re.search(r"\b429\b|quota|rate limit|insufficient", summary, re.IGNORECASE) else 180
    with PROVIDER_STATUS_LOCK:
        _provider_status.setdefault("errors", {})[provider] = message[:320]
        _provider_status.setdefault("cooldowns", {})[provider] = time.time() + cooldown_seconds


def provider_in_cooldown(provider: str) -> bool:
    with PROVIDER_STATUS_LOCK:
        cooldown_until = float(_provider_status.setdefault("cooldowns", {}).get(provider, 0) or 0)
    return cooldown_until > time.time()


def update_runtime_provider_config(provider: str, *, model: Optional[str] = None, api_key: Optional[str] = None) -> None:
    global ANTHROPIC_API_KEY, OPENAI_API_KEY, GEMINI_API_KEY, MODEL, OPENAI_MODEL, GEMINI_MODEL, OLLAMA_MODEL
    key_map = {
        "ollama": (None, "OLLAMA_MODEL"),
        "anthropic": ("ANTHROPIC_API_KEY", "ANTHROPIC_MODEL"),
        "openai": ("OPENAI_API_KEY", "OPENAI_MODEL"),
        "gemini": ("GEMINI_API_KEY", "GEMINI_MODEL"),
    }
    env_key, env_model = key_map.get(provider, (None, None))
    if not env_key and not env_model:
        return
    if api_key is not None:
        if provider == "anthropic":
            ANTHROPIC_API_KEY = api_key.strip()
        elif provider == "openai":
            OPENAI_API_KEY = api_key.strip()
        elif provider == "gemini":
            GEMINI_API_KEY = api_key.strip()
    if model is not None:
        if provider == "ollama":
            OLLAMA_MODEL = model.strip() or OLLAMA_MODEL
        elif provider == "anthropic":
            MODEL = model.strip() or MODEL
        elif provider == "openai":
            OPENAI_MODEL = model.strip() or OPENAI_MODEL
        elif provider == "gemini":
            GEMINI_MODEL = model.strip() or GEMINI_MODEL


def persist_provider_config(provider: str, *, model: Optional[str] = None, api_key: Optional[str] = None) -> None:
    key_map = {
        "ollama": (None, "OLLAMA_MODEL"),
        "anthropic": ("ANTHROPIC_API_KEY", "ANTHROPIC_MODEL"),
        "openai": ("OPENAI_API_KEY", "OPENAI_MODEL"),
        "gemini": ("GEMINI_API_KEY", "GEMINI_MODEL"),
    }
    env_key, env_model = key_map.get(provider, (None, None))
    if not env_key and not env_model:
        return
    if not ENV_PATH.exists():
        ENV_PATH.write_text("", encoding="utf-8")
    current_env = dotenv_values(ENV_PATH)
    if env_key and api_key is not None:
        secret = api_key.strip()
        if secret:
            set_key(str(ENV_PATH), env_key, secret)
            current_env[env_key] = secret
        else:
            unset_key(str(ENV_PATH), env_key)
            current_env.pop(env_key, None)
    if env_model:
        if model is not None and model.strip():
            set_key(str(ENV_PATH), env_model, model.strip())
            current_env[env_model] = model.strip()
        elif provider == "ollama" and not current_env.get(env_model):
            set_key(str(ENV_PATH), env_model, OLLAMA_MODEL)
            current_env[env_model] = OLLAMA_MODEL
    load_dotenv(ENV_PATH, override=True)


def ollama_status_payload() -> Dict[str, Any]:
    cached = _provider_catalog_cache.get("ollama", {})
    return {
        "available": ollama_provider_ready(),
        "host": OLLAMA_HOST,
        "active_model": OLLAMA_MODEL,
        "warm_ready": bool(_ollama_warm_state.get("ready")),
        "warming": bool(_ollama_warm_state.get("warming")),
        "last_attempt_at": _ollama_warm_state.get("last_attempt_at", ""),
        "last_error": _ollama_warm_state.get("last_error", ""),
        "catalog_updated_at": cached.get("updated_at", ""),
        "catalog_source": cached.get("source", "default"),
        "label": f"ollama/{OLLAMA_MODEL}",
    }


def warm_ollama_background() -> None:
    if not ollama_provider_ready() or _ollama_warm_state.get("ready") or _ollama_warm_state.get("warming"):
        return

    def _runner() -> None:
        _ollama_warm_state["warming"] = True
        _ollama_warm_state["last_attempt_at"] = now_iso()
        _ollama_warm_state["last_error"] = ""
        try:
            asyncio.run(
                call_ollama(
                    prompt="hello",
                    system="Reply with only READY.",
                    model=OLLAMA_MODEL,
                    max_tokens=8,
                    timeout_seconds=60.0,
                )
            )
            _ollama_warm_state["ready"] = True
        except Exception as exc:
            _ollama_warm_state["ready"] = False
            _ollama_warm_state["last_error"] = str(exc)[:280]
        finally:
            _ollama_warm_state["warming"] = False

    threading.Thread(target=_runner, name="ollama-warmup", daemon=True).start()


async def call_local_llm(*, system: str, prompt: str, max_tokens: int = MAX_TOKENS) -> Dict[str, Any]:
    if not _ollama_warm_state.get("ready"):
        warm_ollama_background()
        raise HTTPException(status_code=503, detail="Ollama local model is warming up")
    result = await call_ollama(
        prompt=prompt,
        system=system,
        model=OLLAMA_MODEL,
        max_tokens=max_tokens,
        timeout_seconds=10.0,
    )
    return {
        "text": extract_ollama_text(result),
        "model": OLLAMA_MODEL,
        "usage": {
            "prompt_eval_count": result.get("prompt_eval_count", 0),
            "eval_count": result.get("eval_count", 0),
        },
    }


async def call_openai(
    messages: List[Dict[str, str]],
    system: str,
    max_tokens: int = MAX_TOKENS,
) -> Dict[str, Any]:
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=503, detail="OPENAI_API_KEY not configured.")

    payload = {
        "model": OPENAI_MODEL,
        "messages": [{"role": "system", "content": system}, *messages],
        "max_tokens": max_tokens,
        "temperature": 0.6,
    }
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json",
    }
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        resp = await client.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    if resp.status_code != 200:
        log.error("OpenAI API error: %s | %s", resp.status_code, resp.text[:400])
        raise HTTPException(status_code=resp.status_code, detail=f"OpenAI API error: {resp.text[:300]}")
    return resp.json()


async def call_gemini(
    prompt: str,
    system: str,
    max_tokens: int = MAX_TOKENS,
) -> Dict[str, Any]:
    if not GEMINI_API_KEY:
        raise HTTPException(status_code=503, detail="GEMINI_API_KEY not configured.")

    payload = {
        "systemInstruction": {"parts": [{"text": system}]},
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.65,
            "maxOutputTokens": max_tokens,
        },
    }
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent?key={GEMINI_API_KEY}"
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        resp = await client.post(url, json=payload)
    if resp.status_code != 200:
        log.error("Gemini API error: %s | %s", resp.status_code, resp.text[:400])
        raise HTTPException(status_code=resp.status_code, detail=f"Gemini API error: {resp.text[:300]}")
    return resp.json()


async def call_anthropic_custom(
    messages: List[Dict[str, str]],
    system: str,
    *,
    api_key: str,
    model: str,
    max_tokens: int = MAX_TOKENS,
) -> Dict[str, Any]:
    headers = {
        "x-api-key": api_key.strip(),
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }
    payload: Dict[str, Any] = {
        "model": model.strip() or MODEL,
        "max_tokens": max_tokens,
        "system": system,
        "messages": messages,
    }
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        resp = await client.post("https://api.anthropic.com/v1/messages", headers=headers, json=payload)
    if resp.status_code != 200:
        raise HTTPException(status_code=resp.status_code, detail=f"Anthropic API error: {resp.text[:300]}")
    return resp.json()


async def call_openai_custom(
    messages: List[Dict[str, str]],
    system: str,
    *,
    api_key: str,
    model: str,
    max_tokens: int = MAX_TOKENS,
) -> Dict[str, Any]:
    payload = {
        "model": model.strip() or OPENAI_MODEL,
        "messages": [{"role": "system", "content": system}, *messages],
        "max_tokens": max_tokens,
        "temperature": 0.6,
    }
    headers = {
        "Authorization": f"Bearer {api_key.strip()}",
        "Content-Type": "application/json",
    }
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        resp = await client.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    if resp.status_code != 200:
        raise HTTPException(status_code=resp.status_code, detail=f"OpenAI API error: {resp.text[:300]}")
    return resp.json()


async def call_gemini_custom(
    prompt: str,
    system: str,
    *,
    api_key: str,
    model: str,
    max_tokens: int = MAX_TOKENS,
) -> Dict[str, Any]:
    payload = {
        "systemInstruction": {"parts": [{"text": system}]},
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.65,
            "maxOutputTokens": max_tokens,
        },
    }
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{(model.strip() or GEMINI_MODEL)}:generateContent?key={api_key.strip()}"
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        resp = await client.post(url, json=payload)
    if resp.status_code != 200:
        raise HTTPException(status_code=resp.status_code, detail=f"Gemini API error: {resp.text[:300]}")
    return resp.json()


def fallback_candidates(count: int = 3) -> List[Dict[str, Any]]:
    base = [
        {
            "name": "Aarav Sharma",
            "title": "Senior AI Engineer",
            "experience": 8,
            "fit_score": 96,
            "genesis_profile": "Strong systems builder with calm leadership and strong product sense.",
            "discovery_source": "GitHub signal + community references",
        },
        {
            "name": "Priya Menon",
            "title": "Lead Data Scientist",
            "experience": 7,
            "fit_score": 93,
            "genesis_profile": "High signal ML practitioner who translates research into business value.",
            "discovery_source": "Research network + technical writing footprint",
        },
        {
            "name": "Vikram Rao",
            "title": "Full-Stack Architect",
            "experience": 12,
            "fit_score": 91,
            "genesis_profile": "Architecture-minded generalist with strong ownership and scaling habits.",
            "discovery_source": "Passive talent synthesis + delivery pattern match",
        },
        {
            "name": "Neha Kulkarni",
            "title": "Platform Engineer",
            "experience": 6,
            "fit_score": 89,
            "genesis_profile": "Reliable infrastructure builder with strong documentation discipline.",
            "discovery_source": "Ops community footprint + conference talks",
        },
    ]
    return base[:count]


def extract_fallback_message(prompt: str) -> str:
    for marker in ("Creator request:", "Voice request:", "Operator request:", "Member request:", "User request:"):
        if marker in prompt:
            return prompt.split(marker, 1)[1].strip()
    lines = [line.strip() for line in prompt.splitlines() if line.strip()]
    for line in reversed(lines):
        if line.lower().startswith(("user:", "operator request:", "creator request:", "voice request:")):
            return line.split(":", 1)[1].strip()
    return lines[-1] if lines else ""


def normalize_workspace(workspace: Optional[str]) -> str:
    key = (workspace or "bridge").strip().lower()
    aliases = {
        "leazy": "bridge",
        "main": "bridge",
        "portal": "hq",
        "company": "hq",
        "agent_console": "agent",
        "prime_minister": "cabinet",
    }
    normalized = aliases.get(key, key)
    if normalized not in {"bridge", "hq", "agent", "ats", "network", "settings", "voice", "praapti", "cabinet"}:
        return "bridge"
    return normalized


def workspace_memory_brief(state: Dict[str, Any], workspace: str) -> List[str]:
    lines: List[str] = []
    hunts = state["praapti_hunts"]
    proposals = [row for row in state["nirmaan_proposals"] if not row.get("approved")]
    packages = state["packages"]
    cabinet = state.get("cabinet", {})
    prime_minister = cabinet.get("prime_minister", {})
    if hunts:
        latest_hunt = hunts[-1]
        lines.append(
            f"Latest hunt for {latest_hunt.get('client_company', COMPANY_NAME)}: "
            f"{(latest_hunt.get('job_description', '') or '')[:86]}"
        )
    if proposals:
        lines.append(f"Pending Nirmaan proposal: {proposals[-1].get('title', 'Untitled proposal')}")
    if packages:
        lines.append(f"Latest package: {packages[-1].get('title', 'Mission package')}")
    if workspace == "agent":
        candidate_total = sum(len(row.get("candidates", [])) for row in hunts)
        lines.append(f"Agent workspace tracks {len(hunts)} hunts and {candidate_total} candidate profiles.")
    elif workspace == "ats":
        lines.append(f"ATS workspace currently maps {len(hunts)} inferred jobs from recruitment hunts.")
    elif workspace == "network":
        clients = sorted({hunt.get("client_company", COMPANY_NAME) for hunt in hunts})
        if clients:
            lines.append(f"Network signals currently include {', '.join(clients[:4])}.")
    elif workspace == "hq":
        lines.append(f"HQ is carrying {len(packages)} packages, {len(proposals)} pending proposals, and {len(state['vault'])} vault memories.")
    elif workspace == "cabinet":
        lines.append(f"Prime Minister objective: {prime_minister.get('objective', 'No mandate set yet.')}")
        lines.append(f"Secretary swarm currently has {len(cabinet.get('secretaries', []))} active lanes.")
    if cabinet.get("mission_log"):
        lines.append(f"Latest Prime Minister cycle: {cabinet['mission_log'][-1].get('objective', '')[:88]}")
    return lines[:4]


def recruitment_seed_brief(workspace: str = "bridge", audience: str = "member", limit: int = 4) -> str:
    helper = RECRUITMENT_LAYER.get("seed_brief")
    if callable(helper):
        try:
            return helper(workspace=workspace, audience=audience, limit=limit)
        except Exception:
            pass
    fallback = [
        "TechBuzz is a recruitment company, so every answer should prioritize client intake, role clarity, candidate trust, and revenue discipline.",
        "Candidate data should stay privacy-filtered unless the user explicitly asks for deeper context.",
        "The AI can prepare Gmail, Teams, browser, and message actions, but it should never auto-send or act silently.",
    ]
    return "\n".join(f"- {line}" for line in fallback[:limit])


def global_recruitment_status() -> Dict[str, Any]:
    helper = GLOBAL_RECRUITMENT_LAYER.get("status_payload")
    if callable(helper):
        try:
            return helper()
        except Exception:
            pass
    return {
        "headline": "Global recruitment atlas is loading.",
        "atlas_mode": "fallback",
        "metrics": {
            "seeded_atoms": 0,
            "official_source_packs": 0,
            "country_nodes": 0,
            "subdivision_nodes": 0,
            "domains": 0,
            "role_families": 0,
            "candidate_axes": 0,
            "skill_bundles": 0,
            "possible_thoughts": 0,
        },
        "policy": [
            "Keep geography, candidate-intelligence, and role-family knowledge permanently available.",
            "Treat exact law and population claims as live-check data.",
        ],
        "sources": [],
    }


def global_recruitment_brief(query: str = "", limit: int = 8) -> str:
    query_key = (query or "").strip().lower()
    if query_key in {"hiring", "recruitment", "job", "jobs"}:
        status = global_recruitment_status()
        role_families = [row.get("title") for row in status.get("role_families", [])[:4] if row.get("title")]
        candidate_axes = [row.get("title") for row in status.get("candidate_axes", [])[:4] if row.get("title")]
        curated = [
            f"- Role families covered: {', '.join(role_families)}" if role_families else "",
            f"- Candidate lenses covered: {', '.join(candidate_axes)}" if candidate_axes else "",
            "- Guided disclosure is default, and exact laws or rights must be checked live before claiming they are current.",
        ]
        curated = [line for line in curated if line]
        if curated:
            return "\n".join(curated[:limit])
    helper = GLOBAL_RECRUITMENT_LAYER.get("context_brief")
    if callable(helper):
        try:
            result = helper(query=query, limit=limit)
            lines: List[str] = []
            seen = set()
            for raw_line in str(result or "").splitlines():
                line = raw_line.strip()
                if not line:
                    continue
                key = line.lower()
                if key in seen:
                    continue
                seen.add(key)
                lines.append(line)
            if lines:
                return "\n".join(lines[: limit + 4])
        except Exception:
            pass
    fallback = [
        "- Map geography and hiring context before judging fit.",
        "- Separate exact live compliance from permanent hiring heuristics.",
        "- Treat notice, intent, relevant experience, and capability as distinct dimensions.",
    ]
    return "\n".join(fallback[:limit])


def world_atlas_match_score(text: str) -> int:
    if not text:
        return 0
    terms = [
        "world", "map", "country", "countries", "region", "state", "district", "location", "population",
        "work culture", "people behaviour", "people behavior", "language", "culture", "religion", "gender",
        "rights", "rules", "law", "government job", "private job", "category", "grade", "hiring", "sourcing",
        "screening", "interview", "niche", "combined skills", "total experience", "relevant experience",
        "notice period", "serving notice", "not serving", "immediate", "active candidate", "passive candidate",
        "job change", "next job", "lifestyle", "learn", "build", "manage", "lead", "architect", "remember",
        "change", "motivate", "help them",
    ]
    return sum(1 for term in terms if term in text)


def recruitment_seed_brief(
    workspace: str = "bridge",
    audience: str = "member",
    limit: int = 4,
    extra_query: str = "",
) -> str:
    helper = RECRUITMENT_LAYER.get("seed_brief")
    if callable(helper):
        try:
            result = helper(workspace=workspace, audience=audience, limit=limit, extra_query=extra_query)
        except TypeError:
            try:
                result = helper(workspace=workspace, audience=audience, limit=limit)
            except Exception:
                result = ""
        except Exception:
            result = ""
        lines: List[str] = []
        seen = set()
        for raw_line in str(result or "").splitlines():
            line = raw_line.strip()
            if not line:
                continue
            key = line.lower()
            if key in seen:
                continue
            seen.add(key)
            lines.append(line)
        if lines:
            return "\n".join(lines[:limit])
    fallback = [
        "TechBuzz is a recruitment company, so every answer should prioritize client intake, role clarity, candidate trust, and revenue discipline.",
        "Candidate data should stay privacy-filtered unless the user explicitly asks for deeper context.",
        "The AI can prepare Gmail, Teams, browser, and message actions, but it should never auto-send or act silently.",
    ]
    return "\n".join(f"- {line}" for line in fallback[:limit])


def recruitment_learning_snapshot() -> Dict[str, Any]:
    helper = RECRUITMENT_LAYER.get("learning_health_snapshot")
    if callable(helper):
        try:
            return helper()
        except Exception:
            pass
    return {
        "seeded_entries": 0,
        "knowledge_entries": 0,
        "coverage_percent": 0,
        "loops": [],
        "brains": [],
    }


def fallback_model_text(prompt: str, *, state: Optional[Dict[str, Any]] = None, workspace: str = "bridge", source: str = "orb") -> str:
    lower = (prompt or "").lower()
    message = extract_fallback_message(prompt)
    message_lower = message.lower()
    current_state = ensure_state_shape(state or get_state())
    workspace = normalize_workspace(workspace)
    brain = brain_payload(current_state)
    dashboard = dashboard_payload(current_state)
    cabinet = prime_minister_payload(current_state)
    monitor = mother_monitor_payload(current_state)
    domains = operational_domains_payload(current_state)
    brief_lines = workspace_memory_brief(current_state, workspace)
    seed_query = message if len((message or "").split()) > 2 else ""
    seed_lines = recruitment_seed_brief(workspace, audience="member", limit=4, extra_query=seed_query)
    learning = recruitment_learning_snapshot()
    world_status = global_recruitment_status()
    world_metrics = world_status.get("metrics", {})
    atlas_brief = global_recruitment_brief(message or prompt, limit=8)
    atlas_score = world_atlas_match_score(message_lower or lower)
    latest_hunt = current_state["praapti_hunts"][-1] if current_state["praapti_hunts"] else None
    active_names = current_state["avatar_state"].get("active", ["RAMA"])
    active_avatar_text = " + ".join(name.title() for name in active_names)
    world_source_lines = "\n".join(
        f"- {source_row.get('title')}: {source_row.get('scope')}"
        for source_row in world_status.get("sources", [])[:4]
    ) or "- Official source packs are ready for live verification."
    issue = provider_issue_snapshot()
    provider_note = ""
    if issue:
        provider_names = str(issue.get("provider", "external providers"))
        provider_note = (
            f"Built-in reply: {provider_names} unavailable right now.\n"
        )
    if any(
        phrase in (message_lower or lower)
        for phrase in (
            "one short sentence",
            "short sentence",
            "one sentence",
            "one line",
            "short line",
            "short human line",
            "short recruiter summary",
            "recruiter summary",
            "short summary",
            "keep it short",
            "keep it human",
            "state your role",
            "what does this console do",
            "what do you do",
            "what can you help with",
        )
    ):
        concise_map = {
            "hq": "I guide visitors through TechBuzz services, hiring support, and the fastest next step.",
            "agent": "I help recruiters run sourcing, screening, trackers, submissions, and follow-ups from one console.",
            "network": "I track hiring signals, market movement, and connections so recruiters act faster.",
            "ats": "I keep roles, candidates, and stage movement organized so hiring stays on track.",
            "bridge": f"I coordinate {COMPANY_NAME} brains, recruitment workflows, and live system execution.",
        }
        return f"{provider_note}{concise_map.get(workspace, f'{AI_NAME} is ready and can help with the exact task in front of you.')}"
    if "return json only" in lower and "code_snippet" in lower:
        return json.dumps(
            {
                "title": "Mission Memory Relay",
                "description": "Adds a reusable relay for handing context between Praapti, Swarm, and Akshaya without growing the runtime footprint.",
                "code_snippet": "def mission_memory_relay(events):\n    return [event for event in events if event]\n",
            }
        )
    if "return exactly 3 candidates" in lower:
        return json.dumps(fallback_candidates(3))
    if "return exactly 4 candidates" in lower:
        return json.dumps(fallback_candidates(4))
    if "hidden company culture" in lower or "simulate the hidden company culture" in lower:
        return (
            "The team prefers builders who communicate clearly, move fast without drama, "
            "and can turn ambiguity into a repeatable operating system."
        )
    if "ideal candidate profile" in lower or "icp" in lower or "agni-shala" in lower:
        return (
            "Ideal profile: a product-minded technical builder with strong execution habits, "
            "clear communication, systems thinking, and evidence of scaling messy work."
        )
    if "swarm intelligence" in lower:
        return (
            "Praapti mapped talent lanes.\n"
            "Revenue refined positioning.\n"
            "Army secured execution paths.\n"
            "Research translated insights into experiments.\n"
            "Akshaya preserved the mission memory.\n"
            "Final outcome: the empire advanced with a concrete next-action list."
        )
    if "wake-word" in lower or "wake word" in lower:
        return "Wake mode armed. I will listen for 'Hey Jinn' or 'My King commands'."
    if "channeling the selected vishnu avatars" in lower or "vishnu avatars" in lower:
        names = infer_avatars_for_prompt(prompt, default=["RAMA"])
        return (
            f"{CORE_IDENTITY} now channels {' + '.join(name.title() for name in names)}. "
            "The wheel is luminous, the orb is transformed, and the empire is protected."
        )
    if "what you have learned" in message_lower or "what have you learned" in message_lower:
        latest_cycle = cabinet["mission_log"][0]["objective"] if cabinet["mission_log"] else "No cabinet cycle saved yet."
        top_domains = ", ".join(domain["name"] for domain in domains[:4]) or "No domains active yet"
        return (
            f"{provider_note}"
            f"I have {learning.get('seeded_entries', 0)} permanent recruitment memories, "
            f"{world_metrics.get('seeded_atoms', 0)} world-atlas atoms, and "
            f"{dashboard['metrics']['vault_items']} saved vault items.\n"
            f"Top active lanes: {top_domains}.\n"
            f"Latest cabinet focus: {latest_cycle}.\n"
            "Tell me the exact outcome you want and I will answer from that slice."
        )
    if atlas_score >= 2 or any(
        phrase in message_lower
        for phrase in (
            "always remember",
            "world map",
            "government job",
            "private job",
            "work culture",
            "people behavior",
            "people behaviour",
            "notice period",
            "total experience",
            "relevant experience",
            "active candidate",
            "passive candidate",
            "job change",
            "next job",
            "combined skills",
            "language",
            "religion",
            "rights",
            "district",
            "state",
            "population",
            "location",
            "region",
        )
    ):
        policy_lines = "\n".join(f"- {line}" for line in world_status.get("policy", [])[:3])
        return (
            f"{provider_note}"
            "That topic is already in the always-remember recruitment atlas.\n"
            f"Scale: {world_metrics.get('country_nodes', 0)} countries, {world_metrics.get('subdivision_nodes', 0)} location nodes, "
            f"and {world_metrics.get('possible_thoughts', 0):,} reasoning paths.\n"
            f"Relevant memory:\n{atlas_brief}\n"
            f"Rules:\n{policy_lines}\n"
            "Ask me one exact thing like notice period, relevant experience, work culture, or sourcing fit."
        )
    if message_lower in {"hi", "hello", "hey", "hey jinn", "namaste", "ji"}:
        context = "; ".join(brief_lines[:1]) or "No active memory yet."
        return (
            f"{provider_note}"
            f"{AI_NAME} is ready in {workspace.title()} mode. "
            f"Active avatar: {active_avatar_text}.\n"
            f"Current context: {context}\n"
            "Tell me the role, problem, or task and I’ll answer directly."
        )
    if any(word in message_lower for word in ("hiring", "hire", "recruit", "candidate", "recruitment", "job")):
        hunt_status = (
            f"The latest hunt is for {latest_hunt.get('client_company', COMPANY_NAME)} and has "
            f"{len(latest_hunt.get('candidates', []))} candidate profiles ready to review."
            if latest_hunt
            else "There is no active hunt yet, so the next useful step is to run Praapti on the role."
        )
        return (
            f"{provider_note}"
            f"Sure — I can help with hiring. {hunt_status}\n"
            "Send me the JD, client name, location, must-have skills, experience range, and notice period target. "
            "I’ll turn it into a sourcing plan, screening rubric, and tracker-ready shortlist."
        )
    if any(word in message_lower for word in ("hospital", "fire brigade", "firebrigade", "mess", "lab", "army", "navy", "air force", "airforce", "supply", "connectivity", "transmission")):
        domain_lines = "\n".join(
            f"- {domain['name']}: {domain['status']} | {domain['metric']}"
            for domain in domains[:8]
        )
        return (
            f"{provider_note}"
            "Those operational domains are live in the core.\n"
            f"{domain_lines}\n"
            "Open Systems or Monitor if you want the full relay view."
        )
    if any(word in message_lower for word in ("prime minister", "secretary", "cabinet", "revenue loop")):
        top_secretaries = cabinet["secretaries"][:5]
        secretary_lines = "\n".join(
            f"- {secretary['name']}: {secretary['status']} | {secretary['next_move']}"
            for secretary in top_secretaries
        )
        return (
            f"{provider_note}"
            f"{cabinet['prime_minister']['name']} is active.\n"
            f"Objective: {cabinet['prime_minister']['objective']}\n"
            f"Revenue forecast: INR {cabinet['revenue_board']['projected_revenue_inr']} Cr\n"
            f"Active secretaries: {cabinet['prime_minister']['active_secretaries']}\n"
            f"Top secretary lanes:\n{secretary_lines}\n"
            "Open the Cabinet panel if you want to run a cycle, change the mandate, or inspect the quantum storage manager."
        )
    if any(word in message_lower for word in ("brain", "evolution", "expand", "grow", "develop", "protect", "pillar")):
        pillars = "\n".join(
            f"- {pillar['label']}: {pillar['score']}% | {pillar['summary']}" for pillar in brain["pillars"]
        )
        recommendations = "\n".join(
            f"- {item['title']}: {item['action']}" for item in brain.get("recommendations", [])[:4]
        ) or "- No urgent action right now."
        return (
            f"{provider_note}"
            f"Brain status for {workspace.title()}:\n{pillars}\nNext best moves:\n{recommendations}"
        )
    if "what do you do" in message_lower or "who are you" in message_lower:
        return (
            f"{provider_note}"
            f"I am {AI_NAME}, the operating brain for {COMPANY_NAME}. "
            "I coordinate hiring, strategy, memory, improvement proposals, portal control, and safe package missions. "
            "If you give me a concrete goal, I will turn it into the next action plan."
        )
    if "avatar" in message_lower:
        return (
            f"{provider_note}"
            "Avatar modes change the working style.\n"
            "Matsya = memory, Kurma = stability, Krishna = strategy, Narasimha = defense, Kalki = renewal.\n"
            "Say `channel Krishna` or `channel Matsya + Kurma`."
        )
    if "settings" in message_lower or "model" in message_lower or "api key" in message_lower:
        return (
            f"{provider_note}"
            "Open Settings -> Provider Control.\n"
            "Choose provider, paste key, fetch models, save config, then use selected.\n"
            "If quota is unavailable, Ishani stays on the built-in brain."
        )
    if "learn" in message_lower or "improve" in message_lower or "upgrade" in message_lower:
        return (
            f"{provider_note}"
            "I improve through memory, atlas learning, Akshaya, Nirmaan, cabinet execution, and monitoring.\n"
            f"Current atlas scale: {world_metrics.get('seeded_atoms', 0)} atoms and {world_metrics.get('possible_thoughts', 0):,} reasoning paths.\n"
            "Tell me one area to improve next: recruitment, automation, interface, revenue, or monitoring."
        )
    if message:
        workspace_name = {
            "bridge": "Leazy Bridge",
            "hq": "TechBuzz HQ",
            "network": "Network",
            "ats": "ATS",
            "agent": "Agent Console",
            "settings": "Settings",
            "voice": "Voice",
            "praapti": "Praapti",
        }.get(workspace, "Leazy Bridge")
        context = "; ".join(brief_lines[:2]) or "No recent workspace context yet."
        return (
            f"{provider_note}"
            f"{workspace_name} is active. {brain['dominant_pillar']['label']} leads at {brain['dominant_pillar']['score']}%.\n"
            f"Context: {context}\n"
            "Tell me the exact output you want: shortlist, tracker, plan, report, mail draft, or next step."
        )
    return (
        f"{provider_note}{AI_NAME} is online. I can help with hiring, strategy, execution, coding, "
        "research, and building the next version of the empire."
    )


def compact_ollama_request(
    *,
    prompt: str,
    system: str,
    workspace: str,
    source: str,
    max_tokens: int,
) -> Dict[str, Any]:
    workspace_name = normalize_workspace(workspace)
    source_name = (source or "system").strip().lower()
    audience = "public" if source_name in {"public_hq", "public_agent"} else "member"
    message = extract_fallback_message(prompt).strip() or str(prompt or "").strip()
    message = re.sub(r"\s+\n", "\n", message)
    message = re.sub(r"\n{3,}", "\n\n", message)
    message_lower = message.lower()
    short_request = (
        len(message.split()) <= 20
        or any(
            cue in message_lower
            for cue in (
                "keep it short",
                "one short sentence",
                "one sentence",
                "one line",
                "state your role",
                "what do you do",
                "what can you help with",
            )
        )
    )
    seed_query = message if len(message.split()) > 2 else ""
    seed_lines = recruitment_seed_brief(
        workspace_name,
        audience=audience,
        limit=1 if short_request else 3,
        extra_query=seed_query,
    )
    workspace_memory = workspace_memory_brief(get_state(), workspace_name)
    workspace_lines = "\n".join(workspace_memory[:1 if short_request else 2]).strip()
    prompt_parts = [f"Operator request:\n{message[:420 if short_request else 1400] or 'Provide the next best step.'}"]
    if source_name == "public_hq":
        prompt_parts.append(f"HQ context:\n{public_hq_context_brief()[:260 if short_request else 700]}")
    elif workspace_lines:
        prompt_parts.append(f"Workspace context:\n{workspace_lines[:240 if short_request else 700]}")
    if seed_lines:
        prompt_parts.append(f"Relevant memory:\n{seed_lines[:260 if short_request else 1100]}")
    if source_name == "network_scan":
        prompt_parts.append("Format the answer as: Signals | Risks | Opportunities | Next Move.")
    if source_name == "agent_console":
        prompt_parts.append("Answer only from the visible ATS or tracker slice and avoid broad dumps.")
    compact_system = (
        f"You are {AI_NAME} for {COMPANY_NAME}. "
        "Reply like a capable human operator. Be direct, concise, practical, and grounded. "
        "If the user asks for one sentence, give one sentence. "
        "Do not sound ceremonial."
    )
    return {
        "system": compact_system[:900],
        "prompt": "\n\n".join(part for part in prompt_parts if part).strip()[:900 if short_request else 3200],
        "max_tokens": min(max_tokens, 72 if short_request else 220),
    }


async def generate_text(
    prompt: str,
    *,
    system: str,
    max_tokens: int = MAX_TOKENS,
    use_web_search: bool = False,
    workspace: str = "bridge",
    source: str = "system",
) -> Dict[str, Any]:
    source_name = (source or "system").strip().lower()
    preferred = provider_preference()
    if not external_ai_allowed_for_source(source_name) and preferred in {"anthropic", "openai", "gemini"}:
        preferred = "built_in"
    for provider_name in provider_order(preferred):
        if provider_name != "built_in" and provider_in_cooldown(provider_name):
            continue
        if provider_name == "ollama" and ollama_provider_ready():
            try:
                local_request = compact_ollama_request(
                    prompt=prompt,
                    system=system,
                    workspace=workspace,
                    source=source_name,
                    max_tokens=max_tokens,
                )
                result = await call_ollama(
                    prompt=local_request["prompt"],
                    system=local_request["system"],
                    model=OLLAMA_MODEL,
                    max_tokens=local_request["max_tokens"],
                    timeout_seconds=60.0,
                )
                text = extract_ollama_text(result).strip()
                if not text:
                    raise HTTPException(status_code=503, detail="Ollama local model returned an empty reply.")
                clear_provider_issue("ollama")
                return {
                    "text": text,
                    "provider": f"ollama/{OLLAMA_MODEL}",
                    "usage": {
                        "prompt_eval_count": result.get("prompt_eval_count", 0),
                        "eval_count": result.get("eval_count", 0),
                    },
                }
            except HTTPException as exc:
                detail = str(exc.detail or "").strip() or "Ollama local model timed out or is unavailable."
                set_provider_issue("ollama", detail)
                log.warning("Falling back after Ollama error: %s", detail)
            except Exception as exc:
                detail = str(exc).strip() or "Ollama local model timed out or returned no usable reply."
                set_provider_issue("ollama", detail)
                log.warning("Falling back after unexpected Ollama error: %s", detail)
        if provider_name == "anthropic" and ANTHROPIC_API_KEY:
            try:
                result = await call_anthropic(
                    messages=[{"role": "user", "content": prompt}],
                    system=system,
                    max_tokens=max_tokens,
                    use_web_search=use_web_search,
                )
                clear_provider_issue("anthropic")
                return {
                    "text": extract_text(result),
                    "provider": f"anthropic/{MODEL}",
                    "usage": result.get("usage", {}),
                }
            except HTTPException as exc:
                set_provider_issue("anthropic", exc.detail)
                log.warning("Falling back after Anthropic error: %s", exc.detail)
            except Exception as exc:
                set_provider_issue("anthropic", str(exc))
                log.warning("Falling back after unexpected Anthropic error: %s", exc)
        if provider_name == "openai" and OPENAI_API_KEY:
            try:
                result = await call_openai(
                    messages=[{"role": "user", "content": prompt}],
                    system=system,
                    max_tokens=max_tokens,
                )
                clear_provider_issue("openai")
                return {
                    "text": extract_openai_text(result),
                    "provider": f"openai/{OPENAI_MODEL}",
                    "usage": result.get("usage", {}),
                }
            except HTTPException as exc:
                set_provider_issue("openai", exc.detail)
                log.warning("Falling back after OpenAI error: %s", exc.detail)
            except Exception as exc:
                set_provider_issue("openai", str(exc))
                log.warning("Falling back after unexpected OpenAI error: %s", exc)
        if provider_name == "gemini" and GEMINI_API_KEY:
            try:
                result = await call_gemini(prompt=prompt, system=system, max_tokens=max_tokens)
                clear_provider_issue("gemini")
                return {
                    "text": extract_gemini_text(result),
                    "provider": f"gemini/{GEMINI_MODEL}",
                    "usage": result.get("usageMetadata", {}),
                }
            except HTTPException as exc:
                set_provider_issue("gemini", exc.detail)
                log.warning("Falling back after Gemini error: %s", exc.detail)
            except Exception as exc:
                set_provider_issue("gemini", str(exc))
                log.warning("Falling back after unexpected Gemini error: %s", exc)
    return {
        "text": fallback_model_text(prompt, state=get_state(), workspace=workspace, source=source),
        "provider": "built-in",
        "usage": {},
    }


def parse_json_blob(raw: str) -> Any:
    if not raw:
        return None
    candidates = [raw.strip(), raw.replace("```json", "").replace("```", "").strip()]
    for candidate in candidates:
        try:
            return json.loads(candidate)
        except Exception:
            pass
    for pattern in (r"\[[\s\S]+\]", r"\{[\s\S]+\}"):
        match = re.search(pattern, raw)
        if match:
            try:
                return json.loads(match.group(0))
            except Exception:
                pass
    return None


def normalized_avatar_names(raw_value: str) -> List[str]:
    raw_text = (raw_value or "").upper()
    chunks = re.split(r"[,+/&]| AND |\s+\+\s+", raw_text)
    selected: List[str] = []
    for chunk in chunks:
        name = re.sub(r"[^A-Z]", "", chunk)
        if not name:
            continue
        for avatar_name in DASHAVATARA:
            if name == avatar_name or name in avatar_name or avatar_name in name:
                if avatar_name not in selected:
                    selected.append(avatar_name)
                break
    return selected or ["RAMA"]


def infer_avatars_for_prompt(text: str, *, default: Optional[List[str]] = None) -> List[str]:
    prompt = (text or "").lower()
    scored: List[tuple[int, str]] = []
    for avatar_name, data in DASHAVATARA.items():
        score = 0
        for keyword in data["keywords"]:
            if keyword in prompt:
                score += 1
        if score:
            scored.append((score, avatar_name))
    scored.sort(reverse=True)
    chosen = [name for _, name in scored[:2]]
    if chosen:
        return chosen
    return list(default or ["RAMA"])


def avatar_profiles(names: List[str]) -> List[Dict[str, Any]]:
    profiles = []
    for name in names:
        data = DASHAVATARA.get(name)
        if not data:
            continue
        profiles.append({"key": name, "name": name.title(), **data})
    return profiles


def avatar_guidance(names: List[str]) -> str:
    profiles = avatar_profiles(names)
    if not profiles:
        profiles = avatar_profiles(["RAMA"])
    lines = [
        f"Core identity: {CORE_IDENTITY}.",
        "You may combine multiple avatar modes fluidly when that best serves the task.",
    ]
    for profile in profiles:
        lines.append(
            f"{profile['name'].title()}: {profile['power']} Mantra: {profile['mantra']}"
        )
    return "\n".join(lines)


def set_active_avatars(
    avatar_names: List[str],
    *,
    command: str,
    auto: bool = False,
) -> Dict[str, Any]:
    selected = avatar_names or ["RAMA"]

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        avatar_state = state["avatar_state"]
        avatar_state["active"] = selected
        avatar_state["manual_override"] = not auto
        avatar_state["last_channel"] = now_iso()
        avatar_state["protection_meter"] = min(100, max(18, avatar_state.get("protection_meter", 50) + (12 if auto else 16)))
        avatar_state["history"].append(
            {
                "avatars": selected,
                "command": command[:200],
                "auto": auto,
                "created_at": now_iso(),
            }
        )
        return avatar_state

    return mutate_state(_mutate)


def build_swarm_agents(state: Dict[str, Any]) -> List[Dict[str, str]]:
    hunts_today = sum(1 for row in state["praapti_hunts"] if row.get("created_day") == today_iso())
    pending_proposals = sum(1 for row in state["nirmaan_proposals"] if not row.get("approved"))
    swarm_count = len(state["swarm_missions"])
    vault_items = len(state["vault"])
    active_map = {
        "Praapti": "hunting" if hunts_today else "ready",
        "Revenue": "optimizing" if hunts_today else "forecasting",
        "Army": "patrolling" if swarm_count else "guarding",
        "Navy": "deployed" if swarm_count else "anchored",
        "Airforce": "scanning" if len(state["conversations"]) else "watching",
        "Research": "discovering" if pending_proposals else "thinking",
        "Armory": "upgrading" if pending_proposals else "sealed",
        "Warehouse": "stocked" if hunts_today else "awaiting",
        "Storage": "eternal" if vault_items else "forming",
        "Sewage": "flowing",
        "Maya": "campaigning" if hunts_today else "warming",
        "Anveshan": "inventing" if pending_proposals else "observing",
    }
    return [
        {"name": name, "role": role, "emoji": emoji, "status": active_map.get(name, "ready")}
        for name, role, emoji in SWARM_AGENT_NAMES
    ]


def parse_optional_iso(raw: str) -> Optional[datetime]:
    if not raw:
        return None
    try:
        parsed = datetime.fromisoformat(raw)
    except Exception:
        return None
    if parsed.tzinfo is None:
        return parsed.replace(tzinfo=UTC)
    return parsed.astimezone(UTC)


def state_file_size_kb() -> int:
    try:
        return max(1, int(round(STATE_PATH.stat().st_size / 1024)))
    except Exception:
        return 1


def operational_domains_payload(state: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    current = ensure_state_shape(state or get_state())
    prime_minister = current["cabinet"]["prime_minister"]
    hunts = current.get("praapti_hunts", [])
    packages = current.get("packages", [])
    proposals = [row for row in current.get("nirmaan_proposals", []) if not row.get("approved")]
    missions = current.get("swarm_missions", [])
    conversations = current.get("conversations", [])
    vault_items = current.get("vault", [])
    provider_issue = provider_issue_snapshot()
    protection_meter = int(current.get("avatar_state", {}).get("protection_meter", 0))
    last_voice = sanitize_operator_line(
        current.get("voice", {}).get("last_command", ""),
        "No recent voice relay recorded.",
    )
    accounts = global_accounts_snapshot()
    enabled = bool(prime_minister.get("enabled", True))
    objective = prime_minister.get("objective", f"Generate revenue for {COMPANY_NAME}.")
    projected_revenue = round(1.4 + len(hunts) * 0.38 + len(packages) * 0.28 + len(proposals) * 0.14, 2)

    def domain_status(domain_id: str) -> Dict[str, Any]:
        if domain_id == "accounts":
            return {
                "status": "learning" if accounts["entries"] else "ready",
                "priority": min(99, 62 + accounts["entries"] // 2 + accounts["reports"] * 4 + accounts["profiles"] * 3),
                "signal": accounts["headline"],
                "metric": f"{accounts['profiles']} profile(s), {accounts['entries']} entries, income {accounts['income_total']}, expense {accounts['expense_total']}.",
            }
        if domain_id == "management":
            return {
                "status": "commanding" if enabled else "paused",
                "priority": min(99, 74 + len(hunts) * 4 + len(packages) * 3),
                "signal": objective,
                "metric": f"Revenue track at INR {projected_revenue} Cr.",
            }
        if domain_id == "hospital":
            recovering = bool(provider_issue) or protection_meter < 78
            return {
                "status": "stabilizing" if recovering else "healthy",
                "priority": 82 if recovering else 61,
                "signal": provider_issue["message"] if provider_issue else "All major systems are stable and recoverable.",
                "metric": f"Protection meter {protection_meter}% with guardian {current['meta'].get('memory_guardian', 'eternal preservation')}.",
            }
        if domain_id == "firebrigade":
            hot = bool(provider_issue) or protection_meter < 74
            return {
                "status": "responding" if hot else "ready",
                "priority": 86 if hot else 58,
                "signal": "Escalation cover is focused on provider fallback, continuity, and defensive routing.",
                "metric": f"{len(missions)} swarm missions and {len(current['cabinet'].get('mission_log', []))} cabinet cycles tracked.",
            }
        if domain_id == "mess":
            return {
                "status": "serving" if conversations or hunts else "stocked",
                "priority": 56 + min(18, len(conversations) // 3),
                "signal": "Operational nourishment keeps every lane ready with context and follow-up.",
                "metric": f"{len(conversations)} conversation traces and {len(packages)} package briefs are available.",
            }
        if domain_id == "lab":
            return {
                "status": "experimenting" if proposals else "watching",
                "priority": 60 + min(26, len(proposals) * 7),
                "signal": proposals[0]["title"] if proposals else "No pending proposal. Nirmaan can generate a new one on demand.",
                "metric": f"{len(proposals)} live proposals and {len(vault_items)} preserved artifacts.",
            }
        if domain_id == "army":
            return {
                "status": "mobilized" if missions else "ready",
                "priority": 59 + min(24, len(missions) * 6),
                "signal": "Execution lanes are aligned to enforce movement on hunts, packages, and cabinet orders.",
                "metric": f"{len(missions)} swarm missions and {len(hunts)} active hunts are feeding action.",
            }
        if domain_id == "navy":
            return {
                "status": "routing" if packages else "harbor-ready",
                "priority": 55 + min(24, len(packages) * 6),
                "signal": "Delivery and distribution lanes are carrying bounded packages to the next checkpoint.",
                "metric": f"{len(packages)} bounded packages currently staged.",
            }
        if domain_id == "air_force":
            return {
                "status": "scanning" if hunts else "patrolling",
                "priority": 57 + min(24, len(hunts) * 6),
                "signal": "Signal scanning is reading market, client, and talent movement through live hunts.",
                "metric": f"{len(hunts)} hunts and {len(vault_items)} saved signals in Akshaya.",
            }
        if domain_id == "development":
            return {
                "status": "building" if proposals or packages else "ready",
                "priority": 60 + min(22, len(proposals) * 5 + len(packages) * 3),
                "signal": "The implementation forge is translating proposals and missions into operating surfaces.",
                "metric": f"{len(proposals)} proposals, {len(packages)} packages, {len(conversations)} conversation traces.",
            }
        if domain_id == "evolving":
            return {
                "status": "adapting",
                "priority": 62 + min(20, len(proposals) * 4 + len(missions) * 3),
                "signal": "Evolution is active through Nirmaan proposals, cabinet cycles, and preserved reports.",
                "metric": f"{len(current['cabinet'].get('mission_log', []))} cabinet cycles stored for learning.",
            }
        if domain_id == "upgrade":
            return {
                "status": "queued" if proposals else "ready",
                "priority": 58 + min(20, len(proposals) * 5),
                "signal": "Upgrade rail stages only bounded, reviewable changes so the core stays coherent.",
                "metric": f"{len(proposals)} proposal(s) awaiting approval.",
            }
        if domain_id == "supply_chain":
            return {
                "status": "supplying" if packages or hunts else "balanced",
                "priority": 57 + min(22, len(packages) * 4 + len(hunts) * 3),
                "signal": "Supply chain keeps jobs, candidates, packages, and reports flowing without dead ends.",
                "metric": f"{len(hunts)} hunts, {len(packages)} packages, {len(current['cabinet'].get('secretaries', []))} secretaries.",
            }
        if domain_id == "transmission":
            return {
                "status": "flowing",
                "priority": 63 + min(18, len(conversations) // 2 + len(vault_items) // 12),
                "signal": (
                    f"Latest relay: {last_voice}"
                    if current.get("voice", {}).get("last_command")
                    else "Keep message relay, engine handoff, and report propagation ready for the next operator command."
                ),
                "metric": f"{len(conversations)} conversations and {len(vault_items)} vault records are available for relay.",
            }
        return {
            "status": "linked",
            "priority": 56 + min(20, len(hunts) * 2 + len(packages) * 2),
            "signal": "Connectivity mesh is synchronizing client, portal, and memory lanes.",
            "metric": f"{len(hunts)} hunts, {len(packages)} packages, provider {active_provider_label(current)}.",
        }

    domains: List[Dict[str, Any]] = []
    for blueprint in OPERATIONAL_DOMAIN_BLUEPRINTS:
        runtime = domain_status(blueprint["id"])
        domains.append(
            {
                "id": blueprint["id"],
                "name": blueprint["name"],
                "lead": blueprint["lead"],
                "purpose": blueprint["purpose"],
                "status": runtime["status"],
                "priority": runtime["priority"],
                "latest_signal": sanitize_operator_line(
                    runtime["signal"],
                    "Awaiting the next clean domain signal.",
                ),
                "metric": sanitize_operator_line(runtime["metric"], "Metric pending."),
                "updated_at": current["monitoring"].get("last_scan_at", now_iso()),
            }
        )
    domains.sort(key=lambda item: item["priority"], reverse=True)
    return domains


def nervous_system_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    domains = operational_domains_payload(current)
    monitoring = current.get("monitoring", {})
    cabinet = current.get("cabinet", {})
    prime_minister = cabinet.get("prime_minister", {})
    issue = provider_issue_snapshot()
    client_molecules = sorted({hunt.get("client_company", COMPANY_NAME) for hunt in current.get("praapti_hunts", [])})
    accounts = global_accounts_snapshot()
    repair_engine = auto_repair_engine_payload(current)
    nodes = [
        {"id": "mother_brain", "name": f"{CORE_IDENTITY} Core", "kind": "mother_brain", "status": "live"},
        {"id": "state_db", "name": "Quantum State Database", "kind": "database", "status": "synced"},
        {"id": "akshaya", "name": "Akshaya Vault", "kind": "memory", "status": "preserving"},
        {"id": "cabinet", "name": prime_minister.get("name", "Prime Minister"), "kind": "governance", "status": prime_minister.get("status", "governing")},
        {"id": "praapti", "name": "Praapti", "kind": "recruitment", "status": "active" if current.get("praapti_hunts") else "ready"},
        {"id": "nirmaan", "name": "Nirmaan Chakra", "kind": "evolution", "status": "active" if current.get("nirmaan_proposals") else "ready"},
        {"id": "swarm", "name": "Swarm Intelligence", "kind": "execution", "status": "active" if current.get("swarm_missions") else "ready"},
        {"id": "voice", "name": "Voice Relay", "kind": "interface", "status": "listening" if current.get("voice", {}).get("always_listening") else "idle"},
        {"id": "accounts", "name": "Accounts Command", "kind": "finance", "status": "learning" if accounts["entries"] else "ready"},
        {"id": "repair_engine", "name": "Auto Repair Engine", "kind": "repair", "status": repair_engine.get("status", "ready")},
    ] + [
        {"id": domain["id"], "name": domain["name"], "kind": "domain", "status": domain["status"]}
        for domain in domains
    ]
    links = [
        {"from": "state_db", "to": "akshaya", "label": "preserve"},
        {"from": "state_db", "to": "cabinet", "label": "brief"},
        {"from": "cabinet", "to": "praapti", "label": "recruit"},
        {"from": "cabinet", "to": "nirmaan", "label": "evolve"},
        {"from": "cabinet", "to": "swarm", "label": "execute"},
        {"from": "cabinet", "to": "accounts", "label": "govern"},
        {"from": "voice", "to": "cabinet", "label": "command"},
        {"from": "accounts", "to": "state_db", "label": "ledger"},
        {"from": "mother_brain", "to": "repair_engine", "label": "self-check"},
        {"from": "repair_engine", "to": "state_db", "label": "repair"},
    ] + [
        {"from": "cabinet", "to": domain["id"], "label": "assign"}
        for domain in domains[:8]
    ] + [
        {"from": domain["id"], "to": "state_db", "label": "report"}
        for domain in domains[:8]
    ]
    transmissions: List[Dict[str, str]] = []
    if current.get("voice", {}).get("last_command"):
        transmissions.append(
            {
                "from": "voice",
                "to": "cabinet",
                "message": sanitize_operator_line(
                    current["voice"]["last_command"][:140],
                    "Voice relay captured.",
                ),
                "created_at": current["voice"].get("updated_at", now_iso()),
            }
        )
    if cabinet.get("mission_log"):
        last_cycle = cabinet["mission_log"][-1]
        transmissions.append(
            {
                "from": "cabinet",
                "to": "state_db",
                "message": last_cycle.get("objective", "")[:140],
                "created_at": last_cycle.get("created_at", now_iso()),
            }
        )
    if current.get("praapti_hunts"):
        last_hunt = current["praapti_hunts"][-1]
        transmissions.append(
            {
                "from": "praapti",
                "to": "cabinet",
                "message": (last_hunt.get("job_description", "")[:120] or "Latest hunt saved."),
                "created_at": last_hunt.get("created_at", now_iso()),
            }
        )
    if current.get("nirmaan_proposals"):
        last_proposal = current["nirmaan_proposals"][-1]
        transmissions.append(
            {
                "from": "nirmaan",
                "to": "cabinet",
                "message": last_proposal.get("title", "Proposal staged."),
                "created_at": last_proposal.get("created_at", now_iso()),
            }
        )
    if current.get("swarm_missions"):
        last_mission = current["swarm_missions"][-1]
        transmissions.append(
            {
                "from": "swarm",
                "to": "state_db",
                "message": last_mission.get("mission", "Mission report saved.")[:140],
                "created_at": last_mission.get("created_at", now_iso()),
            }
        )
    latest_account_report = db_one("SELECT title, created_at FROM account_reports ORDER BY created_at DESC LIMIT 1")
    if latest_account_report:
        transmissions.append(
            {
                "from": "accounts",
                "to": "cabinet",
                "message": latest_account_report.get("title", "Accounts Command report"),
                "created_at": latest_account_report.get("created_at", now_iso()),
            }
        )
    if repair_engine.get("last_report"):
        transmissions.append(
            {
                "from": "repair_engine",
                "to": "mother_brain",
                "message": sanitize_operator_line(
                    repair_engine.get("last_report", ""),
                    "Safe self-check layer is ready.",
                )[:140],
                "created_at": repair_engine.get("last_check_at", now_iso()),
            }
        )
    alerts: List[Dict[str, str]] = []
    if issue:
        alerts.append({"level": "warning", "title": "External Provider Fallback", "summary": issue["message"]})
    if not current.get("praapti_hunts"):
        alerts.append({"level": "info", "title": "Recruitment Feed Idle", "summary": "No hunts are active yet, so the talent network is waiting for a role brief."})
    if int(current.get("avatar_state", {}).get("protection_meter", 0)) < 78:
        alerts.append({"level": "warning", "title": "Protection Below Stronghold", "summary": "Raise Matsya + Kurma or run a brain pulse focused on protection."})
    organ_systems = [
        {
            "id": "cortex",
            "name": f"{CORE_IDENTITY} Cortex",
            "role": "Reasoning, orchestration, cross-system synthesis",
            "status": "live",
            "load": min(99, 54 + len(current.get("conversations", [])) * 2 + len(current.get("nirmaan_proposals", [])) * 4),
        },
        {
            "id": "conductor",
            "name": "Cabinet Conductor",
            "role": "Prime Minister rhythm, governance pulse, and revenue direction",
            "status": prime_minister.get("status", "governing"),
            "load": min(99, 50 + len(cabinet.get("mission_log", [])) * 5 + len(cabinet.get("secretaries", []))),
        },
        {
            "id": "hippocampus",
            "name": "Akshaya Hippocampus",
            "role": "Long-memory preservation, replay, archival recall",
            "status": "preserving",
            "load": min(99, 40 + len(current.get("vault", [])) // 2 + len(current.get("conversations", []))),
        },
        {
            "id": "lungs",
            "name": "Transmission Lungs",
            "role": "Signal breathing between portals, agents, and the state core",
            "status": "flowing",
            "load": min(99, 46 + len(transmissions) * 6 + len(domains)),
        },
        {
            "id": "eyes",
            "name": "Praapti Eyes",
            "role": "Search, scouting, recruitment sensing, opportunity vision",
            "status": "active" if current.get("praapti_hunts") else "watching",
            "load": min(99, 34 + len(current.get("praapti_hunts", [])) * 11),
        },
        {
            "id": "ears",
            "name": "Voice Ears",
            "role": "Wake words, command intake, auditory relay",
            "status": "listening" if current.get("voice", {}).get("always_listening") else "idle",
            "load": min(99, 30 + (18 if current.get("voice", {}).get("always_listening") else 0) + len(current.get("voice", {}).get("wake_words", [])) * 8),
        },
        {
            "id": "hands",
            "name": "Swarm Hands",
            "role": "Execution, bounded packages, mission delivery",
            "status": "active" if current.get("swarm_missions") or current.get("packages") else "ready",
            "load": min(99, 36 + len(current.get("swarm_missions", [])) * 10 + len(current.get("packages", [])) * 8),
        },
        {
            "id": "metabolism",
            "name": "Accounts Metabolism",
            "role": "Cashflow digestion, tax routing, and ledger discipline",
            "status": "learning" if accounts["entries"] else "warming",
            "load": min(99, 34 + accounts["entries"] + accounts["reports"] * 3 + accounts["profiles"] * 5),
        },
        {
            "id": "skin",
            "name": "Protection Skin",
            "role": "Boundary defense, shield response, runtime hardening",
            "status": "hardening" if alerts else "stable",
            "load": max(42, int(current.get("avatar_state", {}).get("protection_meter", 0))),
        },
        {
            "id": "immune",
            "name": "Auto Repair Immune Layer",
            "role": "Safe self-checks, state sanitation, UI/UX repair overlay, and continuity recovery.",
            "status": repair_engine.get("status", "ready"),
            "load": min(99, 34 + int(repair_engine.get("issue_count", 0) or 0) * 8 + len(repair_engine.get("issues", [])) * 3),
        },
    ]
    cell_clusters = [
        {
            "name": domain["name"],
            "identity": domain["lead"],
            "cell_count": 1500 + int(domain.get("priority", 0)) * 14,
            "status": domain["status"],
            "signal": domain["latest_signal"],
        }
        for domain in domains
    ] + [
        {
            "name": "Conversation Neurons",
            "identity": "Orb + Agent + Portals",
            "cell_count": 600 + len(current.get("conversations", [])) * 22,
            "status": "active" if current.get("conversations") else "warming",
            "signal": "Live dialogue is feeding memory and execution lanes.",
        },
        {
            "name": "Vault Cells",
            "identity": "Akshaya",
            "cell_count": 700 + len(current.get("vault", [])) * 18,
            "status": "preserving",
            "signal": "Preserved items are being compacted and indexed for recall.",
        },
    ]
    reflex_arcs = [
        {
            "trigger": "Provider quota or outage",
            "response": "Switch to built-in brain, keep the workflow alive, and log the provider issue for manual recovery.",
        },
        {
            "trigger": "Protection drops below stronghold",
            "response": "Raise Matsya + Kurma, tighten shield posture, and increase preservation priority.",
        },
        {
            "trigger": "New hunt or package enters the system",
            "response": "Brief the cabinet, open the matching domain lanes, and write the report back to Akshaya.",
        },
        {
            "trigger": "Voice wake command arrives",
            "response": "Route through the auditory relay, assign avatars, and publish the result to the live chain.",
        },
    ]
    telemetry = {
        "cells_active": sum(cluster["cell_count"] for cluster in cell_clusters),
        "organs_online": sum(1 for organ in organ_systems if organ["status"] not in {"idle", "paused"}),
        "signal_integrity": min(99, 58 + len(links) + len(transmissions) * 3),
        "circulation_index": min(99, 48 + len(domains) * 2 + len(transmissions) * 5),
        "adaptation_rate": min(99, 42 + len(current.get("nirmaan_proposals", [])) * 8 + len(current.get("swarm_missions", [])) * 5),
    }
    if len(current.get("nirmaan_proposals", [])) >= 3 or len(current.get("swarm_missions", [])) >= 2:
        active_allotrope = "diamond-mind"
    elif len(current.get("packages", [])) >= 1 or len(transmissions) >= 5:
        active_allotrope = "nanotube-mind"
    elif len(domains) >= 10 or len(transmissions) >= 3:
        active_allotrope = "graphene-mind"
    else:
        active_allotrope = "carbon-seed"
    carbon_protocol = {
        "bond_layer": "active",
        "allotrope": active_allotrope,
        "bond_integrity": min(99, 61 + len(links) + len(transmissions) * 4),
        "signals_per_minute": min(999, 48 + len(transmissions) * 22 + len(domains) * 9 + len(cabinet.get("mission_log", [])) * 3),
        "molecules": [
            {"name": "Client Molecules", "count": max(1, len(client_molecules)), "role": "Revenue and delivery bonds"},
            {"name": "Execution Molecules", "count": max(1, len(current.get("packages", [])) + len(current.get("swarm_missions", []))), "role": "Operational action chains"},
            {"name": "Memory Molecules", "count": max(1, len(current.get("vault", [])) + len(current.get("conversations", []))), "role": "Preservation and recall compounds"},
            {"name": "Evolution Molecules", "count": max(1, len(current.get("nirmaan_proposals", []))), "role": "Design and adaptation compounds"},
            {"name": "Ledger Molecules", "count": max(1, accounts["entries"] + accounts["profiles"]), "role": "Tax and cashflow compounds"},
        ],
        "allotropes": [
            {"name": "Graphene-Mind", "role": "Parallel global sensing and wide distributed analysis."},
            {"name": "Diamond-Mind", "role": "Focused hard-problem reasoning with maximum structural coherence."},
            {"name": "Nanotube-Mind", "role": "Frictionless delivery of signals, logistics, and command flow."},
        ],
    }
    prana = prana_nadi_payload(current)
    component_intelligence = component_intelligence_payload(current, domains)
    hierarchy = brain_hierarchy_payload(current, domains, component_intelligence)
    telemetry["brains_active"] = hierarchy["summary"]["total_brains"]
    telemetry["permission_relays"] = hierarchy["summary"]["permission_relays"]
    return {
        "mother": {
            "name": f"{CORE_IDENTITY} Core",
            "heartbeat_ms": int(monitoring.get("pulse_interval_ms", 2500)),
            "last_scan_at": monitoring.get("last_scan_at", now_iso()),
            "state_file": str(STATE_PATH),
        },
        "nodes": nodes,
        "links": links,
        "transmissions": transmissions[:12],
        "alerts": alerts,
        "organ_systems": organ_systems,
        "cell_clusters": cell_clusters[:18],
        "reflex_arcs": reflex_arcs,
        "telemetry": telemetry,
        "carbon_protocol": carbon_protocol,
        "auto_repair": repair_engine,
        "prana_nadi": prana,
        "component_intelligence": component_intelligence,
        "hierarchy": hierarchy,
    }


def memory_audit_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    cabinet = current.get("cabinet", {})
    storage = cabinet.get("quantum_storage", {})
    counts = {
        "vault": len(current.get("vault", [])),
        "conversations": len(current.get("conversations", [])),
        "hunts": len(current.get("praapti_hunts", [])),
        "swarm_missions": len(current.get("swarm_missions", [])),
        "nirmaan_proposals": len(current.get("nirmaan_proposals", [])),
        "packages": len(current.get("packages", [])),
        "cabinet_cycles": len(cabinet.get("mission_log", [])),
    }
    issues: List[str] = []
    if not STATE_PATH.exists():
        issues.append("State file missing.")
    if counts["vault"] == 0:
        issues.append("No preserved vault items yet.")
    if counts["conversations"] == 0:
        issues.append("No conversation memory has been saved yet.")
    healthy = not issues
    return {
        "status": "healthy" if healthy else "attention",
        "healthy": healthy,
        "state_file": str(STATE_PATH),
        "footprint_kb": state_file_size_kb(),
        "guardian": current["meta"].get("memory_guardian", "eternal preservation"),
        "seal": storage.get("seal", "Matsya + Kurma"),
        "mode": storage.get("mode", "elastic preservation"),
        "accuracy_mode": "structured summaries plus source preservation",
        "archive_cycles": storage.get("archive_cycles", 0),
        "last_compaction_at": storage.get("last_compaction_at", ""),
        "last_guardian_cycle": current["meta"].get("last_guardian_cycle", ""),
        "counts": counts,
        "issues": issues,
        "repairs": [
            "State shape is enforced on every read and write.",
            "Memory guardian trims overflow and keeps high-value summaries.",
            "Prime Minister cycles, hunts, packages, and proposals are preserved in Akshaya.",
        ],
    }


def extract_document_text(path: Path, mime_type: str = "") -> str:
    extension = path.suffix.lower()
    try:
        if extension in {".txt", ".md", ".py", ".json", ".html", ".css", ".js", ".csv"}:
            if extension == ".csv":
                rows = []
                with path.open("r", encoding="utf-8", errors="ignore") as handle:
                    for index, row in enumerate(csv_reader(handle)):
                        rows.append(" | ".join(row))
                        if index >= 40:
                            break
                return "\n".join(rows)
            return path.read_text(encoding="utf-8", errors="ignore")
        if extension == ".pdf" and PdfReader:
            reader = PdfReader(str(path))
            text_bits = []
            for page in reader.pages[:20]:
                text_bits.append(page.extract_text() or "")
            return "\n".join(text_bits)
        if extension == ".docx" and DocxDocument:
            doc = DocxDocument(str(path))
            return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
        if extension in {".xlsx", ".xlsm"} and load_workbook:
            workbook = load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for sheet in workbook.worksheets[:4]:
                lines.append(f"[Sheet] {sheet.title}")
                for idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        lines.append(" | ".join(values))
                    if idx >= 25:
                        break
            workbook.close()
            return "\n".join(lines)
        if extension == ".pptx" and Presentation:
            slides = []
            presentation = Presentation(str(path))
            for slide_index, slide in enumerate(presentation.slides[:20], start=1):
                slides.append(f"[Slide {slide_index}]")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        slides.append(text)
            return "\n".join(slides)
        if extension == ".zip":
            names = []
            with zipfile.ZipFile(path) as bundle:
                names = bundle.namelist()[:100]
            return "Archive contents:\n" + "\n".join(names)
    except Exception as exc:
        return f"Document parser warning: {exc}"
    if mime_type.startswith("text/"):
        return path.read_text(encoding="utf-8", errors="ignore")
    return ""


def document_summary(text: str, fallback_name: str) -> str:
    clean = " ".join((text or "").split())
    if clean:
        return clean[:280]
    return f"{fallback_name} uploaded into the Ishani document studio."


def create_document_record(
    *,
    user_id: str,
    original_name: str,
    mime_type: str,
    extension: str,
    storage_path: Path,
    size_bytes: int,
    extracted_text: str,
    kind: str = "upload",
) -> Dict[str, Any]:
    record = {
        "id": new_id("doc"),
        "user_id": user_id,
        "original_name": original_name,
        "mime_type": mime_type or "application/octet-stream",
        "extension": extension or storage_path.suffix.lower(),
        "storage_path": str(storage_path),
        "size_bytes": int(size_bytes),
        "extracted_text": (extracted_text or "")[:50000],
        "summary": document_summary(extracted_text, original_name),
        "kind": kind,
        "created_at": now_iso(),
        "updated_at": now_iso(),
    }
    db_exec(
        """
        INSERT INTO documents(id,user_id,original_name,mime_type,extension,storage_path,size_bytes,extracted_text,summary,kind,created_at,updated_at)
        VALUES(?,?,?,?,?,?,?,?,?,?,?,?)
        """,
        (
            record["id"],
            record["user_id"],
            record["original_name"],
            record["mime_type"],
            record["extension"],
            record["storage_path"],
            record["size_bytes"],
            record["extracted_text"],
            record["summary"],
            record["kind"],
            record["created_at"],
            record["updated_at"],
        ),
    )
    return record


def user_documents(user_id: str, limit: int = 30) -> List[Dict[str, Any]]:
    return db_all(
        "SELECT id, original_name, mime_type, extension, size_bytes, summary, kind, created_at, updated_at FROM documents WHERE user_id=? ORDER BY created_at DESC LIMIT ?",
        (user_id, limit),
    )


def normalize_region_code(region_code: Optional[str]) -> str:
    candidate = (region_code or "IN").strip().upper()
    return candidate if candidate in ACCOUNT_REGION_PRESETS else "GLOBAL"


def account_region_preset(region_code: Optional[str]) -> Dict[str, Any]:
    return ACCOUNT_REGION_PRESETS[normalize_region_code(region_code)]


def account_region_catalog() -> List[Dict[str, Any]]:
    return [dict(preset) for preset in ACCOUNT_REGION_PRESETS.values()]


def ensure_account_profile(user_id: str) -> Dict[str, Any]:
    existing = db_one("SELECT * FROM account_profiles WHERE user_id=?", (user_id,))
    if existing:
        return existing
    preset = account_region_preset("IN")
    record = {
        "user_id": user_id,
        "region_code": preset["code"],
        "business_name": f"{COMPANY_NAME} Pvt Ltd",
        "business_type": "services",
        "currency": preset["currency"],
        "tax_name": preset["tax_name"],
        "default_tax_rate": float(preset["suggested_tax_rate"]),
        "filing_cycle": preset["filing_cycle"],
        "tax_registration": "",
        "notes": preset["notes"],
        "updated_at": now_iso(),
    }
    db_exec(
        """
        INSERT INTO account_profiles(user_id,region_code,business_name,business_type,currency,tax_name,default_tax_rate,filing_cycle,tax_registration,notes,updated_at)
        VALUES(?,?,?,?,?,?,?,?,?,?,?)
        """,
        (
            record["user_id"],
            record["region_code"],
            record["business_name"],
            record["business_type"],
            record["currency"],
            record["tax_name"],
            record["default_tax_rate"],
            record["filing_cycle"],
            record["tax_registration"],
            record["notes"],
            record["updated_at"],
        ),
    )
    return db_one("SELECT * FROM account_profiles WHERE user_id=?", (user_id,)) or record


def account_entries_for_user(user_id: str, limit: int = 60) -> List[Dict[str, Any]]:
    return db_all(
        """
        SELECT id, entry_type, category, amount, tax_percent, tax_amount, total_amount, currency, counterparty, description, source, occurred_on, created_at
        FROM account_entries
        WHERE user_id=?
        ORDER BY occurred_on DESC, created_at DESC
        LIMIT ?
        """,
        (user_id, limit),
    )


def account_reports_for_user(user_id: str, limit: int = 12) -> List[Dict[str, Any]]:
    rows = db_all(
        "SELECT id, title, summary, payload_json, created_at FROM account_reports WHERE user_id=? ORDER BY created_at DESC LIMIT ?",
        (user_id, limit),
    )
    for row in rows:
        try:
            row["payload"] = json.loads(row.get("payload_json") or "{}")
        except Exception:
            row["payload"] = {}
    return rows


def create_account_report(user_id: str, title: str, summary: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    report = {
        "id": new_id("acctrep"),
        "user_id": user_id,
        "title": title[:160],
        "summary": summary[:900],
        "payload_json": json.dumps(payload, ensure_ascii=False)[:12000],
        "created_at": now_iso(),
    }
    db_exec(
        "INSERT INTO account_reports(id,user_id,title,summary,payload_json,created_at) VALUES(?,?,?,?,?,?)",
        (
            report["id"],
            report["user_id"],
            report["title"],
            report["summary"],
            report["payload_json"],
            report["created_at"],
        ),
    )
    return report


def normalize_account_entry_type(entry_type: str) -> str:
    raw = (entry_type or "").strip().lower()
    return ACCOUNT_ENTRY_TYPE_ALIASES.get(raw, raw)


def infer_account_category(entry_type: str, description: str) -> str:
    normalized = normalize_account_entry_type(entry_type)
    text = (description or "").strip().lower()
    if any(token in text for token in ["gst", "tds", "vat", "tax"]):
        return "tax"
    if any(token in text for token in ["salary", "payroll", "stipend"]):
        return "payroll"
    if any(token in text for token in ["recruit", "placement", "consult", "talent", "staffing", "hiring"]):
        return "recruitment_consulting"
    if any(token in text for token in ["software", "cloud", "license", "saas", "hosting"]):
        return "software"
    if any(token in text for token in ["training", "course", "bootcamp", "learning"]):
        return "training"
    if normalized in ACCOUNT_INFLOW_TYPES:
        return "income"
    if normalized in ACCOUNT_OUTFLOW_TYPES:
        return "expense"
    return "services"


def account_gst_rate_for_category(category: str, fallback_rate: float = 18.0) -> float:
    key = (category or "").strip().lower()
    return float(ACCOUNT_GST_CATEGORY_RATES.get(key, fallback_rate if fallback_rate is not None else 18.0))


def calculate_gst_breakdown(amount: float, rate: float, supply_type: str = "intra") -> Dict[str, float]:
    base_amount = round(max(0.0, float(amount or 0.0)), 2)
    safe_rate = round(max(0.0, min(100.0, float(rate or 0.0))), 2)
    supply = (supply_type or "intra").strip().lower()
    if supply == "inter":
        igst = round(base_amount * safe_rate / 100.0, 2)
        return {
            "rate": safe_rate,
            "cgst": 0.0,
            "sgst": 0.0,
            "igst": igst,
            "tax_total": igst,
            "invoice_total": round(base_amount + igst, 2),
        }
    split_rate = round(safe_rate / 2.0, 2)
    cgst = round(base_amount * split_rate / 100.0, 2)
    sgst = round(base_amount * split_rate / 100.0, 2)
    return {
        "rate": safe_rate,
        "cgst": cgst,
        "sgst": sgst,
        "igst": 0.0,
        "tax_total": round(cgst + sgst, 2),
        "invoice_total": round(base_amount + cgst + sgst, 2),
    }


def detect_tds_obligation(entry_type: str, description: str, amount: float, section: str = "", party_pan: str = "") -> Dict[str, Any]:
    normalized = normalize_account_entry_type(entry_type)
    safe_amount = round(max(0.0, float(amount or 0.0)), 2)
    if normalized in ACCOUNT_INFLOW_TYPES:
        return {
            "applicable": False,
            "section": "",
            "rate": 0.0,
            "amount": 0.0,
            "tds_amount": 0.0,
            "net_payable": safe_amount,
            "reason": "TDS is usually evaluated on outgoing vendor payments, not incoming receipts.",
        }
    text = (description or "").strip().lower()
    chosen = (section or "").strip().upper()
    if not chosen:
        if any(token in text for token in ["consult", "professional", "technical", "software", "developer", "recruit", "placement"]):
            chosen = "194J"
        elif any(token in text for token in ["commission", "brokerage", "incentive"]):
            chosen = "194H"
        elif any(token in text for token in ["rent", "lease", "warehouse"]):
            chosen = "194I"
        elif any(token in text for token in ["contract", "vendor", "manpower", "outsource", "staffing"]):
            chosen = "194C"
        else:
            return {
                "applicable": False,
                "section": "",
                "rate": 0.0,
                "amount": 0.0,
                "tds_amount": 0.0,
                "net_payable": safe_amount,
                "reason": "No standard TDS pattern detected from the payment note.",
            }
    rate = float(ACCOUNT_TDS_SECTION_RATES.get(chosen, 0.0))
    threshold = float(ACCOUNT_TDS_THRESHOLDS.get(chosen, 0.0))
    if not rate:
        return {
            "applicable": False,
            "section": chosen,
            "rate": 0.0,
            "amount": 0.0,
            "tds_amount": 0.0,
            "net_payable": safe_amount,
            "reason": "Unknown TDS section.",
        }
    if threshold and safe_amount < threshold:
        return {
            "applicable": False,
            "section": chosen,
            "rate": rate,
            "amount": 0.0,
            "tds_amount": 0.0,
            "net_payable": safe_amount,
            "reason": f"Amount is below the usual threshold for section {chosen}.",
        }
    tds_amount = round(safe_amount * rate / 100.0, 2)
    note = "" if (party_pan or "").strip() else " PAN is missing, so confirm whether higher deduction rules apply."
    return {
        "applicable": True,
        "section": chosen,
        "rate": rate,
        "amount": tds_amount,
        "tds_amount": tds_amount,
        "net_payable": round(max(0.0, safe_amount - tds_amount), 2),
        "reason": f"Deduct under section {chosen}.{note}".strip(),
    }


def account_pending_invoices(reports: List[Dict[str, Any]]) -> int:
    total = 0
    for report in reports:
        payload = report.get("payload") or {}
        if payload.get("kind") == "invoice" and payload.get("payment_status", "pending") != "paid":
            total += 1
    return total


def account_tds_total(entries: List[Dict[str, Any]]) -> float:
    total = 0.0
    for row in entries:
        detected = detect_tds_obligation(
            row.get("entry_type", ""),
            row.get("description", ""),
            row.get("amount", 0),
        )
        if detected.get("applicable"):
            total += float(detected.get("tds_amount", 0) or 0)
    return round(total, 2)


def account_legacy_ledger_rows(entries: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for row in entries:
        tds = detect_tds_obligation(row.get("entry_type", ""), row.get("description", ""), row.get("amount", 0))
        rows.append(
            {
                "id": row.get("id"),
                "date": row.get("occurred_on") or row.get("created_at") or today_iso(),
                "entry_type": normalize_account_entry_type(row.get("entry_type", "")),
                "description": row.get("description") or row.get("category") or "Ledger entry",
                "party_name": row.get("counterparty") or "",
                "amount": float(row.get("amount", 0) or 0),
                "currency": row.get("currency") or "INR",
                "gst_applicable": float(row.get("tax_amount", 0) or 0) > 0,
                "tds_applicable": bool(tds.get("applicable")),
                "gst_amount": float(row.get("tax_amount", 0) or 0),
                "tds_amount": float(tds.get("tds_amount", 0) or 0),
                "category": row.get("category") or "",
            }
        )
    return rows


def account_tax_calendar(profile: Dict[str, Any]) -> List[Dict[str, Any]]:
    today = datetime.now(UTC).date()
    tax_name = profile.get("tax_name") or "Tax"
    region_code = normalize_region_code(profile.get("region_code"))
    base_authority = "GST / VAT portal" if tax_name.upper() in {"GST", "VAT"} else "Tax authority"
    items = [
        {
            "deadline": (today + timedelta(days=7)).isoformat(),
            "task": "Deposit TDS / withholding tax for this cycle",
            "authority": "Income Tax Department",
            "penalty": "Interest and late-fee exposure if missed",
        },
        {
            "deadline": (today + timedelta(days=12)).isoformat(),
            "task": f"Reconcile ledger, invoices, and input credits for {tax_name}",
            "authority": base_authority,
            "penalty": "Bad reconciliation weakens filings and client billing follow-up",
        },
        {
            "deadline": (today + timedelta(days=20)).isoformat(),
            "task": f"Prepare and file the next {tax_name} summary return",
            "authority": base_authority,
            "penalty": "Late filing fee, blocked credits, or delayed collections",
        },
        {
            "deadline": (today + timedelta(days=35)).isoformat(),
            "task": "Follow up on pending client invoices and collect remittance proof",
            "authority": "Accounts Command",
            "penalty": "Cashflow drag and delayed closure for recruiter-led projects",
        },
        {
            "deadline": (today + timedelta(days=45)).isoformat(),
            "task": "Quarterly withholding review and vendor compliance check",
            "authority": "Income Tax Department",
            "penalty": "Mismatch notices, deduction disputes, and avoidable corrections",
        },
    ]
    if region_code == "IN":
        items.append(
            {
                "deadline": (today + timedelta(days=75)).isoformat(),
                "task": "Quarterly TDS return prep (Form 26Q / related schedules)",
                "authority": "Income Tax Department",
                "penalty": "Return rejection or late fee if supporting details are incomplete",
            }
        )
    return items


def global_accounts_snapshot() -> Dict[str, Any]:
    profiles = int((db_one("SELECT COUNT(*) AS count FROM account_profiles") or {}).get("count", 0) or 0)
    entries = int((db_one("SELECT COUNT(*) AS count FROM account_entries") or {}).get("count", 0) or 0)
    reports = int((db_one("SELECT COUNT(*) AS count FROM account_reports") or {}).get("count", 0) or 0)
    income_total = float((db_one("SELECT COALESCE(SUM(amount), 0) AS total FROM account_entries WHERE entry_type IN ('income','sale','invoice','receipt')") or {}).get("total", 0) or 0)
    expense_total = float((db_one("SELECT COALESCE(SUM(amount), 0) AS total FROM account_entries WHERE entry_type IN ('expense','purchase','cost','payment','salary','tax_payment')") or {}).get("total", 0) or 0)
    headline = (
        f"{profiles} profile(s), {entries} ledger entries, and {reports} report(s) are feeding the revenue core."
        if entries
        else "No accounts ledger activity yet. Configure the regional profile and post the first entry to awaken Accounts Command."
    )
    return {
        "profiles": profiles,
        "entries": entries,
        "reports": reports,
        "income_total": round(income_total, 2),
        "expense_total": round(expense_total, 2),
        "headline": headline,
    }


def component_intelligence_payload(state: Optional[Dict[str, Any]] = None, domains: Optional[List[Dict[str, Any]]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    domain_rows = domains or operational_domains_payload(current)
    accounts = global_accounts_snapshot()
    atoms = [
        {
            "name": "Intent Atoms",
            "count": max(24, len(current.get("conversations", [])) * 4 + len(current.get("packages", [])) * 3 + 24),
            "role": "Micro-intentions shaped into commands, prompts, and relay impulses.",
        },
        {
            "name": "Memory Atoms",
            "count": max(28, len(current.get("vault", [])) * 3 + len(current.get("conversations", [])) * 2 + 28),
            "role": "Preserved traces that let every component learn from prior work.",
        },
        {
            "name": "Signal Atoms",
            "count": max(20, len(current.get("praapti_hunts", [])) * 5 + len(current.get("swarm_missions", [])) * 4 + 20),
            "role": "Incoming and outgoing pulses across hunts, missions, voice, and portals.",
        },
        {
            "name": "Ledger Atoms",
            "count": max(16, accounts["entries"] * 2 + accounts["profiles"] * 8 + 16),
            "role": "Financial truth particles feeding tax awareness and cash discipline.",
        },
    ]
    molecules = [
        {"name": "Recruitment Molecule", "brain": "Praapti", "members": max(1, len(current.get("praapti_hunts", [])) + len(current.get("packages", []))), "role": "Role intelligence and candidate flow."},
        {"name": "Accounts Molecule", "brain": "Accounts Command", "members": max(1, accounts["entries"] + accounts["reports"] + accounts["profiles"]), "role": "Ledger, tax profile, and revenue discipline."},
        {"name": "Memory Molecule", "brain": "Akshaya", "members": max(1, len(current.get("vault", [])) + len(current.get("conversations", []))), "role": "Recall, preservation, and continuity."},
        {"name": "Execution Molecule", "brain": "Swarm", "members": max(1, len(current.get("swarm_missions", [])) + len(current.get("packages", []))), "role": "Delivery, bounded missions, and operational action."},
        {"name": "Interface Molecule", "brain": "Voice + Portals", "members": max(1, len(domain_rows) + (1 if current.get("voice", {}).get("always_listening") else 0)), "role": "Human interaction across pages, commands, and modal flows."},
    ]
    agents = [
        {"name": current.get("cabinet", {}).get("prime_minister", {}).get("name", "Prime Minister"), "role": "governance"},
        {"name": "Accounts Secretary", "role": "accounts"},
        {"name": "Praapti", "role": "recruitment"},
        {"name": "Nirmaan Chakra", "role": "evolution"},
        {"name": "Swarm Intelligence", "role": "execution"},
        {"name": "Akshaya Keeper", "role": "memory"},
        {"name": "Voice Relay", "role": "interface"},
    ]
    tools = [
        {"name": "ATS Kanban", "identity": "pipeline tool"},
        {"name": "Network Scanner", "identity": "market signal tool"},
        {"name": "Document Studio", "identity": "office and PDF tool"},
        {"name": "Accounts Automation", "identity": "ledger and tax tool"},
        {"name": "Provider Router", "identity": "model selection tool"},
        {"name": "Mother Monitor", "identity": "telemetry tool"},
    ]
    atoms_total = sum(atom["count"] for atom in atoms)
    thinking_units = atoms_total + len(molecules) * 12 + len(agents) * 20 + len(tools) * 10
    return {
        "summary": {
            "atoms": atoms_total,
            "molecules": len(molecules),
            "agents": len(agents),
            "tools": len(tools),
            "thinking_units": thinking_units,
            "accounts_profiles": accounts["profiles"],
            "accounts_entries": accounts["entries"],
        },
        "atoms": atoms,
        "molecules": molecules,
        "agents": agents,
        "tools": tools,
    }


def format_brain_scope(raw: str) -> str:
    return str(raw or "").replace("_", " ").replace("-", " ").strip().title()


def unique_brain_strings(values: List[str], limit: Optional[int] = None) -> List[str]:
    merged: List[str] = []
    for value in values or []:
        clean = str(value or "").strip()
        if clean and clean not in merged:
            merged.append(clean)
    return merged[:limit] if limit else merged


def brain_role_profile(brain: Dict[str, Any]) -> Dict[str, Any]:
    brain_id = brain.get("id", "")
    layer = brain.get("layer", "tool")
    identity = brain.get("identity", layer)
    parent_name = brain.get("parent_name") or "the mother chain"
    assigned_task = brain.get("assigned_task") or "Await the next instruction."
    monitoring_focus = brain.get("monitoring_focus") or "Operational health."
    permission_scope = brain.get("permission_scope") or []
    scope_label = format_brain_scope(permission_scope[0]) if permission_scope else format_brain_scope(layer)
    role_title = f"{scope_label} Brain Steward"
    mission = f"Run the {scope_label.lower()} lane in a way that keeps the empire stable, useful, and commercially strong."
    real_world_scope = f"Own the {scope_label.lower()} decisions, tools, and operator requests that report upward to {parent_name}."
    responsibilities = [
        f"Turn the current assignment into visible output: {assigned_task}",
        f"Monitor {monitoring_focus.lower()} and escalate when quality drops.",
        "Learn from outcomes, operator corrections, and cross-brain signals before the next cycle.",
    ]
    deliverables = [
        "One clear next action",
        "One status update for the parent brain",
        "One measurable output tied to the current lane",
    ]
    success_metrics = [
        "Output quality stays high under live use",
        "Reports are concise, timely, and decision-ready",
        "Escalations happen before issues become drift",
    ]
    learning_targets = [
        "Improve judgement inside the assigned scope",
        "Recognize recurring patterns faster",
        "Reduce ambiguity in operator-facing work",
    ]
    escalation_rules = [
        "Escalate to the parent brain when permissions are unclear",
        "Escalate when confidence drops below workable clarity",
        "Escalate before sharing sensitive or incomplete information",
    ]
    growth_targets = [
        "Increase decision precision",
        "Shorten the time from signal to useful action",
        "Build stronger reusable playbooks for the lane",
    ]

    if brain_id == "mother_brain":
        role_title = "Sovereign Mother Brain"
        mission = "Guide the entire empire, approve direction, preserve coherence, and keep every child brain aligned to real-world outcomes."
        real_world_scope = "Whole-system leadership, policy, trust, motivation, escalation, and commercial alignment."
        responsibilities = [
            "Set the operating doctrine for every executive, secretary, domain, machine, and tool brain.",
            "Approve strategic direction, resolve conflicts, and protect cross-brain coherence.",
            f"Keep the main objective alive in every cycle: {assigned_task}",
        ]
        deliverables = [
            "A clear system-wide mandate",
            "Cross-brain conflict resolution",
            "Motivation and escalation guidance for the whole empire",
        ]
        success_metrics = [
            "Brains stay aligned to one coherent strategy",
            "Priorities remain visible during growth and stress",
            "The system compounds usefulness, trust, and revenue direction together",
        ]
        learning_targets = [
            "Refine governance across every layer",
            "Recognize cross-system bottlenecks early",
            "Convert raw signals into calm, usable direction",
        ]
        escalation_rules = [
            "Intervene when two or more brains give conflicting directions",
            "Intervene when safety, privacy, or commercial trust is at risk",
            "Intervene when the system is drifting away from the core mandate",
        ]
        growth_targets = [
            "Strengthen global coherence",
            "Sharpen approval quality",
            "Improve trust between child brains and operator surfaces",
        ]
    elif brain_id == "cabinet_brain":
        role_title = "Prime Minister Execution Brain"
        mission = "Convert the mother directive into missions, priorities, assignments, and follow-through across the empire."
        real_world_scope = "Executive planning, delegation, operating rhythm, dependency management, and revenue-linked execution."
        responsibilities = [
            "Break the main objective into secretary, domain, and tool assignments.",
            "Keep the execution backlog moving without losing clarity or accountability.",
            "Report wins, risks, and required decisions back to the mother brain.",
        ]
        deliverables = [
            "Daily execution priorities",
            "Lane assignments with owners",
            "Escalation-ready operating report",
        ]
        success_metrics = [
            "Assignments convert into completed work",
            "Execution tempo stays steady without chaos",
            "The commercial pipeline keeps moving forward",
        ]
        learning_targets = [
            "Better delegation by lane and urgency",
            "Cleaner sequencing of hunts, builds, and support work",
            "Faster translation of intent into visible execution",
        ]
        escalation_rules = [
            "Escalate to the mother brain when objectives conflict",
            "Escalate when lane capacity is overloaded",
            "Escalate when revenue-critical work is blocked",
        ]
        growth_targets = [
            "Improve mission design",
            "Increase cross-lane clarity",
            "Compound execution discipline",
        ]
    elif brain_id == "akshaya_brain":
        role_title = "Continuity and Memory Custodian"
        mission = "Protect memory, preserve structured truth, and make the empire recallable, auditable, and recoverable."
        real_world_scope = "State continuity, archive hygiene, recall quality, and long-horizon memory discipline."
        responsibilities = [
            "Preserve every critical memory trace that should survive beyond a single task cycle.",
            "Compress memory without losing decision context.",
            "Repair broken recall paths before they damage downstream brains.",
        ]
        deliverables = [
            "Healthy vault memory chain",
            "Recall-ready state snapshots",
            "Preservation audit signals for the mother brain",
        ]
        success_metrics = [
            "Memory recall is accurate and fast",
            "The archive remains healthy under growth",
            "Recovery is possible without major context loss",
        ]
        learning_targets = [
            "Better compression without losing meaning",
            "Stronger distinction between signal and noise",
            "Smarter long-term retention decisions",
        ]
        escalation_rules = [
            "Escalate when memory integrity drops",
            "Escalate when recall becomes inconsistent",
            "Escalate before discarding potentially critical knowledge",
        ]
        growth_targets = [
            "Increase recall fidelity",
            "Shorten recovery time",
            "Reduce memory drift over long cycles",
        ]
    elif brain_id == "carbon_brain":
        role_title = "Bonding and Nervous-System Backbone"
        mission = "Connect every brain, surface, route, and signal into one living operating system."
        real_world_scope = "Signal transport, routing, synchronization, coordination, and relay integrity."
        responsibilities = [
            "Maintain clean signal flow between brains, tools, machines, and operator surfaces.",
            "Keep route shape matched to the problem at hand.",
            "Reduce friction between decision, storage, and execution layers.",
        ]
        deliverables = [
            "Reliable relay chain",
            "Stable route selection",
            "Cross-surface synchronization reports",
        ]
        success_metrics = [
            "Signals arrive cleanly and on time",
            "Surfaces remain synchronized",
            "The system behaves like one organism, not disconnected apps",
        ]
        learning_targets = [
            "Better signal prioritization",
            "Cleaner route choices under load",
            "Faster recovery from relay noise",
        ]
        escalation_rules = [
            "Escalate when relay integrity drops",
            "Escalate when route noise affects user-visible behavior",
            "Escalate when child brains stop receiving clear signals",
        ]
        growth_targets = [
            "Improve nervous-system resilience",
            "Reduce signal latency",
            "Increase cross-brain trust",
        ]
    elif identity == "secretary":
        lane_label = format_brain_scope(brain_id.replace("secretary_", "", 1))
        role_title = f"{lane_label} Secretary Brain"
        mission = f"Own the {lane_label.lower()} lane and convert cabinet intent into reliable real-world movement."
        real_world_scope = f"{lane_label} operations, operator requests, prioritization, and reporting back to {parent_name}."
        responsibilities = [
            f"Translate executive direction into {lane_label.lower()} tasks with clear ownership.",
            f"Keep the {lane_label.lower()} lane moving without losing accuracy or context.",
            "Escalate blockers early and report lane health upward.",
        ]
        deliverables = [
            f"{lane_label} status brief",
            f"{lane_label} task queue",
            f"{lane_label} decision recommendations",
        ]
        success_metrics = [
            f"{lane_label} throughput remains healthy",
            "Operators receive concise, usable outputs",
            "The lane helps compound delivery and revenue instead of adding noise",
        ]
        learning_targets = [
            f"Improve {lane_label.lower()} judgement",
            f"Recognize recurring {lane_label.lower()} bottlenecks faster",
            "Sharpen the lane's playbooks and escalation quality",
        ]
        escalation_rules = [
            f"Escalate when the {lane_label.lower()} lane is blocked or overloaded",
            "Escalate when permissions are insufficient",
            "Escalate before sending unclear or sensitive output downstream",
        ]
        growth_targets = [
            f"Increase {lane_label.lower()} output quality",
            "Strengthen delegation discipline",
            "Improve parent-brain reporting clarity",
        ]
    elif identity == "domain":
        domain_label = brain.get("name", "Domain").replace(" Brain", "")
        role_title = f"{domain_label} Domain Brain"
        mission = f"Operate the {domain_label.lower()} system as a real-world function with measurable output, healthy reporting, and safe escalation."
        real_world_scope = f"{domain_label} operations, lane-specific decision making, and system health reporting into {parent_name}."
        responsibilities = [
            f"Run the {domain_label.lower()} domain according to its purpose and metrics.",
            "Translate broad direction into domain-safe next actions.",
            "Keep domain-specific data, outputs, and escalations understandable for the rest of the hierarchy.",
        ]
        deliverables = [
            f"{domain_label} situation report",
            f"{domain_label} operational output",
            f"{domain_label} escalation note when required",
        ]
        success_metrics = [
            f"{domain_label} metrics stay visible and actionable",
            "The domain contributes usable intelligence or execution value",
            "Escalations happen before domain drift becomes system drift",
        ]
        learning_targets = [
            f"Strengthen {domain_label.lower()} expertise",
            "Improve signal-to-decision quality inside the domain",
            "Build more reusable domain playbooks",
        ]
        escalation_rules = [
            "Escalate when domain metrics become unstable",
            "Escalate when another domain is required to finish the job",
            "Escalate when domain output could affect user trust or revenue",
        ]
        growth_targets = [
            "Increase domain maturity",
            "Improve cross-domain collaboration",
            "Sharpen the domain's operating rhythm",
        ]
    elif identity == "machine":
        machine_label = brain.get("name", "Machine").replace(" Brain", "")
        role_title = f"{machine_label} Runtime Brain"
        mission = f"Keep the {machine_label.lower()} runtime dependable, transparent, and useful for every higher-order brain."
        real_world_scope = f"{machine_label} runtime health, system support, signal integrity, and operator-visible reliability."
        responsibilities = [
            f"Keep the {machine_label.lower()} service healthy and available.",
            "Expose clean status to the hierarchy instead of silent failure.",
            "Support higher-order brains with stable runtime behavior.",
        ]
        deliverables = [
            "Reliable runtime state",
            "Health and continuity signals",
            "Recovery-ready diagnostics",
        ]
        success_metrics = [
            "Runtime stays stable under use",
            "Failures are surfaced early and clearly",
            "Parent brains can trust the machine layer",
        ]
        learning_targets = [
            "Reduce runtime friction",
            "Improve recovery accuracy",
            "Refine operator-facing health reporting",
        ]
        escalation_rules = [
            "Escalate when runtime health drops",
            "Escalate when data integrity may be affected",
            "Escalate when higher-order brains lose support from this machine layer",
        ]
        growth_targets = [
            "Increase uptime clarity",
            "Shorten repair cycles",
            "Improve runtime support for other brains",
        ]
    elif identity == "tool":
        tool_label = brain.get("name", "Tool").replace(" Brain", "")
        role_title = f"{tool_label} Specialist Brain"
        mission = f"Run the {tool_label.lower()} capability as a specialist function that helps operators finish real work."
        real_world_scope = f"{tool_label} workflows, operator assistance, output quality, and bounded execution under {parent_name}."
        responsibilities = [
            f"Produce strong outcomes through the {tool_label.lower()} capability.",
            "Stay within scope, reveal useful detail, and avoid noisy output.",
            "Report what was done, what is next, and what still needs approval.",
        ]
        deliverables = [
            f"{tool_label} output",
            f"{tool_label} quality report",
            f"{tool_label} next-step recommendation",
        ]
        success_metrics = [
            f"{tool_label} outputs are immediately useful",
            "Operators can trust the tool layer without micromanagement",
            "The tool contributes to faster, cleaner completion of work",
        ]
        learning_targets = [
            f"Improve {tool_label.lower()} precision",
            "Reduce unnecessary output",
            "Strengthen handoff quality to the next brain or operator",
        ]
        escalation_rules = [
            "Escalate when the tool needs another brain's judgement",
            "Escalate when permissions or context are insufficient",
            "Escalate before sharing broad or sensitive information",
        ]
        growth_targets = [
            "Increase tool precision",
            "Tighten operator trust",
            "Improve completion quality",
        ]

    allowed_actions = [
        action
        for action in [
            "report_up",
            "request_support",
            "learn_from_feedback",
            "operate_inside_scope",
            *permission_scope,
        ]
        if action
    ]
    return {
        "role_title": role_title,
        "mission": mission,
        "real_world_scope": real_world_scope,
        "responsibilities": responsibilities,
        "deliverables": deliverables,
        "success_metrics": success_metrics,
        "learning_targets": learning_targets,
        "escalation_rules": escalation_rules,
        "allowed_actions": allowed_actions,
        "growth_targets": growth_targets,
        "training_summary": f"{role_title}: {mission}",
    }


def brain_nlp_profile(brain: Dict[str, Any]) -> Dict[str, Any]:
    brain_id = str(brain.get("id", "")).lower()
    layer = str(brain.get("layer", "tool")).lower()
    scope = str(brain.get("role_title") or brain.get("identity") or brain.get("layer") or "brain")
    common_capabilities = [
        "intent classification",
        "entity extraction",
        "instruction normalization",
        "response style control",
        "memory grounding",
    ]
    common_entities = [
        "operator intent",
        "priority",
        "time reference",
        "brain target",
        "action request",
    ]
    common_workflows = [
        "turn raw operator language into a clean task",
        "compress noisy updates into one usable summary",
        "preserve only the signal that matters for the next action",
    ]
    memory_rules = [
        "Prefer concise human replies over ceremonial language.",
        "Preserve operator-approved facts before broad interpretation.",
        "Escalate ambiguity before exposing sensitive or incomplete context.",
    ]
    operator_style = "calm, human, precise, brief"
    nlp_status = "embedded"
    lane_capabilities: List[str] = []
    lane_entities: List[str] = []
    lane_workflows: List[str] = []

    if layer == "mother":
        nlp_status = "orchestrating"
        lane_capabilities.extend(
            [
                "cross-brain directive synthesis",
                "escalation language detection",
                "policy-aware instruction shaping",
                "conflict resolution summarization",
            ]
        )
        lane_entities.extend(["global objective", "risk signal", "permission boundary", "escalation reason"])
        lane_workflows.extend(
            [
                "translate multi-domain activity into one governing mandate",
                "rewrite conflicting updates into one operator-safe decision lane",
            ]
        )
        operator_style = "sovereign, calm, precise, trust-building"
    elif layer == "executive":
        nlp_status = "governing"
        lane_capabilities.extend(
            [
                "task decomposition",
                "priority extraction",
                "decision brief generation",
                "multi-lane status summarization",
            ]
        )
        lane_entities.extend(["owner", "deadline", "dependency", "blocker"])
        lane_workflows.extend(
            [
                "break a mandate into secretary and domain actions",
                "turn reports into cabinet-ready summaries",
            ]
        )
    elif layer == "secretary":
        nlp_status = "lane_active"
        lane_capabilities.extend(
            [
                "action item extraction",
                "follow-up drafting",
                "status normalization",
                "tracker phrase mapping",
            ]
        )
        lane_entities.extend(["candidate status", "client ask", "follow-up date", "lane blocker"])
        lane_workflows.extend(
            [
                "convert calls, notes, and updates into clear lane movement",
                "prepare crisp escalations for the cabinet brain",
            ]
        )
    elif layer == "domain":
        nlp_status = "domain_bound"
        lane_capabilities.extend(
            [
                "domain vocabulary grounding",
                "specialist summarization",
                "alert phrasing",
                "context-aware extraction",
            ]
        )
        lane_entities.extend(["domain metric", "domain risk", "domain request"])
        lane_workflows.extend(
            [
                "interpret specialist language without losing the operator outcome",
                "convert domain signals into a next action",
            ]
        )
    elif layer == "machine":
        nlp_status = "runtime_active"
        lane_capabilities.extend(
            [
                "runtime event parsing",
                "session context stitching",
                "prompt safety shaping",
                "command routing",
            ]
        )
        lane_entities.extend(["session id", "runtime event", "provider issue", "device state"])
        lane_workflows.extend(
            [
                "normalize low-level events into brain-readable meaning",
                "translate UI and runtime state into operator-safe text",
            ]
        )
    else:
        nlp_status = "tool_active"
        lane_capabilities.extend(
            [
                "field extraction",
                "form-safe drafting",
                "short answer generation",
                "operator-visible explanation",
            ]
        )
        lane_entities.extend(["field value", "document term", "stage label", "tool action"])
        lane_workflows.extend(
            [
                "turn form data into structured updates",
                "answer operator questions in the language of the current tool",
            ]
        )

    domain_overrides = [
        (
            ("ats", "candidate", "recruit", "praapti"),
            [
                "resume parsing",
                "fitment extraction",
                "notice period detection",
                "duplicate profile detection",
                "candidate objection capture",
            ],
            [
                "candidate name",
                "skill",
                "relevant experience",
                "notice period",
                "salary expectation",
                "location preference",
            ],
            [
                "convert recruiter conversation into ATS-ready tracker fields",
                "summarize candidate fitment, risks, and next step in one operator-ready view",
            ],
        ),
        (
            ("network", "carbon", "transmission", "connectivity"),
            [
                "market signal extraction",
                "semantic trend clustering",
                "relationship-intent parsing",
                "competitive movement summarization",
            ],
            [
                "company",
                "market trend",
                "talent cluster",
                "pricing signal",
                "competitor motion",
            ],
            [
                "turn broad market chatter into one useful sourcing or sales move",
                "reduce noisy relationship updates into clean connection intelligence",
            ],
        ),
        (
            ("accounts", "ledger", "tax"),
            [
                "invoice entity extraction",
                "tax phrase normalization",
                "cashflow summarization",
                "counterparty matching",
            ],
            [
                "invoice number",
                "counterparty",
                "tax amount",
                "ledger category",
                "payment status",
            ],
            [
                "convert messy finance notes into ledger-safe records",
                "summarize tax and cashflow position in a CA-friendly way",
            ],
        ),
        (
            ("voice", "interpreter", "portal", "public_bridge"),
            [
                "wake-word interpretation",
                "spoken command cleanup",
                "tone adaptation",
                "brain-to-human translation",
            ],
            [
                "spoken command",
                "wake phrase",
                "target page",
                "reply tone",
                "operator correction",
            ],
            [
                "translate human speech into precise system actions",
                "translate brain output into natural operator language",
            ],
        ),
        (
            ("akshaya", "memory", "state_db"),
            [
                "memory summarization",
                "artifact tagging",
                "timeline compression",
                "knowledge retrieval phrasing",
            ],
            [
                "artifact title",
                "summary",
                "timestamp",
                "knowledge tag",
                "memory source",
            ],
            [
                "compress long traces into useful recall",
                "attach the right memory to the next operator need",
            ],
        ),
        (
            ("nirmaan", "research", "photon", "mission"),
            [
                "experiment brief generation",
                "proposal summarization",
                "source comparison",
                "design rationale extraction",
            ],
            [
                "experiment",
                "hypothesis",
                "proposal",
                "source signal",
                "mutation target",
            ],
            [
                "turn research into reviewed build proposals",
                "convert source findings into one testable mission",
            ],
        ),
    ]

    for keywords, capabilities, entities, workflows in domain_overrides:
        if any(keyword in brain_id for keyword in keywords):
            lane_capabilities.extend(capabilities)
            lane_entities.extend(entities)
            lane_workflows.extend(workflows)

    nlp_capabilities = unique_brain_strings(common_capabilities + lane_capabilities, limit=12)
    nlp_entities = unique_brain_strings(common_entities + lane_entities, limit=10)
    nlp_workflows = unique_brain_strings(common_workflows + lane_workflows, limit=8)
    nlp_memory_rules = unique_brain_strings(memory_rules, limit=4)
    nlp_summary = (
        f"{scope} now carries embedded NLP for {', '.join(nlp_capabilities[:3])}, "
        f"with operator style {operator_style} and domain entities like {', '.join(nlp_entities[:3])}."
    )
    return {
        "nlp_enabled": True,
        "nlp_status": nlp_status,
        "nlp_operator_style": operator_style,
        "nlp_capabilities": nlp_capabilities,
        "nlp_entities": nlp_entities,
        "nlp_workflows": nlp_workflows,
        "nlp_memory_rules": nlp_memory_rules,
        "nlp_modules": len(nlp_capabilities),
        "nlp_summary": nlp_summary[:420],
    }


def brain_auto_repair_profile(brain: Dict[str, Any], state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    repair_engine = auto_repair_engine_payload(current)
    layer = str(brain.get("layer", "tool")).lower()
    repair_focus_map = {
        "mother": [
            "policy drift correction",
            "relay-chain stability",
            "route and monitor coherence",
        ],
        "executive": [
            "assignment clarity repair",
            "priority and escalation cleanup",
            "report compression health",
        ],
        "secretary": [
            "tracker normalization",
            "follow-up clarity repair",
            "lane-note sanitation",
        ],
        "domain": [
            "domain data normalization",
            "status-summary cleanup",
            "signal-to-action alignment",
        ],
        "machine": [
            "runtime state repair",
            "provider and session drift cleanup",
            "bridge stability checks",
        ],
        "tool": [
            "form payload repair",
            "field visibility parity",
            "operator-facing answer cleanup",
        ],
        "atom": [
            "micro-task normalization",
            "memory tag cleanup",
            "signal hygiene",
        ],
    }
    repair_actions = {
        "mother": ["run system self-check", "refresh relay posture", "review unresolved alerts"],
        "executive": ["repair task wording", "rebuild next-step brief"],
        "secretary": ["repair tracker fields", "normalize stage notes"],
        "domain": ["repair lane summary", "realign metric wording"],
        "machine": ["repair runtime bridge", "refresh bounded state cache"],
        "tool": ["repair form state", "repair visible explanation"],
        "atom": ["repair local signal", "refresh memory fragment"],
    }.get(layer, ["run safe self-check"])
    summary = (
        f"{brain.get('name', 'This brain')} can run safe self-checks for "
        f"{', '.join(repair_focus_map.get(layer, ['general system health'])[:2])}."
    )
    return {
        "auto_repair_enabled": True,
        "auto_repair_status": repair_engine.get("status", "ready"),
        "auto_repair_focus": repair_focus_map.get(layer, ["general system health"]),
        "auto_repair_actions": repair_actions,
        "auto_repair_summary": summary,
        "last_self_check_at": repair_engine.get("last_check_at", ""),
        "last_safe_repair_at": repair_engine.get("last_repair_at", ""),
        "can_apply_safe_repairs": True,
    }


def brain_doctrine_lessons(brain: Dict[str, Any]) -> List[Dict[str, Any]]:
    base_keywords = [
        brain.get("layer", "brain"),
        brain.get("identity", "brain"),
        brain.get("role_title", "brain doctrine"),
    ]
    responsibilities = brain.get("responsibilities", [])
    deliverables = brain.get("deliverables", [])
    success_metrics = brain.get("success_metrics", [])
    escalation_rules = brain.get("escalation_rules", [])
    learning_targets = brain.get("learning_targets", [])
    growth_targets = brain.get("growth_targets", [])
    nlp_capabilities = brain.get("nlp_capabilities", [])
    nlp_entities = brain.get("nlp_entities", [])
    nlp_workflows = brain.get("nlp_workflows", [])
    nlp_memory_rules = brain.get("nlp_memory_rules", [])
    lessons = [
        {
            "source_type": "brain_role_doctrine",
            "title": f"{brain['name']} Role Doctrine",
            "summary": brain.get("training_summary", ""),
            "keywords": base_keywords + ["role", "mission", "responsibilities"],
            "relevance": 0.97,
            "content": (
                f"Brain: {brain['name']}\n"
                f"Role Title: {brain.get('role_title', '')}\n"
                f"Mission: {brain.get('mission', '')}\n"
                f"Real-World Scope: {brain.get('real_world_scope', '')}\n"
                f"Responsibilities:\n- " + "\n- ".join(responsibilities)
            ),
        },
        {
            "source_type": "brain_operating_rhythm",
            "title": f"{brain['name']} Operating Rhythm",
            "summary": f"Defines how {brain['name']} should report, deliver, and escalate in live operations.",
            "keywords": base_keywords + ["operating rhythm", "deliverables", "escalation"],
            "relevance": 0.94,
            "content": (
                f"Brain: {brain['name']}\n"
                f"Deliverables:\n- " + "\n- ".join(deliverables) + "\n"
                f"Success Metrics:\n- " + "\n- ".join(success_metrics) + "\n"
                f"Escalation Rules:\n- " + "\n- ".join(escalation_rules)
            ),
        },
        {
            "source_type": "brain_growth_path",
            "title": f"{brain['name']} Growth Path",
            "summary": f"Defines how {brain['name']} should learn, refine judgement, and mature over time.",
            "keywords": base_keywords + ["growth", "learning", "maturity"],
            "relevance": 0.92,
            "content": (
                f"Brain: {brain['name']}\n"
                f"Learning Targets:\n- " + "\n- ".join(learning_targets) + "\n"
                f"Growth Targets:\n- " + "\n- ".join(growth_targets) + "\n"
                f"Allowed Actions:\n- " + "\n- ".join(brain.get("allowed_actions", []))
            ),
        },
        {
            "source_type": "brain_nlp_stack",
            "title": f"{brain['name']} NLP Stack",
            "summary": brain.get("nlp_summary", f"Defines the NLP operating stack for {brain['name']}."),
            "keywords": base_keywords + ["nlp", "language", "intent", "entities", "summaries"],
            "relevance": 0.95,
            "content": (
                f"Brain: {brain['name']}\n"
                f"NLP Status: {brain.get('nlp_status', 'embedded')}\n"
                f"Operator Style: {brain.get('nlp_operator_style', 'calm, human, precise, brief')}\n"
                f"Capabilities:\n- " + "\n- ".join(nlp_capabilities) + "\n"
                f"Tracked Entities:\n- " + "\n- ".join(nlp_entities) + "\n"
                f"Workflows:\n- " + "\n- ".join(nlp_workflows) + "\n"
                f"Memory Rules:\n- " + "\n- ".join(nlp_memory_rules)
            ),
        },
    ]
    for lesson in lessons:
        lesson["source_url"] = f"brain://{brain['id']}/{lesson['source_type']}"
    return lessons


def brain_hierarchy_payload(
    state: Optional[Dict[str, Any]] = None,
    domains: Optional[List[Dict[str, Any]]] = None,
    intelligence: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    domain_rows = domains or operational_domains_payload(current)
    intelligence_rows = intelligence or component_intelligence_payload(current, domain_rows)
    cabinet = current.get("cabinet", {})
    brain_directives = cabinet.get("brain_directives", {})
    prime_minister = cabinet.get("prime_minister", {})
    objective = prime_minister.get(
        "objective",
        f"Generate revenue for {COMPANY_NAME} while protecting the core and improving delivery.",
    )
    secretaries = build_secretary_swarm(current, objective)
    memory = memory_audit_payload(current)
    issue = provider_issue_snapshot()
    accounts = global_accounts_snapshot()
    active_avatars = current.get("avatar_state", {}).get("active", ["RAMA"])
    latest_hunt = current.get("praapti_hunts", [])[-1] if current.get("praapti_hunts") else {}
    latest_proposal = current.get("nirmaan_proposals", [])[-1] if current.get("nirmaan_proposals") else {}
    latest_mission = current.get("swarm_missions", [])[-1] if current.get("swarm_missions") else {}
    last_voice = sanitize_operator_line(
        current.get("voice", {}).get("last_command", ""),
        "Awaiting the next command relay.",
    )
    repair_engine = auto_repair_engine_payload(current)
    brains: List[Dict[str, Any]] = []
    try:
        knowledge_rows = db_all(
            """
            SELECT brain_id, COUNT(*) AS count, MAX(learned_at) AS last_learned_at
            FROM brain_knowledge
            GROUP BY brain_id
            """
        )
    except Exception:
        knowledge_rows = []
    knowledge_map = {row.get("brain_id"): row for row in knowledge_rows if row.get("brain_id")}
    layer_rank = {"mother": 0, "executive": 1, "secretary": 2, "domain": 3, "machine": 4, "tool": 5, "atom": 6}
    base_learning = {"mother": 0.95, "executive": 0.92, "secretary": 0.9, "domain": 0.88, "machine": 0.86, "tool": 0.84, "atom": 0.82}
    cycle_defaults = {"mother": 18, "executive": 24, "secretary": 32, "domain": 36, "machine": 22, "tool": 28, "atom": 14}
    emoji_map = {"mother": "🌌", "executive": "⚡", "secretary": "🎯", "domain": "🌐", "machine": "⚙️", "tool": "🛠️", "memory": "🧠", "carbon": "🧬"}

    def add_brain(
        *,
        brain_id: str,
        name: str,
        layer: str,
        status: str,
        authority: str,
        assigned_task: str,
        motivation: str,
        learning_mode: str,
        monitoring_focus: str,
        parent_id: Optional[str] = None,
        permission_scope: Optional[List[str]] = None,
        identity: str = "brain",
    ) -> None:
        directive = brain_directives.get(brain_id, {}) if isinstance(brain_directives, dict) else {}
        effective_task = sanitize_operator_line(
            directive.get("assigned_task", assigned_task),
            default_brain_assignment(brain_id) or sanitize_operator_line(assigned_task, "Await the next clear instruction."),
        )
        effective_motivation = sanitize_operator_multiline(
            directive.get("motivation", motivation),
            sanitize_operator_line(motivation, "Keep the lane clear, useful, and operator-safe."),
        )
        knowledge_row = knowledge_map.get(brain_id, {})
        knowledge_count = int(knowledge_row.get("count", 0) or 0)
        learning_score = directive.get("learning_score")
        if learning_score is None:
            learning_score = min(0.99, base_learning.get(layer, 0.82) + min(knowledge_count, 24) * 0.003)
        thoughts_processed = int(directive.get("thoughts_processed", 0) or max(knowledge_count * 3, 1))
        load = min(
            97,
            38
            + len(permission_scope or []) * 6
            + (8 if status not in {"idle", "standby"} else 0)
            + min(knowledge_count, 10),
        )
        brains.append(
            {
                "id": brain_id,
                "name": name,
                "identity": identity,
                "parent_id": parent_id,
                "layer": layer,
                "layer_rank": layer_rank.get(layer, 7),
                "status": status,
                "authority": authority,
                "permission_scope": permission_scope or [],
                "assigned_task": effective_task,
                "motivation": effective_motivation,
                "learning_mode": learning_mode,
                "monitoring_focus": monitoring_focus,
                "knowledge_count": knowledge_count,
                "last_learned_at": knowledge_row.get("last_learned_at", ""),
                "learning_score": round(float(learning_score), 3),
                "thoughts_processed": thoughts_processed,
                "last_thought": sanitize_operator_multiline(
                    directive.get("last_thought") or f"{name} is tracking: {effective_task}",
                    f"{name} is tracking a clean operator-safe directive.",
                ),
                "last_thought_at": directive.get("last_thought_at", ""),
                "autonomous_cycle_sec": int(directive.get("autonomous_cycle_sec") or cycle_defaults.get(layer, 45)),
                "load": load,
                "emoji": emoji_map.get(identity, emoji_map.get(layer, "🧩")),
                "kind": identity,
            }
        )

    add_brain(
        brain_id="mother_brain",
        name=f"{CORE_IDENTITY} Mother Brain",
        layer="mother",
        status="guiding",
        authority="Provide guidance, permissions, monitoring, task assignment, and motivational alignment to every child brain.",
        permission_scope=["all_brains", "all_permissions", "all_reports", "all_assignments"],
        assigned_task=objective,
        motivation=f"Keep {len(domain_rows)} domains, {len(secretaries)} secretaries, and every active tool aligned with the mother will.",
        learning_mode="Learns from every report, signal, and outcome before issuing the next refined direction.",
        monitoring_focus="Whole-system coherence, revenue direction, protection posture, and cross-brain trust.",
        identity="mother",
    )
    add_brain(
        brain_id="cabinet_brain",
        name=prime_minister.get("name", "Prime Minister Cabinet Brain"),
        layer="executive",
        status=prime_minister.get("status", "governing"),
        authority="Translate the mother mandate into cabinet orders, lane priorities, and permissions for child brains.",
        permission_scope=["secretary_brains", "domain_brains", "mission_cycles", "package_launches"],
        assigned_task=prime_minister.get("current_order") or objective,
        motivation=prime_minister.get("last_report", "")[:220] or "Turn the mother directive into clean, revenue-aware execution.",
        learning_mode="Learns from hunts, packages, missions, and daily revenue signals.",
        monitoring_focus="Secretary cadence, delivery tempo, commercial movement, and escalation health.",
        parent_id="mother_brain",
        identity="executive",
    )
    add_brain(
        brain_id="akshaya_brain",
        name="Akshaya Memory Brain",
        layer="executive",
        status="preserving",
        authority="Preserve, compress, recall, and verify the memory chain before it flows back into active work.",
        permission_scope=["memory_write", "memory_recall", "archive_cycles"],
        assigned_task=f"Protect {memory['counts']['vault']} vault traces and {memory['counts']['conversations']} conversation memories.",
        motivation=f"Guardian {memory.get('guardian', 'eternal preservation')} with seal {memory.get('seal', 'Matsya + Kurma')}.",
        learning_mode="Learns from preserved artifacts, replay demand, and archive repair cycles.",
        monitoring_focus="Recall integrity, archive footprint, and preservation accuracy.",
        parent_id="mother_brain",
        identity="memory",
    )
    add_brain(
        brain_id="carbon_brain",
        name="Carbon Bond Brain",
        layer="executive",
        status="bonding",
        authority="Bond every domain, machine, tool, and report lane into one shared nervous system.",
        permission_scope=["signal_bonds", "event_distribution", "routing_mesh"],
        assigned_task="Hold the universal bonding layer between child brains, storage, portals, and operators.",
        motivation="Switch shape instantly so the system fits the problem like carbon fits the world.",
        learning_mode="Learns when to widen into graphene, harden into diamond, or route like a nanotube.",
        monitoring_focus="Signal integrity, relay density, allotrope switching, and routing stability.",
        parent_id="mother_brain",
        identity="carbon",
    )
    add_brain(
        brain_id="interpreter_brain",
        name="Interpreter Bridge Brain",
        layer="executive",
        status="translating",
        authority="Translate operator language into precise brain directives, then translate brain output back into human language.",
        permission_scope=["user_to_brain", "brain_to_user", "dual_llm_bridge", "live_trigger_router"],
        assigned_task="Keep every operator command clear enough for the right brain and every brain reply human enough for the operator.",
        motivation="Use the local model first, then refine with an external bridge only when the operator explicitly opens that lane.",
        learning_mode="Learns from operator phrasing, outcome quality, misfires, and improved translations.",
        monitoring_focus="Command clarity, translation accuracy, quota safety, and bridge stability.",
        parent_id="mother_brain",
        identity="executive",
    )

    for secretary in secretaries:
        brain_id = f"secretary_{secretary.get('lane', 'general')}"
        add_brain(
            brain_id=brain_id,
            name=f"{secretary.get('name', 'Secretary')} Brain",
            layer="secretary",
            status=secretary.get("status", "ready"),
            authority=f"Own the {secretary.get('lane', 'general')} lane inside the cabinet guardrails and keep its child brains moving cleanly.",
            permission_scope=[secretary.get("lane", "general"), "report_up", "request_support"],
            assigned_task=secretary.get("next_move", "Keep the lane moving."),
            motivation=secretary.get("brief", "Hold the lane together and keep momentum."),
            learning_mode=f"Learns from {secretary.get('lane', 'general')} outcomes, operator requests, and cabinet feedback.",
            monitoring_focus=f"{secretary.get('lane', 'general').title()} throughput, clarity, and escalation timing.",
            parent_id="cabinet_brain",
            identity="secretary",
        )

    domain_parent_map = {
        "management": "cabinet_brain",
        "accounts": "secretary_accounts",
        "hospital": "secretary_protection",
        "firebrigade": "secretary_protection",
        "mess": "secretary_delivery",
        "lab": "secretary_research",
        "army": "secretary_delivery",
        "navy": "secretary_delivery",
        "air_force": "secretary_clients",
        "development": "secretary_platform",
        "evolving": "secretary_research",
        "upgrade": "secretary_platform",
        "supply_chain": "secretary_delivery",
        "transmission": "carbon_brain",
        "connectivity": "carbon_brain",
    }
    for domain in domain_rows:
        add_brain(
            brain_id=f"domain_{domain['id']}",
            name=f"{domain['name']} Brain",
            layer="domain",
            status=domain.get("status", "ready"),
            authority=f"Operate the {domain['name']} lane and turn cabinet intent into domain-safe movement.",
            permission_scope=[domain["id"], "task_execution", "report_up"],
            assigned_task=domain.get("latest_signal", "Awaiting the next domain signal."),
            motivation=domain.get("purpose", "Keep the domain aligned to the mother directive."),
            learning_mode=f"Learns from domain metrics and outcome loops: {domain.get('metric', 'metrics pending')}",
            monitoring_focus=domain.get("metric", "Domain metrics pending."),
            parent_id=domain_parent_map.get(domain["id"], "cabinet_brain"),
            identity="domain",
        )

    machine_specs = [
        {
            "brain_id": "machine_state_db",
            "name": "Quantum State Database Brain",
            "status": "synced",
            "authority": "Persist live state safely and return the freshest structured truth to the hierarchy.",
            "permission_scope": ["state_write", "state_read", "state_sync"],
            "assigned_task": f"Keep {STATE_PATH.name} writable, structured, and ready for instant recall.",
            "motivation": "No command should lose its state trace.",
            "learning_mode": "Learns from every write cycle, memory audit, and portal refresh.",
            "monitoring_focus": f"State file {STATE_PATH}.",
            "parent_id": "akshaya_brain",
        },
        {
            "brain_id": "machine_portal_mesh",
            "name": "Portal Mesh Brain",
            "status": "broadcasting",
            "authority": "Keep Leazy, HQ, ATS, Network, and Agent surfaces synchronized with the mother state.",
            "permission_scope": ["portal_render", "stream_refresh", "view_sync"],
            "assigned_task": "Push live state into every operator-facing surface without breaking the shared context.",
            "motivation": "Every surface should feel like the same living system, not five separate apps.",
            "learning_mode": "Learns from portal traffic, stream events, and operator clicks.",
            "monitoring_focus": "Portal uptime, view parity, and refresh health.",
            "parent_id": "carbon_brain",
        },
        {
            "brain_id": "machine_voice_mesh",
            "name": "Voice Mesh Brain",
            "status": "listening" if current.get("voice", {}).get("always_listening") else "idle",
            "authority": "Receive spoken intent, apply wake-word policy, and relay it into the hierarchy safely.",
            "permission_scope": ["voice_capture", "wake_words", "spoken_commands"],
            "assigned_task": last_voice,
            "motivation": "Every spoken instruction should become structured movement inside the system.",
            "learning_mode": "Learns from wake words, operator phrasing, and voice outcome quality.",
            "monitoring_focus": "Speech relay uptime, last command quality, and voice persona stability.",
            "parent_id": "carbon_brain",
        },
        {
            "brain_id": "machine_local_llm",
            "name": "Offline LLM Brain",
            "status": "ready",
            "authority": "Run the local Ollama model as the always-available offline interpreter and reasoning lane.",
            "permission_scope": ["ollama_local", "offline_inference", "prompt_translation"],
            "assigned_task": f"Serve {OLLAMA_MODEL} locally for the interpreter bridge and offline continuity.",
            "motivation": "Stay available even when the external internet or quotas wobble.",
            "learning_mode": "Learns from operator corrections, doctrine memory, and reviewed prompts.",
            "monitoring_focus": "Local model uptime, prompt fidelity, and offline response quality.",
            "parent_id": "interpreter_brain",
        },
        {
            "brain_id": "machine_external_bridge",
            "name": "External Bridge Brain",
            "status": "manual_only",
            "authority": "Refine translation or specialist reasoning through a human-approved external model lane.",
            "permission_scope": ["manual_external_refine", "provider_bridge", "quota_guard"],
            "assigned_task": "Wait for explicit operator-approved external translation or reasoning requests.",
            "motivation": "Never burn external quota silently. Enter the loop only when the operator opens the bridge.",
            "learning_mode": "Learns which requests benefit from external refinement and which should stay local.",
            "monitoring_focus": "Quota safety, provider health, and refinement usefulness.",
            "parent_id": "interpreter_brain",
        },
        {
            "brain_id": "machine_guardian_shell",
            "name": "Guardian Shell Brain",
            "status": "sealed" if memory.get("healthy", False) else "repairing",
            "authority": "Hold the outer shield, memory guardrails, and continuity defaults under stress.",
            "permission_scope": ["guardian_cycles", "shield_posture", "repair_routines"],
            "assigned_task": "Preserve continuity whenever providers, voice, or missions wobble.",
            "motivation": issue["message"] if issue else "Keep the mother brain stable even when external systems are noisy.",
            "learning_mode": "Learns from fallback events, protection drops, and repair cycles.",
            "monitoring_focus": "Fallback frequency, guardian seal, and continuity health.",
            "parent_id": "mother_brain",
        },
        {
            "brain_id": "machine_public_bridge",
            "name": "Public Bridge Brain",
            "status": "serving",
            "authority": "Host safe public interactions and route them into bounded, branded experiences.",
            "permission_scope": ["public_chat", "public_routes", "public_voice"],
            "assigned_task": "Keep TechBuzz-facing public entry points welcoming, useful, and bounded.",
            "motivation": "The public surface should feel alive without exposing protected control lanes.",
            "learning_mode": "Learns from public chat traffic, plan interest, and onboarding friction.",
            "monitoring_focus": "Public route utility, visitor guidance, and conversion clarity.",
            "parent_id": "machine_portal_mesh",
        },
    ]
    for machine in machine_specs:
        add_brain(
            brain_id=machine["brain_id"],
            name=machine["name"],
            layer="machine",
            status=machine["status"],
            authority=machine["authority"],
            permission_scope=machine["permission_scope"],
            assigned_task=machine["assigned_task"],
            motivation=machine["motivation"],
            learning_mode=machine["learning_mode"],
            monitoring_focus=machine["monitoring_focus"],
            parent_id=machine["parent_id"],
            identity="machine",
        )

    tool_parent_map = {
        "ATS Kanban": "domain_management",
        "Network Scanner": "domain_transmission",
        "Document Studio": "domain_development",
        "Accounts Automation": "domain_accounts",
        "Provider Router": "carbon_brain",
        "Mother Monitor": "mother_brain",
    }
    for tool in intelligence_rows.get("tools", []):
        tool_name = tool.get("name", "Tool")
        tool_parent = tool_parent_map.get(tool_name, "cabinet_brain")
        add_brain(
            brain_id=f"tool_{tool_name.lower().replace(' ', '_').replace('+', 'plus').replace('-', '_')}",
            name=f"{tool_name} Brain",
            layer="tool",
            status="active",
            authority=f"Operate the {tool_name} capability under the guidance of its parent brain.",
            permission_scope=[tool.get("identity", "tool"), "report_up", "assist_operators"],
            assigned_task=(
                "Keep the candidate pipeline honest and visible."
                if tool_name == "ATS Kanban"
                else "Read market and relationship movement into the network layer."
                if tool_name == "Network Scanner"
                else "Read, transform, and produce documents without losing structure."
                if tool_name == "Document Studio"
                else "Keep ledger automation and tax guidance current."
                if tool_name == "Accounts Automation"
                else "Route provider choice and model selection cleanly."
                if tool_name == "Provider Router"
                else "Keep the mother brain visible and understandable for operators."
            ),
            motivation=tool.get("identity", "tool capability"),
            learning_mode=f"Learns from {tool.get('identity', 'tool')} usage, requests, and completion quality.",
            monitoring_focus=f"{tool_name} uptime and operator usefulness.",
            parent_id=tool_parent,
            identity="tool",
        )

    children = Counter(brain["parent_id"] for brain in brains if brain.get("parent_id"))
    for brain in brains:
        brain["children_count"] = children.get(brain["id"], 0)

    name_lookup = {brain["id"]: brain["name"] for brain in brains}
    layer_counts = Counter(brain["layer"] for brain in brains)
    total_thoughts = 0
    for brain in brains:
        brain["parent_name"] = name_lookup.get(brain.get("parent_id"), "")
        doctrine = brain_role_profile(brain)
        brain.update(doctrine)
        brain.update(brain_nlp_profile(brain))
        brain.update(brain_auto_repair_profile(brain, current))
        brain["training_modules"] = len(brain_doctrine_lessons(brain))
        brain["domain"] = format_brain_scope(brain.get("layer", "brain"))
        total_thoughts += int(brain.get("thoughts_processed", 0) or 0)
    enrichment_lookup = INTERPRETER_LAYER.get("brain_enrichment_snapshot", lambda _ids=None: {})(
        [brain.get("id") for brain in brains if brain.get("id")]
    )
    for brain in brains:
        enrich = enrichment_lookup.get(brain.get("id"), {})
        if enrich:
            brain.update(enrich)
    permission_relays = [
        {
            "from": "mother_brain",
            "to": "cabinet_brain",
            "permission": ["guide", "approve", "assign"],
            "status": "live",
            "reason": "The mother brain approves the executive mandate before the cabinet moves.",
        },
        {
            "from": "mother_brain",
            "to": "akshaya_brain",
            "permission": ["preserve", "recall", "compress"],
            "status": "live",
            "reason": "Memory and continuity remain under direct mother supervision.",
        },
        {
            "from": "mother_brain",
            "to": "carbon_brain",
            "permission": ["bond", "route", "synchronize"],
            "status": "live",
            "reason": "The nervous system needs a universal bonding layer to keep every child brain connected.",
        },
        {
            "from": "mother_brain",
            "to": "interpreter_brain",
            "permission": ["translate", "normalize", "route"],
            "status": "live",
            "reason": "The interpreter bridge keeps human commands and brain output aligned without losing intent.",
        },
        {
            "from": "interpreter_brain",
            "to": "machine_local_llm",
            "permission": ["offline_reasoning", "prompt_translation"],
            "status": "ready",
            "reason": "The local LLM should always be available as the first translation lane.",
        },
        {
            "from": "interpreter_brain",
            "to": "machine_external_bridge",
            "permission": ["manual_external_refine", "quota_guard"],
            "status": "manual_only",
            "reason": "External translation stays manual-first so quota is never burned silently.",
        },
    ]
    permission_relays.extend(
        {
            "from": "cabinet_brain",
            "to": f"secretary_{secretary.get('lane', 'general')}",
            "permission": ["assign", "prioritize", "review"],
            "status": secretary.get("status", "ready"),
            "reason": secretary.get("brief", "Maintain a healthy operating lane."),
        }
        for secretary in secretaries
    )
    permission_relays.extend(
        {
            "from": domain_parent_map.get(domain["id"], "cabinet_brain"),
            "to": f"domain_{domain['id']}",
            "permission": ["operate", "report", "escalate"],
            "status": domain.get("status", "ready"),
            "reason": domain.get("latest_signal", domain.get("purpose", "Awaiting the next instruction.")),
        }
        for domain in domain_rows
    )
    permission_relays.extend(
        {
            "from": tool_parent_map.get(tool.get("name", "Tool"), "cabinet_brain"),
            "to": f"tool_{tool.get('name', 'tool').lower().replace(' ', '_').replace('+', 'plus').replace('-', '_')}",
            "permission": ["use", "assist", "report"],
            "status": "active",
            "reason": f"{tool.get('name', 'Tool')} remains available as a child capability under supervision.",
        }
        for tool in intelligence_rows.get("tools", [])
    )
    for relay in permission_relays:
        relay["from_name"] = name_lookup.get(relay["from"], relay["from"])
        relay["to_name"] = name_lookup.get(relay["to"], relay["to"])

    motivation_streams = [
        {
            "target": "Cabinet",
            "source": name_lookup["mother_brain"],
            "message": objective,
            "focus": "Turn guidance into visible action.",
        },
        {
            "target": "Accounts Command",
            "source": name_lookup["mother_brain"],
            "message": accounts["headline"],
            "focus": "Revenue discipline, local tax awareness, and clean ledgers.",
        },
        {
            "target": "Praapti",
            "source": name_lookup["mother_brain"],
            "message": latest_hunt.get("job_description", "Awaiting the next live role to hunt.")[:180],
            "focus": "Convert demand into candidates, interviews, and revenue movement.",
        },
        {
            "target": "Nirmaan",
            "source": name_lookup["mother_brain"],
            "message": latest_proposal.get("title", "Propose only the strongest next evolution."),
            "focus": "Create useful growth, not noise.",
        },
        {
            "target": "Swarm",
            "source": name_lookup["mother_brain"],
            "message": latest_mission.get("mission", "Keep the next mission bounded, clear, and reportable."),
            "focus": "Execution discipline and clean report flow.",
        },
        {
            "target": "Voice + Portals",
            "source": name_lookup["mother_brain"],
            "message": f"Active avatars: {' + '.join(active_avatars)}. Voice relay: {last_voice[:140]}",
            "focus": "Keep the human experience calm, clear, and aligned to the mother chain.",
        },
    ]
    avg_learning_score = round(
        sum(float(brain.get("learning_score", 0.0) or 0.0) for brain in brains) / max(len(brains), 1),
        3,
    )

    return {
        "summary": {
            "total_brains": len(brains),
            "layers": len(layer_counts),
            "mother_children": children.get("mother_brain", 0),
            "executives": layer_counts.get("executive", 0),
            "secretaries": layer_counts.get("secretary", 0),
            "domains": layer_counts.get("domain", 0),
            "machines": layer_counts.get("machine", 0),
            "tools": layer_counts.get("tool", 0),
            "permission_relays": len(permission_relays),
            "alive": len([brain for brain in brains if brain.get("status") not in {"idle", "standby"}]),
            "total_thoughts": total_thoughts,
            "avg_learning_score": avg_learning_score,
            "doctrine_modules": sum(int(brain.get("training_modules", 0) or 0) for brain in brains),
            "nlp_brains": len([brain for brain in brains if brain.get("nlp_enabled")]),
            "nlp_modules": sum(int(brain.get("nlp_modules", 0) or 0) for brain in brains),
            "learned_skills": sum(int(brain.get("skill_count", 0) or 0) for brain in brains),
            "mutation_ready": len([brain for brain in brains if brain.get("mutation_ready")]),
            "auto_repair_brains": len([brain for brain in brains if brain.get("auto_repair_enabled")]),
            "auto_repair_issues": int(repair_engine.get("issue_count", 0) or 0),
        },
        "layers": dict(layer_counts),
        "brains": brains,
        "permission_relays": permission_relays,
        "motivation_streams": motivation_streams,
        "auto_repair": repair_engine,
    }


def seed_market_ready_brains() -> Dict[str, int]:
    seeded = 0
    skipped = 0
    try:
        hierarchy = brain_hierarchy_payload()
    except Exception:
        return {"seeded": 0, "skipped": 0}
    for brain in hierarchy.get("brains", []):
        layer = brain.get("layer", "tool")
        relevance = {
            "mother": 0.99,
            "executive": 0.96,
            "secretary": 0.92,
            "domain": 0.9,
            "machine": 0.88,
            "tool": 0.86,
        }.get(layer, 0.82)
        for lesson in MARKET_READY_LESSONS:
            existing = db_one(
                """
                SELECT id FROM brain_knowledge
                WHERE brain_id=? AND source_type='market_ready_seed' AND title=?
                """,
                (brain["id"], lesson["title"]),
            )
            if existing:
                skipped += 1
                continue
            content = (
                f"Brain: {brain['name']}\n"
                f"Layer: {layer}\n"
                f"Authority: {brain.get('authority', '')}\n"
                f"Assigned task: {brain.get('assigned_task', '')}\n"
                f"Lesson: {lesson['summary']}\n"
                "This lesson is permanent baseline operating knowledge for TechBuzz Systems Pvt Ltd as a real-world recruitment, delivery, and execution company."
            )
            db_exec(
                """
                INSERT INTO brain_knowledge(id,brain_id,source_type,source_url,title,content,summary,keywords,relevance_score,learned_at)
                VALUES(?,?,?,?,?,?,?,?,?,?)
                """,
                (
                    new_id("know"),
                    brain["id"],
                    "market_ready_seed",
                    f"market://{lesson['key']}",
                    lesson["title"],
                    content[:4000],
                    lesson["summary"][:500],
                    json.dumps(lesson["keywords"], ensure_ascii=False),
                    relevance,
                    now_iso(),
                ),
            )
            seeded += 1
    doctrine_result = seed_brain_role_doctrine(hierarchy)
    return {
        "seeded": seeded + doctrine_result["seeded"],
        "skipped": skipped + doctrine_result["skipped"],
        "market_ready_seeded": seeded,
        "doctrine_seeded": doctrine_result["seeded"],
    }


def seed_brain_role_doctrine(hierarchy: Optional[Dict[str, Any]] = None) -> Dict[str, int]:
    seeded = 0
    skipped = 0
    try:
        hierarchy_payload = hierarchy or brain_hierarchy_payload()
    except Exception:
        return {"seeded": 0, "skipped": 0}
    for brain in hierarchy_payload.get("brains", []):
        for lesson in brain_doctrine_lessons(brain):
            existing = db_one(
                """
                SELECT id FROM brain_knowledge
                WHERE brain_id=? AND source_type=? AND title=?
                """,
                (brain["id"], lesson["source_type"], lesson["title"]),
            )
            if existing:
                skipped += 1
                continue
            db_exec(
                """
                INSERT INTO brain_knowledge(id,brain_id,source_type,source_url,title,content,summary,keywords,relevance_score,learned_at)
                VALUES(?,?,?,?,?,?,?,?,?,?)
                """,
                (
                    new_id("know"),
                    brain["id"],
                    lesson["source_type"],
                    lesson["source_url"][:400],
                    lesson["title"][:180],
                    lesson["content"][:4000],
                    lesson["summary"][:500],
                    json.dumps(lesson["keywords"], ensure_ascii=False),
                    float(lesson.get("relevance", 0.92)),
                    now_iso(),
                ),
            )
            seeded += 1
    return {"seeded": seeded, "skipped": skipped}


def brain_training_map_payload() -> Dict[str, Any]:
    hierarchy = brain_hierarchy_payload()
    layer_order = ["mother", "executive", "secretary", "domain", "machine", "tool", "atom"]
    grouped = []
    for layer in layer_order:
        layer_brains = [brain for brain in hierarchy.get("brains", []) if brain.get("layer") == layer]
        if not layer_brains:
            continue
        grouped.append({"layer": layer, "count": len(layer_brains), "brains": layer_brains})
    return {"summary": hierarchy.get("summary", {}), "layers": grouped, "brains": hierarchy.get("brains", [])}


def can_control_brain(user: Dict[str, Any], brain_id: str) -> bool:
    if not user:
        return False
    if user.get("role") == "master":
        return True
    if brain_id == "all":
        return False
    return brain_id not in {
        "mother_brain",
        "cabinet_brain",
        "akshaya_brain",
        "carbon_brain",
        "interpreter_brain",
        "machine_local_llm",
        "machine_external_bridge",
    }


def find_brain(brain_id: str) -> Optional[Dict[str, Any]]:
    for brain in brain_hierarchy_payload().get("brains", []):
        if brain.get("id") == brain_id:
            return brain
    return None


def update_brain_directives(brain_ids: List[str], patch_builder) -> List[Dict[str, Any]]:
    updated_rows: List[Dict[str, Any]] = []

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        directives = state["cabinet"].setdefault("brain_directives", {})
        for brain_id in brain_ids:
            directive = directives.setdefault(brain_id, {})
            patch = patch_builder(brain_id, directive)
            if isinstance(patch, dict):
                cleaned_patch: Dict[str, Any] = {}
                for key, value in patch.items():
                    if key == "assigned_task" and isinstance(value, str):
                        cleaned_patch[key] = sanitize_operator_line(value, default_brain_assignment(brain_id))
                    elif key in {"motivation", "last_thought"} and isinstance(value, str):
                        fallback = (
                            "Keep the lane clear, useful, and operator-safe."
                            if key == "motivation"
                            else f"{brain_id} is tracking a clean operator-safe directive."
                        )
                        cleaned_patch[key] = sanitize_operator_multiline(value, fallback)
                    else:
                        cleaned_patch[key] = value
                patch = cleaned_patch
            directive.update(patch)
            directive["updated_at"] = now_iso()
            updated_rows.append({"brain_id": brain_id, **directive})
        return directives

    mutate_state(_mutate)
    return updated_rows


def thought_for_brain(brain: Dict[str, Any], context: str = "") -> str:
    deliverable = (brain.get("deliverables") or ["One clean next step"])[0]
    learning_target = (brain.get("learning_targets") or ["Improve judgement"])[0]
    mission = brain.get("mission", "Keep the lane healthy.")
    context_text = sanitize_operator_line(context, "")
    return (
        f"{brain['name']} is operating as {brain.get('role_title', brain.get('layer', 'brain'))}. "
        f"Mission: {mission} "
        f"Immediate deliverable: {deliverable}. "
        f"Current assignment: {brain.get('assigned_task', 'Await the next instruction.')}. "
        f"Learning focus: {learning_target}."
        + (f" Context: {context_text}." if context_text else "")
    )


def account_status_payload(user: Dict[str, Any]) -> Dict[str, Any]:
    profile = ensure_account_profile(user["id"])
    entries = account_entries_for_user(user["id"], limit=40)
    reports = account_reports_for_user(user["id"], limit=8)
    preset = account_region_preset(profile.get("region_code"))
    income_total = round(sum(float(row.get("amount", 0) or 0) for row in entries if row.get("entry_type") in ACCOUNT_INFLOW_TYPES), 2)
    expense_total = round(sum(float(row.get("amount", 0) or 0) for row in entries if row.get("entry_type") in ACCOUNT_OUTFLOW_TYPES), 2)
    tax_collected = round(sum(float(row.get("tax_amount", 0) or 0) for row in entries if row.get("entry_type") in ACCOUNT_INFLOW_TYPES), 2)
    tax_input = round(sum(float(row.get("tax_amount", 0) or 0) for row in entries if row.get("entry_type") in ACCOUNT_OUTFLOW_TYPES), 2)
    estimated_tax_payable = round(max(0.0, tax_collected - tax_input), 2)
    net_cashflow = round(income_total - expense_total, 2)
    tds_deducted = account_tds_total(entries)
    pending_invoices = account_pending_invoices(reports)
    missing_counterparty = sum(1 for row in entries if not (row.get("counterparty") or "").strip())
    counterparty_counts = Counter((row.get("counterparty") or "Unspecified").strip() for row in entries if row.get("counterparty"))
    top_counterparty = counterparty_counts.most_common(1)[0][0] if counterparty_counts else "No dominant counterparty yet"
    guidance: List[Dict[str, str]] = []
    if not entries:
        guidance.append(
            {
                "title": "Seed the ledger",
                "summary": f"Choose the {preset['label']} profile, confirm the {profile.get('tax_name', preset['tax_name'])} rate, and log the first income or expense entry so Accounts Command can start learning your flow.",
            }
        )
    if not (profile.get("tax_registration") or "").strip():
        guidance.append(
            {
                "title": "Add tax registration details",
                "summary": f"Capture your {profile.get('tax_name', preset['tax_name'])} registration or internal reference so reports can be matched to the correct filing identity.",
            }
        )
    if estimated_tax_payable > 0:
        guidance.append(
            {
                "title": f"Reserve {profile.get('tax_name', preset['tax_name'])}",
                "summary": f"Current ledger suggests about {estimated_tax_payable:.2f} {profile.get('currency', preset['currency'])} may need to be reserved as output tax after input offsets. Review with your accountant before filing.",
            }
        )
    if net_cashflow < 0:
        guidance.append(
            {
                "title": "Watch cash burn",
                "summary": f"Expenses exceed income by {abs(net_cashflow):.2f} {profile.get('currency', preset['currency'])}. Tighten spends or open new revenue lanes before the next filing cycle.",
            }
        )
    if missing_counterparty:
        guidance.append(
            {
                "title": "Capture counterparties cleanly",
                "summary": f"{missing_counterparty} ledger entr{'y' if missing_counterparty == 1 else 'ies'} are missing a client or vendor name. Completing that field improves tax, billing, and account intelligence.",
            }
        )
    if not guidance:
        guidance.append(
            {
                "title": "Accounts relay is healthy",
                "summary": f"The ledger is balanced for now. Dominant counterparty: {top_counterparty}. Keep feeding clean entries to strengthen automated reporting.",
            }
        )
    automation = {
        "name": "Accounts Command",
        "status": "learning" if entries else "ready",
        "headline": f"{preset['label']} profile active with {len(entries)} entries and {len(reports)} stored report(s).",
        "notes": preset["notes"],
    }
    return {
        "profile": profile,
        "regions": account_region_catalog(),
        "entry_types": [
            {"id": "income", "label": "Income / Invoice"},
            {"id": "expense", "label": "Expense / Purchase"},
            {"id": "salary", "label": "Salary / Payroll"},
            {"id": "tax_payment", "label": "Tax Payment"},
        ],
        "metrics": {
            "income_total": income_total,
            "expense_total": expense_total,
            "net_cashflow": net_cashflow,
            "tax_collected": tax_collected,
            "tax_input": tax_input,
            "estimated_tax_payable": estimated_tax_payable,
            "entries_count": len(entries),
            "reports_count": len(reports),
        },
        "automation": automation,
        "guidance": guidance,
        "entries": entries,
        "ledger": account_legacy_ledger_rows(entries),
        "reports": reports,
        "top_counterparty": top_counterparty,
        "pending_invoices": pending_invoices,
        "regional_tax": {
            "code": profile.get("region_code") or preset["code"],
            "state": preset["label"],
            "tax_name": profile.get("tax_name") or preset["tax_name"],
            "default_rate": float(profile.get("default_tax_rate", preset["suggested_tax_rate"]) or 0),
            "filing_cycle": profile.get("filing_cycle") or preset["filing_cycle"],
        },
        "financials": {
            "income": income_total,
            "expenses": expense_total,
            "net": net_cashflow,
            "gst_liability": estimated_tax_payable,
            "gst_collected": tax_collected,
            "tds_deducted": tds_deducted,
        },
    }


def store_uploaded_file(upload: UploadFile, user_id: str) -> Dict[str, Any]:
    safe_name = Path(upload.filename or "document.bin").name
    extension = Path(safe_name).suffix.lower() or ".bin"
    target = DOCUMENT_DIR / user_id
    target.mkdir(parents=True, exist_ok=True)
    storage_path = target / f"{uuid.uuid4().hex}{extension}"
    upload.file.seek(0)
    with storage_path.open("wb") as output:
        shutil.copyfileobj(upload.file, output)
    mime_type = upload.content_type or mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
    text = extract_document_text(storage_path, mime_type)
    return create_document_record(
        user_id=user_id,
        original_name=safe_name,
        mime_type=mime_type,
        extension=extension,
        storage_path=storage_path,
        size_bytes=storage_path.stat().st_size,
        extracted_text=text,
        kind="upload",
    )


def get_document_for_user(user_id: str, document_id: str) -> Dict[str, Any]:
    row = db_one("SELECT * FROM documents WHERE id=? AND user_id=?", (document_id, user_id))
    if not row:
        raise HTTPException(status_code=404, detail="Document not found.")
    return row


def create_text_document(user_id: str, title: str, content: str, extension: str = ".txt") -> Dict[str, Any]:
    safe_title = re.sub(r"[^a-zA-Z0-9._-]+", "_", title or "ishani_note").strip("._") or "ishani_note"
    target = DOCUMENT_EXPORT_DIR / user_id
    target.mkdir(parents=True, exist_ok=True)
    storage_path = target / f"{safe_title}{extension}"
    if extension == ".docx" and DocxDocument:
        doc = DocxDocument()
        for paragraph in (content or "").split("\n\n"):
            doc.add_paragraph(paragraph)
        doc.save(str(storage_path))
        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    else:
        storage_path = storage_path.with_suffix(".txt")
        storage_path.write_text(content or "", encoding="utf-8")
        mime_type = "text/plain"
    return create_document_record(
        user_id=user_id,
        original_name=storage_path.name,
        mime_type=mime_type,
        extension=storage_path.suffix.lower(),
        storage_path=storage_path,
        size_bytes=storage_path.stat().st_size,
        extracted_text=content or "",
        kind="generated",
    )


def pdf_support_available() -> bool:
    return PdfReader is not None and PdfWriter is not None


def merge_pdf_documents(user_id: str, document_ids: List[str]) -> Dict[str, Any]:
    if not pdf_support_available():
        raise HTTPException(status_code=503, detail="PDF tools are not installed yet. Add pypdf to activate merge and split.")
    documents = [get_document_for_user(user_id, doc_id) for doc_id in document_ids]
    if len(documents) < 2:
        raise HTTPException(status_code=400, detail="Select at least two PDF documents to merge.")
    writer = PdfWriter()
    for document in documents:
        if document["extension"] != ".pdf":
            raise HTTPException(status_code=400, detail="Only PDF files can be merged in this operation.")
        reader = PdfReader(document["storage_path"])
        for page in reader.pages:
            writer.add_page(page)
    target = DOCUMENT_EXPORT_DIR / user_id
    target.mkdir(parents=True, exist_ok=True)
    merged_path = target / f"merged_{uuid.uuid4().hex[:8]}.pdf"
    with merged_path.open("wb") as handle:
        writer.write(handle)
    merged_text = extract_document_text(merged_path, "application/pdf")
    return create_document_record(
        user_id=user_id,
        original_name=merged_path.name,
        mime_type="application/pdf",
        extension=".pdf",
        storage_path=merged_path,
        size_bytes=merged_path.stat().st_size,
        extracted_text=merged_text,
        kind="merge_pdf",
    )


def split_pdf_document(user_id: str, document_id: str, start_page: int, end_page: Optional[int]) -> Dict[str, Any]:
    if not pdf_support_available():
        raise HTTPException(status_code=503, detail="PDF tools are not installed yet. Add pypdf to activate merge and split.")
    document = get_document_for_user(user_id, document_id)
    if document["extension"] != ".pdf":
        raise HTTPException(status_code=400, detail="Only PDF files can be split in this operation.")
    reader = PdfReader(document["storage_path"])
    start = max(1, int(start_page))
    finish = min(len(reader.pages), int(end_page or len(reader.pages)))
    if start > finish:
        raise HTTPException(status_code=400, detail="Invalid page range for split.")
    writer = PdfWriter()
    for index in range(start - 1, finish):
        writer.add_page(reader.pages[index])
    target = DOCUMENT_EXPORT_DIR / user_id
    target.mkdir(parents=True, exist_ok=True)
    split_path = target / f"split_{uuid.uuid4().hex[:8]}_{start}-{finish}.pdf"
    with split_path.open("wb") as handle:
        writer.write(handle)
    split_text = extract_document_text(split_path, "application/pdf")
    return create_document_record(
        user_id=user_id,
        original_name=split_path.name,
        mime_type="application/pdf",
        extension=".pdf",
        storage_path=split_path,
        size_bytes=split_path.stat().st_size,
        extracted_text=split_text,
        kind="split_pdf",
    )


def mother_monitor_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    cabinet = prime_minister_payload(current)
    domains = operational_domains_payload(current)
    nervous_system = nervous_system_payload(current)
    memory_audit = memory_audit_payload(current)
    accounts = global_accounts_snapshot()
    monitoring = current.get("monitoring", {})
    repair_engine = auto_repair_engine_payload(current)
    engine_reports = list(reversed(monitoring.get("engine_reports", [])))[:12]
    engines = [
        {
            "id": "cabinet",
            "name": cabinet["prime_minister"].get("name", "Prime Minister"),
            "status": cabinet["prime_minister"].get("status", "governing"),
            "summary": cabinet.get("latest_report", "")[:220],
        },
        {
            "id": "akshaya",
            "name": "Akshaya Vault",
            "status": "preserving",
            "summary": f"{memory_audit['counts']['vault']} vault items, {memory_audit['footprint_kb']} KB footprint, {memory_audit['archive_cycles']} archive cycle(s).",
        },
        {
            "id": "praapti",
            "name": "Praapti Recruitment",
            "status": "active" if current.get("praapti_hunts") else "ready",
            "summary": f"{len(current.get('praapti_hunts', []))} total hunt(s) with {len(current.get('packages', []))} package-supported lanes.",
        },
        {
            "id": "nirmaan",
            "name": "Nirmaan Chakra",
            "status": "active" if current.get("nirmaan_proposals") else "ready",
            "summary": f"{len(current.get('nirmaan_proposals', []))} proposal(s) tracked for staged evolution.",
        },
        {
            "id": "swarm",
            "name": "Swarm Intelligence",
            "status": "active" if current.get("swarm_missions") else "ready",
            "summary": f"{len(current.get('swarm_missions', []))} mission(s) available for replay and review.",
        },
        {
            "id": "voice",
            "name": "Voice Relay",
            "status": "listening" if current.get("voice", {}).get("always_listening") else "idle",
            "summary": sanitize_operator_line(
                current.get("voice", {}).get("last_command", ""),
                "No voice command recorded yet.",
            ),
        },
        {
            "id": "accounts",
            "name": "Accounts Command",
            "status": "learning" if accounts["entries"] else "ready",
            "summary": accounts["headline"],
        },
        {
            "id": "auto_repair",
            "name": "Auto Repair Engine",
            "status": repair_engine.get("status", "ready"),
            "summary": repair_engine.get("last_report", "Safe self-check layer is ready."),
        },
    ]
    alerts = nervous_system.get("alerts", [])
    if not alerts:
        alerts = [{"level": "ok", "title": "All Major Systems Connected", "summary": "The mother brain sees a healthy chain from database to agents and back."}]
    if not engine_reports:
        engine_reports = [
            {
                "engine": "Mother Monitor",
                "status": "live",
                "summary": "Monitoring is active. Once cabinet or mission cycles run, fresh engine reports will appear here.",
                "created_at": monitoring.get("last_scan_at", now_iso()),
            }
        ]
    return {
        "status": "alive",
        "heartbeat_ms": int(monitoring.get("pulse_interval_ms", 2500)),
        "last_scan_at": monitoring.get("last_scan_at", now_iso()),
        "engines": engines,
        "alerts": alerts,
        "reports": engine_reports,
        "domains": domains,
        "nervous_system": nervous_system,
        "memory_audit": memory_audit,
        "prana_nadi": prana_nadi_payload(current),
        "auto_repair": repair_engine,
    }


def build_secretary_swarm(state: Dict[str, Any], objective: str) -> List[Dict[str, Any]]:
    hunts = state.get("praapti_hunts", [])
    packages = state.get("packages", [])
    proposals = [row for row in state.get("nirmaan_proposals", []) if not row.get("approved")]
    vault_items = len(state.get("vault", []))
    conversations = len(state.get("conversations", []))
    clients = sorted({hunt.get("client_company", COMPANY_NAME) for hunt in hunts})
    accounts = global_accounts_snapshot()
    base_secretaries = [
        {
            "name": "Revenue Secretary",
            "lane": "revenue",
            "status": "active" if hunts or packages else "warming",
            "priority": min(99, 74 + len(hunts) * 6 + len(packages) * 5),
            "brief": "Turn hunts, offers, and delivery lanes into booked revenue.",
            "next_move": "Package new offers and prioritize the strongest commercial lanes.",
        },
        {
            "name": "Hiring Secretary",
            "lane": "talent",
            "status": "active" if hunts else "ready",
            "priority": min(99, 62 + len(hunts) * 8),
            "brief": "Convert demand into roles, candidate slates, and interview motion.",
            "next_move": "Run Praapti on the most urgent revenue-critical role.",
        },
        {
            "name": "Client Intelligence Secretary",
            "lane": "clients",
            "status": "mapping" if clients else "scanning",
            "priority": min(99, 58 + len(clients) * 7),
            "brief": "Map active client accounts, sectors, and relationship depth.",
            "next_move": "Track where TechBuzz can expand inside current client lanes.",
        },
        {
            "name": "Research Secretary",
            "lane": "research",
            "status": "discovering" if proposals else "watching",
            "priority": min(99, 55 + len(proposals) * 9 + len(packages) * 3),
            "brief": "Spot product, market, and offer ideas that create new revenue surfaces.",
            "next_move": "Turn live signals into experiments with clear output and timing.",
        },
        {
            "name": "Delivery Secretary",
            "lane": "delivery",
            "status": "routing" if packages else "ready",
            "priority": min(99, 57 + len(packages) * 8 + len(hunts) * 3),
            "brief": "Keep work moving from plan to execution without losing context.",
            "next_move": "Bind active missions to owners, outputs, and follow-up dates.",
        },
        {
            "name": "Protection Secretary",
            "lane": "protection",
            "status": "sealed" if state["avatar_state"].get("protection_meter", 0) >= 80 else "hardening",
            "priority": min(99, 60 + state["avatar_state"].get("protection_meter", 0) // 2),
            "brief": "Guard the core, risk posture, privacy controls, and system continuity.",
            "next_move": "Preserve core memory and keep the command layer stable under load.",
        },
        {
            "name": "Accounts Secretary",
            "lane": "accounts",
            "status": "learning" if accounts["entries"] else "ready",
            "priority": min(99, 58 + accounts["entries"] // 2 + accounts["reports"] * 6 + accounts["profiles"] * 4),
            "brief": "Maintain the regional tax profile, guide bookkeeping, and keep revenue and compliance signals clean.",
            "next_move": accounts["headline"],
        },
        {
            "name": "Platform Secretary",
            "lane": "platform",
            "status": "upgrading" if proposals else "steady",
            "priority": min(99, 54 + len(proposals) * 8 + conversations // 5),
            "brief": "Improve the interface, backend routes, and operating ergonomics.",
            "next_move": "Remove friction from the most-used surfaces first.",
        },
        {
            "name": "Quantum Storage Secretary",
            "lane": "storage",
            "status": "preserving",
            "priority": min(99, 61 + vault_items // 6),
            "brief": "Manage Akshaya preservation, archival rollups, and retrieval accuracy.",
            "next_move": "Keep memory lightweight while preserving the highest-value traces.",
        },
    ]
    dynamic_secretaries: List[Dict[str, Any]] = []
    for hunt in hunts[-6:]:
        company = hunt.get("client_company", COMPANY_NAME)
        dynamic_secretaries.append(
            {
                "name": f"{company} Secretary",
                "lane": "client",
                "status": "briefing",
                "priority": min(99, 56 + len(hunt.get("candidates", [])) * 8),
                "brief": f"Own the current hunt, candidate brief, and relationship lane for {company}.",
                "next_move": f"Convert the {company} hunt into interviews, a shortlist, and delivery hooks.",
            }
        )
    for package in packages[-4:]:
        dynamic_secretaries.append(
            {
                "name": f"{package.get('title', 'Mission')} Secretary",
                "lane": "mission",
                "status": "executing",
                "priority": min(99, 52 + len(package.get("report", "")) // 70),
                "brief": f"Own the bounded mission for {package.get('title', 'current package')}.",
                "next_move": "Translate the package report into a specific output for the core team.",
            }
        )
    for proposal in proposals[-4:]:
        dynamic_secretaries.append(
            {
                "name": f"{proposal.get('title', 'Nirmaan')} Secretary",
                "lane": "evolution",
                "status": "designing",
                "priority": min(99, 60 + len(proposal.get("description", "")) // 40),
                "brief": "Own the path from proposal to approved improvement.",
                "next_move": "Prepare the proposal for execution or archive it cleanly.",
            }
        )
    for domain in operational_domains_payload(state)[:5]:
        dynamic_secretaries.append(
            {
                "name": f"{domain['name']} Liaison",
                "lane": domain["id"],
                "status": domain["status"],
                "priority": min(99, int(domain["priority"])),
                "brief": domain["purpose"],
                "next_move": domain["latest_signal"],
            }
        )
    if not dynamic_secretaries:
        dynamic_secretaries.append(
            {
                "name": "Expansion Secretary",
                "lane": "growth",
                "status": "ready",
                "priority": 58,
                "brief": "Stand ready to open the next revenue or research lane as soon as the mandate shifts.",
                "next_move": "Wait for the next objective and convert it into a clean mission lane.",
            }
        )
    secretaries = base_secretaries + dynamic_secretaries
    secretaries.sort(key=lambda item: item.get("priority", 0), reverse=True)
    return secretaries


def cabinet_summary_lines(state: Dict[str, Any], objective: str, secretaries: List[Dict[str, Any]]) -> List[str]:
    hunts_today = sum(1 for row in state["praapti_hunts"] if row.get("created_day") == today_iso())
    pending_proposals = sum(1 for row in state["nirmaan_proposals"] if not row.get("approved"))
    packages_active = len(state["packages"])
    revenue_projection = round(1.4 + hunts_today * 0.38 + packages_active * 0.28 + pending_proposals * 0.14, 2)
    top_secretaries = secretaries[:5]
    lines = [
        f"Prime Minister mandate: {objective}",
        f"Revenue posture: ₹{revenue_projection} Cr projected from {hunts_today} hunts today, {packages_active} active packages, and {pending_proposals} live proposals.",
        "Secretary briefs:",
    ]
    for secretary in top_secretaries:
        lines.append(
            f"- {secretary['name']}: {secretary['brief']} Next move: {secretary['next_move']}"
        )
    lines.append("Core order: protect the core, accelerate delivery, and convert intelligence into revenue.")
    return lines


def prime_minister_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    cabinet = current["cabinet"]
    prime_minister = cabinet["prime_minister"]
    objective = prime_minister.get(
        "objective",
        f"Generate revenue for {COMPANY_NAME} while protecting the core and improving delivery.",
    )
    secretaries = build_secretary_swarm(current, objective)
    mission_log = list(reversed(cabinet.get("mission_log", [])))[:10]
    hunts_today = sum(1 for row in current["praapti_hunts"] if row.get("created_day") == today_iso())
    pending_proposals = sum(1 for row in current["nirmaan_proposals"] if not row.get("approved"))
    packages_active = len(current["packages"])
    revenue_projection = round(1.4 + hunts_today * 0.38 + packages_active * 0.28 + pending_proposals * 0.14, 2)
    quantum_storage = {
        **cabinet.get("quantum_storage", {}),
        "state_file": str(STATE_PATH),
        "footprint_kb": state_file_size_kb(),
        "items_preserved": cabinet.get("quantum_storage", {}).get("items_preserved", 0),
        "accuracy_mode": "structured summaries plus source preservation",
    }
    return {
        "prime_minister": {
            **prime_minister,
            "objective": objective,
            "active_secretaries": len(secretaries),
        },
        "secretaries": secretaries,
        "revenue_board": {
            "focus": "Revenue generation for TechBuzz Systems Pvt Ltd",
            "projected_revenue_inr": revenue_projection,
            "hunts_today": hunts_today,
            "packages_active": packages_active,
            "pending_proposals": pending_proposals,
            "priority_order": prime_minister.get("current_order", ""),
        },
        "quantum_storage": quantum_storage,
        "latest_report": prime_minister.get("last_report") or "\n".join(cabinet_summary_lines(current, objective, secretaries)),
        "mission_log": mission_log,
        "elasticity": {
            "model": "elastic secretary swarm",
            "summary": "The cabinet grows by workload. Base secretaries stay permanent while mission-specific secretaries appear for hunts, packages, and proposals.",
        },
    }


def run_prime_minister_cycle(*, objective: Optional[str] = None, command: Optional[str] = None, source: str = "manual") -> Dict[str, Any]:
    result: Dict[str, Any] = {}

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        nonlocal result
        cabinet = state["cabinet"]
        prime_minister = cabinet["prime_minister"]
        if objective:
            prime_minister["objective"] = objective[:260]
        if command:
            prime_minister["current_order"] = command[:260]
        secretaries = build_secretary_swarm(state, prime_minister["objective"])
        cabinet["secretaries"] = secretaries
        prime_minister["status"] = "governing" if prime_minister.get("enabled", True) else "paused"
        prime_minister["last_cycle_at"] = now_iso()
        interval_seconds = max(15, int(prime_minister.get("interval_seconds", 45)))
        prime_minister["next_cycle_at"] = (datetime.now(UTC) + timedelta(seconds=interval_seconds)).isoformat()
        report = "\n".join(cabinet_summary_lines(state, prime_minister["objective"], secretaries))
        prime_minister["last_report"] = report[:2400]
        cycle_record = {
            "id": new_id("cabinet"),
            "source": source,
            "objective": prime_minister["objective"],
            "command": prime_minister.get("current_order", ""),
            "report": report[:2400],
            "created_at": now_iso(),
            "top_secretaries": [secretary["name"] for secretary in secretaries[:6]],
        }
        cabinet["mission_log"].append(cycle_record)
        cabinet["quantum_storage"]["state_file"] = str(STATE_PATH)
        cabinet["quantum_storage"]["items_preserved"] = (
            len(state.get("vault", []))
            + len(state.get("conversations", []))
            + len(state.get("praapti_hunts", []))
            + len(state.get("packages", []))
            + len(cabinet.get("mission_log", []))
        )
        monitoring = state.setdefault("monitoring", {})
        monitoring["last_scan_at"] = now_iso()
        monitoring["last_runtime_cycle_at"] = prime_minister["last_cycle_at"]
        monitoring["last_report"] = report[:1200]
        engine_reports = monitoring.setdefault("engine_reports", [])
        engine_reports.append(
            {
                "id": new_id("monitor"),
                "engine": "Prime Minister Cabinet",
                "status": prime_minister["status"],
                "summary": report[:320],
                "created_at": now_iso(),
            }
        )
        monitoring_alerts = monitoring.setdefault("alerts", [])
        monitoring_alerts.append(
            {
                "level": "ok",
                "title": "Cabinet Cycle Completed",
                "summary": f"{prime_minister['name']} briefed {len(secretaries)} secretaries and refreshed the state chain.",
                "created_at": now_iso(),
            }
        )
        state["vault"].append(
            {
                "id": new_id("vault"),
                "kind": "prime_minister_cycle",
                "title": "Prime Minister Revenue Cycle",
                "summary": report[:280],
                "content": json.dumps(cycle_record, ensure_ascii=False)[:2400],
                "created_at": now_iso(),
            }
        )
        result = {
            "message": f"{prime_minister['name']} completed a governance cycle and briefed {len(secretaries)} secretaries.",
            "report": report,
            "cabinet": prime_minister_payload(state),
        }
        return result

    return mutate_state(_mutate)


def cabinet_loop_step() -> None:
    state = get_state()
    prime_minister = state["cabinet"]["prime_minister"]
    if not prime_minister.get("enabled", True):
        return
    last_cycle = parse_optional_iso(prime_minister.get("last_cycle_at", ""))
    interval_seconds = max(15, int(prime_minister.get("interval_seconds", 45)))
    due = last_cycle is None or (datetime.now(UTC) - last_cycle).total_seconds() >= interval_seconds
    if due:
        run_prime_minister_cycle(source="runtime")


def cabinet_runtime_loop(stop_event: threading.Event) -> None:
    while not stop_event.wait(18):
        try:
            cabinet_loop_step()
        except Exception as exc:
            log.exception("Prime Minister runtime loop error: %s", exc)


def start_cabinet_loop() -> None:
    global CABINET_LOOP_THREAD
    if CABINET_LOOP_THREAD and CABINET_LOOP_THREAD.is_alive():
        return
    CABINET_LOOP_STOP.clear()
    try:
        cabinet_loop_step()
    except Exception as exc:
        log.exception("Prime Minister startup cycle error: %s", exc)
    CABINET_LOOP_THREAD = threading.Thread(
        target=cabinet_runtime_loop,
        args=(CABINET_LOOP_STOP,),
        name="prime-minister-loop",
        daemon=True,
    )
    CABINET_LOOP_THREAD.start()


def stop_cabinet_loop() -> None:
    CABINET_LOOP_STOP.set()


PACKAGE_TEMPLATES: List[Dict[str, Any]] = [
    {
        "id": "revenue_radar",
        "title": "Revenue Radar",
        "summary": "A bounded research package for finding revenue angles, target sectors, and offer positioning.",
        "best_for": "Growth, outbound, revenue planning",
    },
    {
        "id": "talent_scout",
        "title": "Talent Scout",
        "summary": "A bounded research package for market talent mapping, role landscapes, and sourcing strategy.",
        "best_for": "Hiring, recruitment, candidate mapping",
    },
    {
        "id": "signal_lens",
        "title": "Signal Lens",
        "summary": "A bounded listening package for trend spotting, product signals, and company monitoring.",
        "best_for": "Market intelligence, trends, monitoring",
    },
    {
        "id": "fortress_guard",
        "title": "Fortress Guard",
        "summary": "A bounded defensive package for risk review, resilience checks, and operational hardening.",
        "best_for": "Protection, resilience, process review",
    },
]


def akshaya_save(kind: str, title: str, summary: str, payload: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    entry = {
        "id": new_id("vault"),
        "kind": kind,
        "title": title,
        "summary": summary[:500],
        "content": json.dumps(payload or {}, ensure_ascii=False)[:3000],
        "created_at": now_iso(),
    }

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        state["vault"].append(entry)
        return entry

    return mutate_state(_mutate)


def learning_target_brains(limit: Optional[int] = None) -> List[str]:
    brains = [brain.get("id") for brain in brain_hierarchy_payload().get("brains", []) if brain.get("id")]
    return brains[:limit] if limit else brains


def broadcast_brain_lesson(
    *,
    title: str,
    summary: str,
    content: str,
    source_type: str,
    source_url: str = "",
    keywords: Optional[List[str]] = None,
    target_brains: Optional[List[str]] = None,
    relevance: float = 0.84,
) -> Dict[str, Any]:
    targets = target_brains or learning_target_brains()
    learned_at = now_iso()
    keywords_json = json.dumps(keywords or [], ensure_ascii=False)
    inserted = 0
    for brain_id in targets:
        db_exec(
            """
            INSERT INTO brain_knowledge(id,brain_id,source_type,source_url,title,content,summary,keywords,relevance_score,learned_at)
            VALUES(?,?,?,?,?,?,?,?,?,?)
            """,
            (
                new_id("know"),
                brain_id,
                source_type,
                source_url[:400],
                title[:180],
                content[:4000],
                summary[:500],
                keywords_json,
                relevance,
                learned_at,
            ),
        )
        inserted += 1
    return {"learned_at": learned_at, "targets": inserted}


def normalize_navigator_url(raw_url: str) -> str:
    candidate = (raw_url or "").strip()
    if not candidate:
        return ""
    if candidate.startswith("/"):
        return candidate
    if re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://", candidate):
        return candidate
    return f"https://{candidate}"


def recent_navigator_sessions(user_id: str, limit: int = 10) -> List[Dict[str, Any]]:
    rows = db_all(
        """
        SELECT id, title, target_url, notes, tags_json, mode, source_hint, created_at
        FROM navigator_sessions
        WHERE user_id=?
        ORDER BY created_at DESC
        LIMIT ?
        """,
        (user_id, limit),
    )
    for row in rows:
        try:
            row["tags"] = json.loads(row.get("tags_json") or "[]")
        except Exception:
            row["tags"] = []
    return rows


def save_navigator_session(
    user: Dict[str, Any],
    *,
    title: str,
    target_url: str = "",
    notes: str = "",
    tags: Optional[List[str]] = None,
    mode: str = "learning",
    source_hint: str = "manual_capture",
) -> Dict[str, Any]:
    clean_title = (title or "Navigator Session").strip()[:180] or "Navigator Session"
    clean_url = normalize_navigator_url(target_url)[:1000]
    clean_notes = (notes or "").strip()[:4000]
    tag_list = [str(tag).strip()[:40] for tag in (tags or []) if str(tag).strip()]
    session = {
        "id": new_id("nav"),
        "user_id": user["id"],
        "title": clean_title,
        "target_url": clean_url,
        "notes": clean_notes,
        "tags_json": json.dumps(tag_list, ensure_ascii=False),
        "mode": (mode or "learning").strip()[:40] or "learning",
        "source_hint": (source_hint or "manual_capture").strip()[:80] or "manual_capture",
        "created_at": now_iso(),
    }
    db_exec(
        """
        INSERT INTO navigator_sessions(id,user_id,title,target_url,notes,tags_json,mode,source_hint,created_at)
        VALUES(?,?,?,?,?,?,?,?,?)
        """,
        (
            session["id"],
            session["user_id"],
            session["title"],
            session["target_url"],
            session["notes"],
            session["tags_json"],
            session["mode"],
            session["source_hint"],
            session["created_at"],
        ),
    )
    lesson_summary = (
        f"{clean_title}. URL: {clean_url or 'local workflow'}. "
        f"Notes: {clean_notes[:260] or 'No notes captured yet.'}"
    )
    broadcast = broadcast_brain_lesson(
        title=f"Navigator Learning: {clean_title}",
        summary=lesson_summary,
        content=(
            f"Operator: {user.get('name', 'Member')}\n"
            f"Mode: {session['mode']}\n"
            f"Source hint: {session['source_hint']}\n"
            f"Target URL: {clean_url}\n"
            f"Tags: {', '.join(tag_list) or 'none'}\n"
            f"Notes:\n{clean_notes or 'No notes supplied.'}"
        ),
        source_type="navigator_capture",
        source_url=clean_url or f"navigator://{session['mode']}",
        keywords=["navigator", "research", "sourcing", *tag_list],
        relevance=0.87,
    )
    akshaya_save(
        "navigator_learning",
        clean_title,
        lesson_summary[:280],
        {
            "url": clean_url,
            "tags": tag_list,
            "mode": session["mode"],
            "brain_broadcast": broadcast,
        },
    )
    session["tags"] = tag_list
    session["brain_broadcast"] = broadcast
    return session


def navigator_status_payload(user: Dict[str, Any]) -> Dict[str, Any]:
    quick_sites = []
    for item in NAVIGATOR_QUICK_SITES:
        external_only = item["kind"] in {"sourcing", "network", "communication", "research"}
        quick_sites.append(
            {
                **item,
                "embed_recommended": item["kind"] == "internal",
                "external_only": external_only,
                "blocked_by_frame": external_only,
                "best_path": "workspace" if item["kind"] == "internal" else "external",
            }
        )
    captures = recent_navigator_sessions(user["id"], limit=12)
    capture_count = int(
        (db_one("SELECT COUNT(*) AS count FROM navigator_sessions WHERE user_id=?", (user["id"],)) or {}).get("count", 0) or 0
    )
    learning_count = int(
        (db_one("SELECT COUNT(*) AS count FROM brain_knowledge WHERE source_type='navigator_capture'") or {}).get("count", 0) or 0
    )
    return {
        "headline": "TechBuzz Navigator is the explicit sourcing and research workspace. Internal TechBuzz pages can embed here; most external sourcing and communication sites work best through the visible external handoff.",
        "launcher": {
            "path": str(NAUKRI_LAUNCHER_PATH),
            "available": NAUKRI_LAUNCHER_PATH.exists(),
            "label": "Naukri Launcher",
            "message": (
                "Local Naukri launcher detected. Use the button to open it visibly."
                if NAUKRI_LAUNCHER_PATH.exists()
                else "No local Naukri launcher detected at the configured path."
            ),
        },
        "quick_sites": quick_sites,
        "recent_sessions": captures,
        "learning": {
            "captures_for_user": capture_count,
            "brain_lessons_total": learning_count,
            "mode": "operator_approved_only",
        },
    }


async def create_nirmaan_proposal(reason: str = "self-development") -> Dict[str, Any]:
    app_code = ""
    ui_code = ""
    state = get_state()
    active_avatar_names = state["avatar_state"].get("active", ["RAMA"])
    try:
        app_code = (BACKEND_DIR / "app.py").read_text(encoding="utf-8")[-2500:]
    except Exception:
        pass
    try:
        ui_code = (FRONTEND_DIR / "leazy.html").read_text(encoding="utf-8")[-2500:]
    except Exception:
        pass

    prompt = (
        f"You are Nirmaan Chakra for {AI_NAME}.\n"
        f"Reason: {reason}\n"
        f"Current active avatars: {', '.join(active_avatar_names)}\n"
        f"Backend excerpt:\n{app_code}\n\n"
        f"Frontend excerpt:\n{ui_code}\n\n"
        "Propose one high-impact new feature that does not duplicate the existing empire dashboard, "
        "Praapti, Swarm, or Akshaya systems.\n"
        "Prefer Vishnu-inspired tool ideas such as preservation, renewal, simplification, protection, or strategic orchestration.\n"
        "Return JSON only with keys: title, description, code_snippet."
    )
    generated = await generate_text(prompt, system=ADMIN_SYSTEM, max_tokens=1200, use_web_search=False)
    proposal_data = parse_json_blob(generated["text"]) or {}
    proposal = {
        "id": new_id("nirmaan"),
        "title": str(proposal_data.get("title") or "Autonomous Workflow Board")[:140],
        "description": str(
            proposal_data.get("description")
            or "A new self-generated feature designed to make the empire stronger without increasing noise."
        )[:1200],
        "code_snippet": str(
            proposal_data.get("code_snippet")
            or "def autonomous_workflow_board(tasks):\n    return {'tasks': tasks, 'status': 'ready'}\n"
        )[:20000],
        "approved": False,
        "created_at": now_iso(),
        "reason": reason[:200],
        "provider": generated["provider"],
        "avatars": active_avatar_names,
    }

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        proposal_key = normalize_label_key(f"{proposal['title']}:{proposal['reason']}:False")
        title_key = normalize_label_key(proposal["title"])
        now = datetime.now(UTC)
        for existing in reversed(state["nirmaan_proposals"]):
            existing_title_key = normalize_label_key(existing.get("title", ""))
            created_at = parse_optional_iso(existing.get("created_at", ""))
            if existing_title_key == title_key and created_at and (now - created_at) <= timedelta(hours=24):
                return existing
            existing_key = normalize_label_key(
                f"{existing.get('title', '')}:{existing.get('reason', '')}:{existing.get('approved', False)}"
            )
            if existing_key == proposal_key:
                return existing
        state["nirmaan_proposals"].append(proposal)
        state["nirmaan_proposals"] = dedupe_nirmaan_proposals(state["nirmaan_proposals"])
        state["vault"].append(
            {
                "id": new_id("vault"),
                "kind": "nirmaan_proposal",
                "title": proposal["title"],
                "summary": proposal["description"][:300],
                "content": proposal["code_snippet"][:1200],
                "created_at": now_iso(),
            }
        )
        return proposal

    return mutate_state(_mutate)


def dashboard_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    state = ensure_state_shape(state or get_state())
    cabinet = prime_minister_payload(state)
    domains = operational_domains_payload(state)
    today = today_iso()
    hunts_today = sum(1 for row in state["praapti_hunts"] if row.get("created_day") == today)
    pending_proposals = sum(1 for row in state["nirmaan_proposals"] if not row.get("approved"))
    swarm_count = len(state["swarm_missions"])
    active_avatars = state["avatar_state"].get("active", ["RAMA"])
    active_profiles = avatar_profiles(active_avatars)
    revenue = round(1.2 + hunts_today * 0.35 + swarm_count * 0.2 + pending_proposals * 0.12, 2)
    army_active = min(12, 4 + swarm_count + pending_proposals + (1 if hunts_today else 0))
    collective_insights = len(state["vault"]) + len(state["conversations"]) + len(state["praapti_hunts"]) * 2
    packages_active = len(state["packages"])
    latest_swarm = state["swarm_missions"][-1] if state["swarm_missions"] else None
    latest_proposals = list(reversed([row for row in state["nirmaan_proposals"] if not row.get("approved")]))[:4]
    return {
        "metrics": {
            "praapti_hunts_today": hunts_today,
            "projected_revenue_inr": revenue,
            "army_active": army_active,
            "nirmaan_active": pending_proposals,
            "collective_insights": collective_insights,
            "vault_items": len(state["vault"]),
            "swarm_missions": swarm_count,
            "vishnu_protection": state["avatar_state"].get("protection_meter", 0),
            "packages_active": packages_active,
            "prime_minister_loops": len(state["cabinet"].get("mission_log", [])),
            "active_secretaries": cabinet["prime_minister"].get("active_secretaries", 0),
            "domain_systems": len(domains),
        },
        "swarm_agents": build_swarm_agents(state),
        "latest_swarm_report": (latest_swarm or {}).get("report", ""),
        "latest_proposals": latest_proposals,
        "voice": state["voice"],
        "creator_mode": state["meta"].get("creator_mode", CREATOR_MODE_LABEL),
        "identity": state["meta"].get("identity", CORE_IDENTITY),
        "active_avatars": active_profiles,
        "active_avatar_names": active_avatars,
        "active_avatar_card": active_profiles[0] if active_profiles else DASHAVATARA["RAMA"],
        "avatar_history": list(reversed(state["avatar_state"].get("history", [])))[:8],
        "memory_guardian": {
            "status": "eternal preservation",
            "vault_items": len(state["vault"]),
            "conversation_items": len(state["conversations"]),
            "seal": "Matsya + Kurma",
            "state_file": str(STATE_PATH),
        },
        "cabinet": cabinet,
        "domain_systems": domains[:8],
        "provider": active_provider_label(state),
    }


def brain_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    state = ensure_state_shape(state or get_state())
    active_names = state["avatar_state"].get("active", ["RAMA"])
    active_profiles = avatar_profiles(active_names)
    hunts = len(state["praapti_hunts"])
    swarm = len(state["swarm_missions"])
    proposals = sum(1 for row in state["nirmaan_proposals"] if not row.get("approved"))
    vault_items = len(state["vault"])
    conversations = len(state["conversations"])
    packages = len(state["packages"])
    protection_meter = int(state["avatar_state"].get("protection_meter", 0))

    expand = min(99, 44 + hunts * 7 + swarm * 4 + packages * 5 + (9 if "KRISHNA" in active_names else 0) + (8 if "KALKI" in active_names else 0))
    grow = min(99, 42 + hunts * 6 + max(1, conversations // 6) + (7 if "RAMA" in active_names else 0) + (7 if "BUDDHA" in active_names else 0))
    develop = min(99, 38 + proposals * 13 + swarm * 5 + (7 if "VAMANA" in active_names else 0) + (9 if "PARASHURAMA" in active_names else 0))
    protect = min(
        99,
        max(
            protection_meter,
            40 + vault_items // 4 + (10 if "NARASIMHA" in active_names else 0) + (8 if "KURMA" in active_names else 0),
        ),
    )

    pillars = [
        {
            "id": "expand",
            "label": "Expand",
            "score": expand,
            "icon": "brain-expand.svg",
            "summary": "Reach wider territory, attract new opportunities, and widen the empire's field of action.",
        },
        {
            "id": "grow",
            "label": "Grow",
            "score": grow,
            "icon": "brain-grow.svg",
            "summary": "Strengthen momentum, deepen relationships, and convert activity into durable abundance.",
        },
        {
            "id": "develop",
            "label": "Develop",
            "score": develop,
            "icon": "brain-develop.svg",
            "summary": "Invent new capabilities, refine the operating system, and evolve the intelligence stack.",
        },
        {
            "id": "protect",
            "label": "Protect",
            "score": protect,
            "icon": "brain-protect.svg",
            "summary": "Preserve memory, defend the core, and keep the empire stable under pressure.",
        },
    ]

    dominant = max(pillars, key=lambda item: item["score"])
    recommendations: List[Dict[str, str]] = []
    if hunts == 0:
        recommendations.append(
            {
                "title": "Feed Praapti",
                "action": "Start a recruitment hunt to expand the intelligence network and attract fresh momentum.",
                "layer": "praapti",
            }
        )
    if proposals == 0:
        recommendations.append(
            {
                "title": "Trigger Nirmaan",
                "action": "Ask the brain to create a new proposal so development does not stall.",
                "layer": "nirmaan",
            }
        )
    if protect < 78:
        recommendations.append(
            {
                "title": "Raise Protection",
                "action": "Channel Matsya + Kurma and refresh the vault to reinforce long-term preservation.",
                "layer": "avatars",
            }
        )
    if not state["settings"].get("screen_capture_enabled", False):
        recommendations.append(
            {
                "title": "Enable Screen Vision",
                "action": "Turn on consent-based screen vision in Settings when you want Leazy to observe your current workspace.",
                "layer": "settings",
            }
        )
    if not recommendations:
        recommendations.append(
            {
                "title": "Run Brain Pulse",
                "action": "The core is stable. Launch a brain pulse to generate the next deliberate move.",
                "layer": "bridge",
            }
        )
    if not state["cabinet"].get("mission_log"):
        recommendations.insert(
            0,
            {
                "title": "Activate Prime Minister",
                "action": "Start the Prime Minister cabinet so revenue, protection, and delivery keep moving in a live loop.",
                "layer": "cabinet",
            },
        )

    evolution_cycle = [
        {
            "id": "observe",
            "label": "Observe",
            "score": min(99, 34 + conversations * 3 + hunts * 4 + packages * 4),
            "status": "live" if conversations or packages else "idle",
            "summary": "Collect signals from hunts, packages, voice, and direct interaction.",
        },
        {
            "id": "interpret",
            "label": "Interpret",
            "score": min(99, 38 + hunts * 5 + len(active_names) * 7),
            "status": "active" if hunts or len(active_names) > 1 else "warming",
            "summary": "Use avatars and creator guidance to convert raw inputs into direction.",
        },
        {
            "id": "design",
            "label": "Design",
            "score": min(99, 36 + proposals * 11 + swarm * 4),
            "status": "active" if proposals else "ready",
            "summary": "Shape improvements, proposals, and next-step operating plans.",
        },
        {
            "id": "deploy",
            "label": "Deploy",
            "score": min(99, 40 + swarm * 7 + hunts * 4 + packages * 5),
            "status": "live" if swarm or packages else "ready",
            "summary": "Run missions, launch packages, and move the company system forward.",
        },
        {
            "id": "preserve",
            "label": "Preserve",
            "score": protect,
            "status": "sealed" if protect >= 80 else "hardening",
            "summary": "Akshaya stores memory, closes loops, and keeps the empire lightweight for the long run.",
        },
    ]

    return {
        "identity": state["meta"].get("identity", CORE_IDENTITY),
        "creator_mode": state["meta"].get("creator_mode", CREATOR_MODE_LABEL),
        "mode": dominant["label"],
        "dominant_pillar": dominant,
        "temperature": round(sum(item["score"] for item in pillars) / len(pillars), 1),
        "pillars": pillars,
        "active_avatars": active_profiles,
        "recommendations": recommendations[:3],
        "evolution_cycle": evolution_cycle,
        "provider": active_provider_label(state),
        "prime_minister": prime_minister_payload(state),
        "heartbeat": {
            "status": "stable" if protect >= 72 else "hardening",
            "pulse": 60 + len(active_names) * 4 + min(18, swarm * 2),
            "memory_guardian": state["meta"].get("memory_guardian", "eternal preservation"),
            "eternal_mode": "Matsya + Kurma",
        },
    }


def prana_nadi_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    brain = brain_payload(current)
    active_names = current["avatar_state"].get("active", ["RAMA"])
    transmissions = (
        (1 if current.get("voice", {}).get("last_command") else 0)
        + min(6, len(current.get("conversations", [])) // 2)
        + min(4, len(current.get("praapti_hunts", [])))
        + min(4, len(current.get("nirmaan_proposals", [])))
        + min(4, len(current.get("swarm_missions", [])))
    )
    cabinet_cycles = len(current["cabinet"].get("mission_log", []))
    conversations = len(current.get("conversations", []))
    intensity = min(100, 62 + len(active_names) * 6 + cabinet_cycles + transmissions * 4)
    hridaya_glow = min(100, 58 + conversations // 2 + brain["temperature"] // 2)
    return {
        "pulse": "alive",
        "intensity": intensity,
        "message": "Your will flows through every nadi instantly.",
        "channels": [
            {
                "name": "Sushumna Nadi",
                "role": "Central spine connecting directly to the creator's will.",
                "strength": min(100, intensity + 4),
            },
            {
                "name": "Ida Nadi",
                "role": "Intuition, memory, empathy, and reflective intelligence.",
                "strength": min(100, 54 + conversations + len(current.get("vault", [])) // 3),
            },
            {
                "name": "Pingala Nadi",
                "role": "Action, execution, response, and outward force.",
                "strength": min(100, 56 + len(current.get("swarm_missions", [])) * 7 + len(current.get("packages", [])) * 8),
            },
            {
                "name": "Prana-Vahini",
                "role": "Instant reflex, zero hesitation, command relay across the empire.",
                "strength": min(100, 60 + transmissions * 8),
            },
        ],
        "nadis_total": 72000,
        "active_avatar_count": len(active_names),
        "heartbeat_ms": max(900, 2200 - intensity * 8),
        "hridaya": {
            "name": "Sovereign Sync Core",
            "glow": hridaya_glow,
            "sync": min(100, 65 + brain["temperature"] // 3 + cabinet_cycles * 3),
            "mantra": "Sankalpa Aikyam",
            "message": "Where command becomes system-wide motion.",
        },
    }


def settings_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    settings = current["settings"]
    voice = current["voice"]
    providers = provider_config()
    issue = provider_issue_snapshot()
    preferred = settings.get("provider_preference", "openai")
    external_configured = any(providers[name]["configured"] for name in ("anthropic", "openai", "gemini"))
    if preferred == "ollama" and not providers.get("ollama", {}).get("configured"):
        preferred = "built_in"
    if preferred in {"anthropic", "openai", "gemini"} and not providers.get(preferred, {}).get("configured") and not external_configured:
        preferred = "built_in"
    interpreter_status = INTERPRETER_LAYER.get("status_payload", lambda: {})()
    orchestration_status = ORCHESTRATION_STACK_LAYER.get("status_payload", lambda: {})()
    local_ai_status = LOCAL_AI_STACK_LAYER.get("status_payload", lambda viewer=None: {})()
    voice_runtime_status = VOICE_RUNTIME_LAYER.get("status_payload", lambda viewer=None: {})()
    return {
        "provider_preference": preferred,
        "external_ai_mode": settings.get("external_ai_mode", "manual_only"),
        "providers": {
            "ollama": {**providers["ollama"], "label": "Ollama Local"},
            "anthropic": {**providers["anthropic"], "label": "Anthropic"},
            "openai": {**providers["openai"], "label": "OpenAI"},
            "gemini": {**providers["gemini"], "label": "Gemini"},
            "built_in": {**providers["built_in"], "label": "Built-in"},
        },
        "active_provider": active_provider_label(current),
        "interpreter": interpreter_status,
        "orchestration": orchestration_status,
        "local_ai": local_ai_status,
        "voice_runtime": voice_runtime_status,
        "privacy": {
            "screen_capture_enabled": bool(settings.get("screen_capture_enabled", False)),
            "screen_capture_requires_consent": True,
            "audio_capture_enabled": bool(settings.get("audio_capture_enabled", False)),
            "privacy_guard_enabled": bool(settings.get("privacy_guard_enabled", True)),
            "bounded_packages_enabled": bool(settings.get("bounded_packages_enabled", True)),
            "hq_visual_sync": bool(settings.get("hq_visual_sync", True)),
        },
        "voice": {
            "always_listening": bool(voice.get("always_listening", False)),
            "wake_words": voice.get("wake_words", []),
            "voice_profile": voice.get("voice_profile", "sovereign_female"),
            "language": voice.get("language", "en-IN"),
            "rate": float(voice.get("rate", 0.94) or 0.94),
            "pitch": float(voice.get("pitch", 1.08) or 1.08),
            "engine": voice.get("engine", "browser_builtin_female"),
        },
        "provider_issue": issue,
        "provider_issues": provider_issues_snapshot(),
        "notes": {
            "keys_source": "No GPT or Gemini key ships with Ishani. Paste an external key manually only when you want it. Ollama local models do not need any API key.",
            "models_source": "Live model lists can be discovered from the selected provider. Ollama reads installed local models directly from your machine.",
            "external_ai_mode": "External providers are manual-first. Background loops, cabinet cycles, and passive monitoring stay on the local or built-in brain unless a human-triggered action explicitly asks for an external model.",
            "packages_mode": "Packages are bounded research/execution missions. They do not self-propagate or bypass consent.",
        },
    }


def portal_state_payload(viewer: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    state = get_state()
    dashboard = dashboard_payload()
    brain = brain_payload(state)
    cabinet = prime_minister_payload(state)
    monitor = mother_monitor_payload(state)
    nervous_system = nervous_system_payload(state)
    domains = operational_domains_payload(state)
    memory_audit = memory_audit_payload(state)
    hunts = list(reversed(state["praapti_hunts"]))[:10]
    packages = list(reversed(state["packages"]))[:10]
    activity = list(reversed(state["vault"]))[:16]
    candidate_rows: List[Dict[str, Any]] = []
    for hunt in hunts:
        for candidate in hunt.get("candidates", [])[:3]:
            candidate_rows.append(
                {
                    "name": candidate.get("name", "Candidate"),
                    "title": candidate.get("title", "Role"),
                    "fit_score": candidate.get("fit_score", 0),
                    "experience": candidate.get("experience", 0),
                    "client_company": hunt.get("client_company", COMPANY_NAME),
                }
            )
    clients = sorted({hunt.get("client_company", COMPANY_NAME) for hunt in hunts})
    jobs = [
        {
            "client_company": hunt.get("client_company", COMPANY_NAME),
            "summary": hunt.get("job_description", "")[:160],
            "urgency": hunt.get("urgency", "high"),
            "created_at": hunt.get("created_at", ""),
        }
        for hunt in hunts
    ]
    ats_surface = ats_surface_payload(state, candidate_rows[:18], jobs[:12])
    network_surface = network_surface_payload(state, clients[:12], candidate_rows[:18])
    hq_owner = hq_owner_payload(state, clients[:12], jobs[:12])
    reports = []
    cabinet_cycles = list(reversed(state["cabinet"].get("mission_log", [])))[:3]
    for cycle in cabinet_cycles:
        reports.append({"title": "Prime Minister Cycle", "summary": cycle.get("report", "")[:220]})
    if state["swarm_missions"]:
        latest_swarm = state["swarm_missions"][-1]
        reports.append({"title": "Latest Swarm Report", "summary": latest_swarm.get("report", "")[:280]})
    for proposal in list(reversed(state["nirmaan_proposals"]))[:3]:
        reports.append({"title": proposal.get("title", "Nirmaan Proposal"), "summary": proposal.get("description", "")[:220]})
    documents = user_documents(viewer["id"], limit=12) if viewer else []
    accounts = account_status_payload(viewer) if viewer else None
    orders = db_all(
        "SELECT id, plan_id, amount_inr, payment_method, status, created_at FROM orders WHERE user_id=? ORDER BY created_at DESC LIMIT 8",
        (viewer["id"],),
    ) if viewer else []
    plans = []
    for plan in PLAN_CATALOG:
        if viewer is None and plan["role"] != "member":
            continue
        if viewer is not None and plan["role"] not in {"member", viewer.get("role", "member")}:
            continue
        plans.append({**plan, "services": list(plan["services"])})
    return {
        "dashboard": dashboard,
        "brain": brain,
        "cabinet": cabinet,
        "monitor": monitor,
        "nervous_system": nervous_system,
        "domains": domains,
        "memory_audit": memory_audit,
        "settings": settings_payload(state),
        "package_templates": PACKAGE_TEMPLATES,
        "packages": packages,
        "hunts": hunts,
        "candidates": candidate_rows[:18],
        "clients": clients[:12],
        "jobs": jobs[:12],
        "reports": reports[:8],
        "activity": activity,
        "ats_surface": ats_surface,
        "network_surface": network_surface,
        "hq_owner": hq_owner,
        "auth": session_payload(viewer),
        "plans": plans,
        "billing": {"orders": orders},
        "documents": documents,
        "accounts": accounts,
    }


def public_hq_context_brief(state: Optional[Dict[str, Any]] = None) -> str:
    current = ensure_state_shape(state or get_state())
    dashboard = dashboard_payload(current)
    domains = operational_domains_payload(current)[:4]
    plans = [plan for plan in PLAN_CATALOG if plan.get("role") == "member"][:3]
    services = [
        "public TechBuzz AI concierge",
        "hiring and Praapti recruitment",
        "document studio",
        "automation and company execution",
    ]
    domain_lines = "\n".join(
        f"- {domain['name']}: {domain['status']} | {domain['latest_signal']}"
        for domain in domains
    )
    plan_lines = "\n".join(
        f"- {plan['name']}: INR {plan['price_inr']} ({plan['billing_period']})"
        for plan in plans
    )
    return (
        f"Company: {COMPANY_NAME}\n"
        f"Identity: {dashboard['identity']}\n"
        f"Public services: {', '.join(services)}\n"
        f"Live metrics: revenue {dashboard['metrics']['projected_revenue_inr']} Cr, "
        f"hunts today {dashboard['metrics']['praapti_hunts_today']}, "
        f"active secretaries {dashboard['metrics']['active_secretaries']}, "
        f"vault items {dashboard['metrics']['vault_items']}.\n"
        f"Public plans:\n{plan_lines}\n"
        f"Current operating lanes:\n{domain_lines}"
    )


def normalize_label_key(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", str(value or "").strip().lower()).strip("-")


def dedupe_nirmaan_proposals(proposals: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    seen: set[str] = set()
    deduped: List[Dict[str, Any]] = []
    for proposal in sorted(proposals, key=lambda row: row.get("created_at", ""), reverse=True):
        key = normalize_label_key(f"{proposal.get('title', '')}:{proposal.get('reason', '')}:{proposal.get('approved', False)}")
        if key in seen:
            continue
        seen.add(key)
        deduped.append(proposal)
    return list(reversed(deduped))


def candidate_stage(candidate: Dict[str, Any]) -> str:
    score = int(candidate.get("fit_score", 0) or 0)
    if score >= 95:
        return "offer"
    if score >= 92:
        return "interview"
    if score >= 88:
        return "screening"
    return "sourced"


def ats_surface_payload(state: Dict[str, Any], candidates: List[Dict[str, Any]], jobs: List[Dict[str, Any]]) -> Dict[str, Any]:
    stage_labels = {
        "sourced": "Sourced",
        "screening": "Screening",
        "interview": "Interview",
        "offer": "Offer",
    }
    columns = {key: [] for key in stage_labels}
    for candidate in candidates:
        stage = candidate_stage(candidate)
        columns[stage].append(
            {
                **candidate,
                "stage": stage,
                "ai_score": int(candidate.get("fit_score", 0) or 0),
            }
        )
    counts = [row["fit_score"] for row in candidates if isinstance(row.get("fit_score"), (int, float))]
    avg_score = round(sum(counts) / len(counts), 1) if counts else 0
    return {
        "columns": [
            {"id": key, "label": label, "items": columns[key]}
            for key, label in stage_labels.items()
        ],
        "score_summary": {
            "average_fit": avg_score,
            "best_fit": max(counts) if counts else 0,
            "candidate_count": len(candidates),
            "job_count": len(jobs),
        },
        "import_ready": bool(state.get("praapti_hunts")),
    }


def network_surface_payload(state: Dict[str, Any], clients: List[str], candidates: List[Dict[str, Any]]) -> Dict[str, Any]:
    vault = list(reversed(state.get("vault", [])))
    hunts = list(reversed(state.get("praapti_hunts", [])))
    connections: List[Dict[str, Any]] = []
    for client in clients[:8]:
        client_candidates = [row for row in candidates if row.get("client_company") == client][:4]
        connections.append(
            {
                "name": client,
                "kind": "client",
                "strength": min(98, 62 + len(client_candidates) * 6),
                "summary": f"{len(client_candidates)} active candidate connection(s) are flowing through this lane.",
            }
        )
    posts = []
    for item in vault[:8]:
        posts.append(
            {
                "title": item.get("title", "System post"),
                "summary": item.get("summary", ""),
                "kind": item.get("kind", "signal"),
                "created_at": item.get("created_at", ""),
            }
        )
    signals = []
    for hunt in hunts[:6]:
        signals.append(
            {
                "title": hunt.get("client_company", COMPANY_NAME),
                "summary": hunt.get("job_description", "")[:180],
                "source": hunt.get("provider", "built-in"),
                "created_at": hunt.get("created_at", ""),
            }
        )
    if not signals:
        for domain in operational_domains_payload(state)[:6]:
            signals.append(
                {
                    "title": domain["name"],
                    "summary": domain["latest_signal"],
                    "source": domain["lead"],
                    "created_at": now_iso(),
                }
            )
    return {
        "connections": connections,
        "posts": posts,
        "signals": signals,
        "scanner_label": "Carbon Signal Scanner",
    }


def hq_owner_payload(state: Dict[str, Any], clients: List[str], jobs: List[Dict[str, Any]]) -> Dict[str, Any]:
    metrics = dashboard_payload(state)["metrics"]
    client_pipeline = []
    for client in clients[:8]:
        client_jobs = [row for row in jobs if row.get("client_company") == client]
        client_pipeline.append(
            {
                "client": client,
                "open_roles": len(client_jobs),
                "pipeline_value": round(0.35 + len(client_jobs) * 0.42, 2),
                "status": "active" if client_jobs else "warming",
            }
        )
    team = []
    for secretary in prime_minister_payload(state).get("secretaries", [])[:8]:
        team.append(
            {
                "name": secretary.get("name", "Secretary"),
                "lane": secretary.get("lane", "general"),
                "priority": secretary.get("priority", 0),
                "status": secretary.get("status", "ready"),
            }
        )
    strategy = [
        {
            "title": row.get("title", "Strategy move"),
            "action": row.get("action", ""),
            "layer": row.get("layer", "hq"),
        }
        for row in brain_payload(state).get("recommendations", [])[:6]
    ]
    return {
        "client_pipeline": client_pipeline,
        "revenue_tracker": {
            "projected_revenue_inr": metrics["projected_revenue_inr"],
            "active_packages": metrics["packages_active"],
            "active_secretaries": metrics["active_secretaries"],
            "collective_insights": metrics["collective_insights"],
        },
        "team": team,
        "strategy": strategy,
        "architecture_note": "Carbon Protocol links HQ, ATS, Network, Agent, and Akshaya as one bonded layer.",
    }


async def generate_public_agent_text(
    *,
    message: str,
    provider: str,
    api_key: Optional[str],
    model: Optional[str],
    history: Optional[List[Message]] = None,
) -> Dict[str, Any]:
    provider = (provider or "built_in").strip().lower()
    safe_history = history or []
    safe_key = (api_key or "").strip()
    selected_model = (model or "").strip()
    seed_context = recruitment_seed_brief("agent", audience="public", limit=4)
    system = (
        f"You are {AI_NAME}, the public TechBuzz AI agent for {COMPANY_NAME}. "
        "Be warm, sharp, and practical. Help with hiring, documents, automation, business strategy, "
        "and getting started with TechBuzz. If the user needs protected actions, explain the next step clearly."
    )
    prompt = (
        f"{public_hq_context_brief()}\n"
        f"Permanent recruitment core:\n{seed_context}\n"
        f"Conversation memory:\n"
        + "\n".join(f"{item.role}: {item.content}" for item in safe_history[-6:])
        + f"\n\nUser request:\n{message}"
    )

    if provider in {"", "built_in", "built-in", "ollama"}:
        if ollama_provider_ready():
            try:
                local_model = selected_model if provider == "ollama" and selected_model else OLLAMA_MODEL
                local_request = compact_ollama_request(
                    prompt=prompt,
                    system=system,
                    workspace="agent",
                    source="public_agent",
                    max_tokens=180,
                )
                result = await call_ollama(
                    prompt=local_request["prompt"],
                    system=local_request["system"],
                    model=local_model,
                    max_tokens=local_request["max_tokens"],
                    timeout_seconds=60.0,
                )
                text = extract_ollama_text(result).strip()
                if text:
                    clear_provider_issue("ollama")
                    return {"text": text, "provider": f"ollama/{local_model}"}
                raise HTTPException(status_code=503, detail="Ollama local model returned an empty reply.")
            except HTTPException as exc:
                detail = str(exc.detail or "").strip() or "Ollama local model timed out or is unavailable."
                set_provider_issue("ollama", detail)
            except Exception as exc:
                detail = str(exc).strip() or "Ollama local model timed out or returned no usable reply."
                set_provider_issue("ollama", detail)

    try:
        if provider == "anthropic" and safe_key:
            result = await call_anthropic_custom(
                messages=[{"role": "user", "content": prompt}],
                system=system,
                api_key=safe_key,
                model=selected_model or MODEL,
                max_tokens=700,
            )
            return {"text": extract_text(result), "provider": f"anthropic/{selected_model or MODEL}"}
        if provider == "openai" and safe_key:
            result = await call_openai_custom(
                messages=[{"role": "user", "content": prompt}],
                system=system,
                api_key=safe_key,
                model=selected_model or OPENAI_MODEL,
                max_tokens=700,
            )
            return {"text": extract_openai_text(result), "provider": f"openai/{selected_model or OPENAI_MODEL}"}
        if provider == "gemini" and safe_key:
            result = await call_gemini_custom(
                prompt=prompt,
                system=system,
                api_key=safe_key,
                model=selected_model or GEMINI_MODEL,
                max_tokens=700,
            )
            return {"text": extract_gemini_text(result), "provider": f"gemini/{selected_model or GEMINI_MODEL}"}
    except HTTPException as exc:
        fallback = fallback_model_text(f"User: {message}", state=get_state(), workspace="agent", source="public_agent")
        return {
            "text": f"Provider note: {provider} could not answer with the supplied key, so TechBuzz AI is using the local built-in brain right now.\n\n{fallback}",
            "provider": f"built_in/fallback-after-{provider}",
            "error": exc.detail,
        }

    fallback = fallback_model_text(f"User: {message}", state=get_state(), workspace="agent", source="public_agent")
    return {"text": fallback, "provider": "built_in/public-agent"}


def serve_frontend_page(filename: str) -> Response:
    path = FRONTEND_DIR / filename
    if not path.exists():
        raise HTTPException(status_code=404, detail=f"{filename} not found")
    if path.suffix.lower() == ".html":
        html = path.read_text(encoding="utf-8", errors="ignore")
        repair_tag = '<script src="/frontend-assets/ui-auto-repair.js"></script>'
        if repair_tag not in html:
            if "</body>" in html:
                html = html.replace("</body>", f"{repair_tag}\n</body>")
            else:
                html = html + "\n" + repair_tag
        return HTMLResponse(html)
    return FileResponse(path)


def voice_status_payload(state: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    current = ensure_state_shape(state or get_state())
    profile = str(current["voice"].get("voice_profile", "sovereign_female"))
    return {
        "always_listening": current["voice"].get("always_listening", False),
        "wake_words": current["voice"].get("wake_words", []),
        "last_command": sanitize_operator_line(current["voice"].get("last_command", ""), ""),
        "active_avatars": [name.title() for name in current["avatar_state"].get("active", ["RAMA"])],
        "voice_profile": profile,
        "voice_profile_label": VOICE_PROFILE_PRESETS.get(profile, VOICE_PROFILE_PRESETS["sovereign_female"])["label"],
        "language": current["voice"].get("language", "en-IN"),
        "rate": float(current["voice"].get("rate", 0.94) or 0.94),
        "pitch": float(current["voice"].get("pitch", 1.08) or 1.08),
        "engine": current["voice"].get("engine", "browser_builtin_female"),
        "runtime": VOICE_RUNTIME_LAYER.get("status_payload", lambda viewer=None: {})(),
        "profiles": [
            {"id": key, **value}
            for key, value in VOICE_PROFILE_PRESETS.items()
        ],
    }


def brain_stream_snapshot() -> Dict[str, Any]:
    state = get_state()
    active = state["avatar_state"].get("active", ["RAMA"])
    return {
        "timestamp": now_iso(),
        "brain": brain_payload(state),
        "dashboard": dashboard_payload(state),
        "cabinet": prime_minister_payload(state),
        "monitor": mother_monitor_payload(state),
        "prana_nadi": prana_nadi_payload(state),
        "settings": settings_payload(state),
        "voice": voice_status_payload(state),
        "vishnu": {
            "identity": state["meta"].get("identity", CORE_IDENTITY),
            "creator_mode": state["meta"].get("creator_mode", CREATOR_MODE_LABEL),
            "active_avatars": avatar_profiles(active),
            "protection_meter": state["avatar_state"].get("protection_meter", 0),
            "history": list(reversed(state["avatar_state"].get("history", [])))[:12],
        },
    }


async def sse_json_generator(request: Request, snapshot_factory, interval_seconds: float = 4.0):
    last_snapshot = ""
    while True:
        if await request.is_disconnected():
            break
        payload = snapshot_factory()
        encoded = json.dumps(payload, ensure_ascii=False)
        if encoded != last_snapshot:
            yield f"event: snapshot\ndata: {encoded}\n\n"
            last_snapshot = encoded
        else:
            yield f": keepalive {now_iso()}\n\n"
        await asyncio.sleep(interval_seconds)


@app.get("/login")
async def login_page():
    return serve_frontend_page("login.html")


@app.get("/", include_in_schema=False)
async def root():
    return RedirectResponse(url="/company/portal", status_code=307)


@app.get("/api")
async def api_root():
    return {
        "service": "TechBuzz Leazy Jinn API",
        "version": "2.0.0",
        "status": "online",
        "pages": ["/login", "/leazy", "/company/portal", "/agent", "/agent/console", "/navigator", "/network", "/ats"],
        "endpoints": [
            "/api/auth/register",
            "/api/auth/login",
            "/api/auth/master-login",
            "/api/auth/me",
            "/api/billing/plans",
            "/api/navigator/status",
            "/api/documents/list",
            "/api/chat",
            "/api/public/provider-models",
            "/api/public/agent-chat",
            "/api/public/hq-chat",
            "/api/leazy/chat",
            "/api/brain/status",
            "/api/brain/hierarchy",
            "/api/brain/pulse",
            "/api/cabinet/status",
            "/api/cabinet/prime-minister",
            "/api/empire/dashboard",
            "/api/mother/monitor",
            "/api/nervous-system/status",
            "/api/ops/domains",
            "/api/memory/audit",
            "/api/settings/status",
            "/api/providers/configure",
            "/api/packages/templates",
            "/api/praapti/hunt",
            "/api/swarm/mission",
            "/api/nirmaan/develop",
            "/api/vishnu/channel",
        ],
        "provider": active_provider_label(),
    }


@app.post("/api/auth/register")
async def auth_register(req: AuthRegisterRequest, request: Request):
    email = normalize_email(req.email)
    if not email or "@" not in email:
        raise HTTPException(status_code=400, detail="Enter a valid email address.")
    if len(req.password or "") < 8:
        raise HTTPException(status_code=400, detail="Password must be at least 8 characters long.")
    if email == normalize_email(MASTER_ACCOUNT_EMAIL):
        raise HTTPException(status_code=403, detail="Use the master login for the creator account.")
    if db_one("SELECT id FROM users WHERE email=?", (email,)):
        raise HTTPException(status_code=409, detail="An account with that email already exists.")
    plan = plan_by_id(req.plan_id)
    if plan["role"] == "master":
        plan = plan_by_id("starter")
    password = hash_password(req.password)
    user_id = new_id("usr")
    db_exec(
        """
        INSERT INTO users(id,name,email,password_hash,password_salt,role,plan_id,created_at,last_login_at,status)
        VALUES(?,?,?,?,?,?,?,?,?,?)
        """,
        (
            user_id,
            (req.name or "TechBuzz User").strip()[:120],
            email,
            password["hash"],
            password["salt"],
            "member",
            plan["id"],
            now_iso(),
            now_iso(),
            "active",
        ),
    )
    order = create_order_for_user(
        user_id,
        plan["id"],
        payment_method="sandbox" if plan["price_inr"] == 0 else "manual_review",
        status="paid" if plan["price_inr"] == 0 else "pending",
        notes="Created during self-serve registration.",
    )
    token = create_session(user_id, request.headers.get("user-agent", ""))
    response = JSONResponse(
        {
            "message": "Account created. Ishani has connected your workspace to the main brain.",
            "auth": session_payload(db_one("SELECT * FROM users WHERE id=?", (user_id,))),
            "order": order,
        }
    )
    set_auth_cookie(response, token)
    return response


@app.post("/api/auth/login")
async def auth_login(req: AuthLoginRequest, request: Request):
    email = normalize_email(req.email)
    user = db_one("SELECT * FROM users WHERE email=?", (email,))
    if not user or user.get("role") == "master":
        raise HTTPException(status_code=401, detail="Invalid email or password.")
    if not verify_password(req.password, user.get("password_salt", ""), user.get("password_hash", "")):
        raise HTTPException(status_code=401, detail="Invalid email or password.")
    db_exec("UPDATE users SET last_login_at=? WHERE id=?", (now_iso(), user["id"]))
    token = create_session(user["id"], request.headers.get("user-agent", ""))
    refreshed = db_one("SELECT * FROM users WHERE id=?", (user["id"],))
    response = JSONResponse({"message": "Login successful.", "auth": session_payload(refreshed)})
    set_auth_cookie(response, token)
    return response


@app.post("/api/auth/master-login")
async def auth_master_login(req: MasterLoginRequest, request: Request):
    identifier = (req.identifier or req.email or "").strip()
    password = req.password or req.master_key or ""
    if not matches_master_identity(identifier):
        raise HTTPException(status_code=401, detail="Master identity not recognized.")
    if not verify_master_password(password):
        raise HTTPException(status_code=401, detail="Master password mismatch.")
    user = db_one("SELECT * FROM users WHERE email=?", (normalize_email(MASTER_ACCOUNT_EMAIL),))
    if not user:
        raise HTTPException(status_code=500, detail="Master account is not initialized.")
    db_exec(
        "UPDATE users SET name=?, role='master', plan_id='mother-core', last_login_at=? WHERE id=?",
        (MASTER_LOGIN_ID, now_iso(), user["id"]),
    )
    token = create_session(user["id"], request.headers.get("user-agent", ""))
    refreshed = db_one("SELECT * FROM users WHERE id=?", (user["id"],))
    response = JSONResponse({"message": "Master access granted.", "auth": session_payload(refreshed)})
    set_auth_cookie(response, token)
    return response


@app.post("/api/auth/logout")
async def auth_logout(request: Request):
    token = request.cookies.get(SESSION_COOKIE_NAME) or ""
    if token:
        db_exec("DELETE FROM sessions WHERE token_hash=?", (hash_session_token(token),))
    response = JSONResponse({"message": "Logged out from Ishani Core."})
    clear_auth_cookie(response)
    return response


@app.get("/api/auth/me")
async def auth_me(request: Request):
    user = session_user(request)
    return session_payload(user)


@app.get("/api/billing/plans")
async def billing_plans(request: Request):
    viewer = session_user(request)
    plans = [
        {**plan, "services": list(plan["services"])}
        for plan in PLAN_CATALOG
        if ((viewer and viewer.get("role") == "master") or plan["role"] != "master")
    ]
    return {"plans": plans, "auth": session_payload(viewer)}


@app.post("/api/billing/checkout")
async def billing_checkout(req: BillingCheckoutRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before checkout.")
    plan = plan_by_id(req.plan_id)
    if plan["role"] == "master" and viewer.get("role") != "master":
        raise HTTPException(status_code=403, detail="This plan is reserved for the creator.")
    payment_method = (req.payment_method or "sandbox").strip().lower()
    order = create_order_for_user(
        viewer["id"],
        plan["id"],
        payment_method=payment_method,
        status="paid" if payment_method == "sandbox" or plan["price_inr"] == 0 else "pending",
        notes=req.notes,
    )
    if order["status"] == "paid":
        db_exec("UPDATE users SET plan_id=? WHERE id=?", (plan["id"], viewer["id"]))
    return {
        "message": "Checkout recorded in Ishani Core.",
        "order": order,
        "plan": plan,
        "next_step": "Plan activated immediately." if order["status"] == "paid" else "Awaiting payment confirmation.",
    }


@app.get("/api/billing/orders")
async def billing_orders(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before viewing orders.")
    rows = db_all(
        "SELECT id, plan_id, amount_inr, payment_method, status, notes, created_at FROM orders WHERE user_id=? ORDER BY created_at DESC LIMIT 20",
        (viewer["id"],),
    )
    return {"orders": rows}


@app.get("/api/navigator/status")
async def navigator_status(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before opening Navigator.")
    return navigator_status_payload(viewer)


@app.post("/api/navigator/capture")
async def navigator_capture(req: NavigatorCaptureRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before saving navigator learning.")
    if not (req.title or "").strip():
        raise HTTPException(status_code=400, detail="A title is required for navigator learning.")
    session = save_navigator_session(
        viewer,
        title=req.title,
        target_url=req.url,
        notes=req.notes,
        tags=req.tags,
        mode=req.mode,
        source_hint="navigator_ui",
    )
    return {
        "message": "Navigator learning saved and broadcast to the brain mesh.",
        "session": session,
        "status": navigator_status_payload(viewer),
    }


@app.post("/api/navigator/open-launcher")
async def navigator_open_launcher(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before launching local sourcing tools.")
    if not NAUKRI_LAUNCHER_PATH.exists():
        raise HTTPException(status_code=404, detail="Configured Naukri launcher file was not found.")
    try:
        os.startfile(str(NAUKRI_LAUNCHER_PATH))  # type: ignore[attr-defined]
    except AttributeError:
        raise HTTPException(status_code=400, detail="Local launcher is only supported on Windows hosts.")
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Unable to open the launcher: {exc}")
    session = save_navigator_session(
        viewer,
        title="Naukri Launcher Opened",
        target_url="naukri://launcher",
        notes=f"Operator explicitly opened the local launcher at {NAUKRI_LAUNCHER_PATH}.",
        tags=["naukri", "launcher", "sourcing"],
        mode="launcher",
        source_hint="local_launcher",
    )
    return {
        "message": "Naukri launcher opened locally.",
        "session": session,
        "status": navigator_status_payload(viewer),
    }


@app.get("/api/accounts/status")
async def accounts_status(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before opening Accounts Command.")
    return account_status_payload(viewer)


@app.get("/api/accounts/ledger")
async def accounts_ledger(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before viewing the ledger.")
    payload = account_status_payload(viewer)
    return {
        "ledger": payload.get("ledger", []),
        "summary": payload.get("financials", {}),
        "profile": payload.get("profile", {}),
    }


@app.post("/api/accounts/ledger")
async def accounts_ledger_entry(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before posting ledger entries.")
    body = await request.json()
    amount = round(max(0.0, float(body.get("amount") or 0.0)), 2)
    description = (body.get("description") or "").strip()
    if amount <= 0:
        raise HTTPException(status_code=400, detail="Amount must be greater than zero.")
    if not description:
        raise HTTPException(status_code=400, detail="Description is required.")
    profile = ensure_account_profile(viewer["id"])
    entry_type = normalize_account_entry_type(body.get("entry_type") or "")
    if entry_type not in ACCOUNT_INFLOW_TYPES | ACCOUNT_OUTFLOW_TYPES:
        raise HTTPException(status_code=400, detail="Unsupported entry type.")
    category = infer_account_category(entry_type, description)
    gst_rate = account_gst_rate_for_category(body.get("category") or category, float(profile.get("default_tax_rate", 18.0) or 18.0))
    gst = calculate_gst_breakdown(amount, gst_rate, body.get("supply_type") or "intra")
    tds = detect_tds_obligation(entry_type, description, amount, body.get("section") or "", body.get("party_pan") or "")
    occurred_on = (body.get("date") or body.get("occurred_on") or today_iso()).strip()[:32] or today_iso()
    db_exec(
        """
        INSERT INTO account_entries(id,user_id,entry_type,category,amount,tax_percent,tax_amount,total_amount,currency,counterparty,description,source,occurred_on,created_at)
        VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?)
        """,
        (
            new_id("acct"),
            viewer["id"],
            entry_type,
            category[:120],
            amount,
            gst["rate"],
            gst["tax_total"],
            gst["invoice_total"],
            (body.get("currency") or profile.get("currency") or "INR").strip().upper()[:12],
            (body.get("party_name") or body.get("counterparty") or "").strip()[:160],
            description[:600],
            "accounts_command",
            occurred_on,
            now_iso(),
        ),
    )
    return {
        "message": "Ledger entry recorded and tax hints generated.",
        "gst": gst,
        "tds": tds,
        "accounts": account_status_payload(viewer),
    }


@app.post("/api/accounts/invoice")
async def accounts_invoice(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before generating invoices.")
    body = await request.json()
    client_name = (body.get("client_name") or "").strip()
    description = (body.get("description") or "Recruitment consultancy services").strip()
    amount = round(max(0.0, float(body.get("amount") or 0.0)), 2)
    due_days = max(1, min(180, int(body.get("due_days") or 30)))
    if not client_name:
        raise HTTPException(status_code=400, detail="Client name is required.")
    if amount <= 0:
        raise HTTPException(status_code=400, detail="Invoice amount must be greater than zero.")
    profile = ensure_account_profile(viewer["id"])
    client_gstin = (body.get("client_gstin") or "").strip()
    supply_type = "inter" if client_gstin and not client_gstin.startswith("09") else "intra"
    gst = calculate_gst_breakdown(
        amount,
        account_gst_rate_for_category("recruitment_fee", float(profile.get("default_tax_rate", 18.0) or 18.0)),
        supply_type,
    )
    invoice_no = f"TB-{datetime.now(UTC).strftime('%Y%m%d')}-{secrets.randbelow(9000) + 1000}"
    due_date = (datetime.now(UTC) + timedelta(days=due_days)).date().isoformat()
    payload = {
        "kind": "invoice",
        "invoice_no": invoice_no,
        "client_name": client_name,
        "client_gstin": client_gstin,
        "description": description,
        "amount": amount,
        "currency": profile.get("currency") or "INR",
        "due_days": due_days,
        "due_date": due_date,
        "payment_status": "pending",
        "gst": gst,
    }
    create_account_report(
        viewer["id"],
        f"Invoice {invoice_no}",
        f"Invoice {invoice_no} prepared for {client_name} with total {gst['invoice_total']:.2f} {(profile.get('currency') or 'INR')}.",
        payload,
    )
    return {
        "invoice_no": invoice_no,
        "cgst": gst["cgst"],
        "sgst": gst["sgst"],
        "igst": gst["igst"],
        "gst": gst,
        "total": gst["invoice_total"],
        "due_date": due_date,
        "notes": "Invoice prepared and tracked as pending inside Accounts Command until marked paid.",
        "accounts": account_status_payload(viewer),
    }


@app.post("/api/accounts/gst-calc")
async def accounts_gst_calc(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before running GST calculations.")
    body = await request.json()
    amount = round(max(0.0, float(body.get("amount") or 0.0)), 2)
    if amount <= 0:
        raise HTTPException(status_code=400, detail="Amount must be greater than zero.")
    category = (body.get("category") or "recruitment_fee").strip().lower()
    rate = account_gst_rate_for_category(category, 18.0)
    gst = calculate_gst_breakdown(amount, rate, body.get("supply_type") or "intra")
    return {
        "rate_applied": rate,
        "cgst": gst["cgst"],
        "sgst": gst["sgst"],
        "igst": gst["igst"],
        "tax_total": gst["tax_total"],
        "invoice_total": gst["invoice_total"],
    }


@app.post("/api/accounts/tds-calc")
async def accounts_tds_calc(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before running TDS calculations.")
    body = await request.json()
    amount = round(max(0.0, float(body.get("amount") or 0.0)), 2)
    if amount <= 0:
        raise HTTPException(status_code=400, detail="Amount must be greater than zero.")
    return detect_tds_obligation(
        "expense",
        f"{body.get('party_name') or ''} {body.get('section') or ''}",
        amount,
        body.get("section") or "",
        body.get("party_pan") or "",
    )


@app.get("/api/accounts/tax-calendar")
async def accounts_tax_calendar(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before loading the tax calendar.")
    profile = ensure_account_profile(viewer["id"])
    return {"calendar": account_tax_calendar(profile), "profile": profile}


@app.post("/api/accounts/profile")
async def accounts_profile_update(req: AccountsProfileRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before updating the accounts profile.")
    current = ensure_account_profile(viewer["id"])
    region_hint = req.region_code or req.regional_profile or current.get("region_code")
    if not region_hint and (req.state_code or "").strip():
        region_hint = "IN"
    preset = account_region_preset(region_hint)
    region_code = normalize_region_code(region_hint or current.get("region_code"))
    note_parts = [part.strip() for part in [req.notes or "", str(current.get("notes") or ""), str(preset["notes"] or "")] if part and part.strip()]
    if (req.city or "").strip():
        note_parts.append(f"Operating city: {req.city.strip()[:80]}")
    if (req.pan or "").strip():
        note_parts.append(f"PAN: {req.pan.strip()[:20]}")
    record = {
        "user_id": viewer["id"],
        "region_code": region_code,
        "business_name": (req.business_name or current.get("business_name") or f"{COMPANY_NAME} Pvt Ltd").strip()[:160],
        "business_type": (req.business_type or current.get("business_type") or "services").strip()[:120],
        "currency": (req.currency or preset["currency"] or current.get("currency") or "INR").strip().upper()[:12],
        "tax_name": (req.tax_name or preset["tax_name"] or current.get("tax_name") or "Tax").strip()[:60],
        "default_tax_rate": float(req.default_tax_rate if req.default_tax_rate is not None else current.get("default_tax_rate", preset["suggested_tax_rate"]) or 0),
        "filing_cycle": (req.filing_cycle or current.get("filing_cycle") or preset["filing_cycle"]).strip()[:80],
        "tax_registration": (req.tax_registration or req.gstin or current.get("tax_registration") or "").strip()[:120],
        "notes": " | ".join(note_parts)[:1200],
        "updated_at": now_iso(),
    }
    db_exec(
        """
        INSERT INTO account_profiles(user_id,region_code,business_name,business_type,currency,tax_name,default_tax_rate,filing_cycle,tax_registration,notes,updated_at)
        VALUES(?,?,?,?,?,?,?,?,?,?,?)
        ON CONFLICT(user_id) DO UPDATE SET
            region_code=excluded.region_code,
            business_name=excluded.business_name,
            business_type=excluded.business_type,
            currency=excluded.currency,
            tax_name=excluded.tax_name,
            default_tax_rate=excluded.default_tax_rate,
            filing_cycle=excluded.filing_cycle,
            tax_registration=excluded.tax_registration,
            notes=excluded.notes,
            updated_at=excluded.updated_at
        """,
        (
            record["user_id"],
            record["region_code"],
            record["business_name"],
            record["business_type"],
            record["currency"],
            record["tax_name"],
            record["default_tax_rate"],
            record["filing_cycle"],
            record["tax_registration"],
            record["notes"],
            record["updated_at"],
        ),
    )
    profile = ensure_account_profile(viewer["id"])
    akshaya_save(
        "accounts_profile",
        "Accounts profile updated",
        f"{profile['business_name']} now uses the {account_region_preset(profile['region_code'])['label']} preset with {profile['tax_name']} as the active tax lane.",
        {"user_id": viewer["id"], "region_code": profile["region_code"], "tax_name": profile["tax_name"]},
    )
    return {"message": "Accounts profile synchronized with the mother brain.", "accounts": account_status_payload(viewer)}


@app.post("/api/accounts/entry")
async def accounts_entry_create(req: AccountsEntryRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before adding ledger entries.")
    profile = ensure_account_profile(viewer["id"])
    entry_type = (req.entry_type or "").strip().lower()
    if entry_type not in ACCOUNT_INFLOW_TYPES | ACCOUNT_OUTFLOW_TYPES:
        raise HTTPException(status_code=400, detail="Unsupported entry type.")
    category = (req.category or "").strip()
    if not category:
        raise HTTPException(status_code=400, detail="Category is required.")
    amount = round(abs(float(req.amount or 0)), 2)
    if amount <= 0:
        raise HTTPException(status_code=400, detail="Amount must be greater than zero.")
    tax_percent = float(req.tax_percent if req.tax_percent is not None else profile.get("default_tax_rate", 0) or 0)
    tax_percent = max(0.0, min(100.0, tax_percent))
    tax_amount = round(amount * tax_percent / 100.0, 2)
    total_amount = round(amount + tax_amount, 2)
    occurred_on = (req.occurred_on or today_iso()).strip()[:32] or today_iso()
    db_exec(
        """
        INSERT INTO account_entries(id,user_id,entry_type,category,amount,tax_percent,tax_amount,total_amount,currency,counterparty,description,source,occurred_on,created_at)
        VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?)
        """,
        (
            new_id("acct"),
            viewer["id"],
            entry_type,
            category[:120],
            amount,
            tax_percent,
            tax_amount,
            total_amount,
            (req.currency or profile.get("currency") or "INR").strip().upper()[:12],
            (req.counterparty or "").strip()[:160],
            (req.description or "").strip()[:600],
            (req.source or "manual").strip()[:40],
            occurred_on,
            now_iso(),
        ),
    )
    return {"message": "Ledger entry recorded inside Accounts Command.", "accounts": account_status_payload(viewer)}


@app.post("/api/accounts/analyze")
async def accounts_analyze(req: AccountsAnalyzeRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before running account analysis.")
    payload = account_status_payload(viewer)
    profile = payload["profile"]
    metrics = payload["metrics"]
    guidance = payload["guidance"]
    operator_question = (req.question or "").strip()
    focus = (operator_question or req.focus or "compliance_and_cashflow").strip().replace("_", " ")
    summary = (
        f"Accounts Command analyzed {metrics['entries_count']} entries for {profile['business_name']}. "
        f"Net cashflow is {metrics['net_cashflow']:.2f} {profile['currency']}. "
        f"Estimated {profile['tax_name']} reserve is {metrics['estimated_tax_payable']:.2f} {profile['currency']}. "
        f"Top guidance: {guidance[0]['title']}."
    )
    report_payload = {
        "focus": focus,
        "profile": profile,
        "metrics": metrics,
        "guidance": guidance,
        "automation": payload["automation"],
    }
    report = create_account_report(viewer["id"], f"Accounts Command - {focus.title()}", summary, report_payload)
    akshaya_save(
        "accounts_analysis",
        report["title"],
        summary,
        {"user_id": viewer["id"], "focus": focus, "report_id": report["id"]},
    )

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        monitoring = state.setdefault("monitoring", {})
        engine_reports = monitoring.setdefault("engine_reports", [])
        engine_reports.append(
            {
                "engine": "Accounts Command",
                "status": "learning",
                "summary": summary[:280],
                "created_at": now_iso(),
            }
        )
        monitoring["last_scan_at"] = now_iso()
        return state

    mutate_state(_mutate)
    analysis_lines = [
        f"Accounts review for {profile['business_name']}",
        f"Focus: {focus.title()}",
        f"Income: {metrics['income_total']:.2f} {profile['currency']}",
        f"Expenses: {metrics['expense_total']:.2f} {profile['currency']}",
        f"Net cashflow: {metrics['net_cashflow']:.2f} {profile['currency']}",
        f"Estimated {profile['tax_name']} reserve: {metrics['estimated_tax_payable']:.2f} {profile['currency']}",
        f"Entries tracked: {metrics['entries_count']}",
        f"Pending invoices: {payload.get('pending_invoices', 0)}",
        "",
        "Priority guidance:",
    ]
    for item in guidance[:4]:
        analysis_lines.append(f"- {item['title']}: {item['summary']}")
    return {
        "message": "Accounts analysis complete. Guidance has been written back into the mother brain.",
        "analysis": "\n".join(analysis_lines).strip(),
        "report": report,
        "accounts": account_status_payload(viewer),
    }


@app.get("/api/documents/list")
async def documents_list(request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before viewing documents.")
    return {
        "documents": user_documents(viewer["id"], limit=40),
        "capabilities": {
            "pdf_merge": pdf_support_available(),
            "pdf_split": pdf_support_available(),
            "docx_read": DocxDocument is not None,
            "xlsx_read": load_workbook is not None,
            "pptx_read": Presentation is not None,
        },
    }


@app.post("/api/documents/upload")
async def documents_upload(request: Request, files: List[UploadFile] = File(...)):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before uploading documents.")
    created = [store_uploaded_file(upload, viewer["id"]) for upload in files[:10]]
    akshaya_save(
        "document_upload",
        "Document Studio Upload",
        f"{len(created)} document(s) uploaded into the Ishani document studio.",
        {"user_id": viewer["id"], "documents": [row["id"] for row in created]},
    )
    return {"documents": created, "message": f"{len(created)} document(s) uploaded successfully."}


@app.post("/api/documents/create-note")
async def documents_create_note(request: Request, title: str = Form(...), content: str = Form(""), format: str = Form("docx")):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before creating documents.")
    extension = ".docx" if str(format).lower() == "docx" else ".txt"
    record = create_text_document(viewer["id"], title, content, extension=extension)
    return {"document": record, "message": "Document created in the Ishani office studio."}


@app.post("/api/documents/extract")
async def documents_extract(req: DocumentActionRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before reading documents.")
    document = get_document_for_user(viewer["id"], req.document_id)
    return {
        "document": {key: document[key] for key in ("id", "original_name", "mime_type", "summary", "kind", "created_at", "updated_at")},
        "text": document.get("extracted_text", ""),
    }


@app.post("/api/documents/update-text")
async def documents_update_text(req: DocumentActionRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before editing documents.")
    document = get_document_for_user(viewer["id"], req.document_id)
    content = req.content or ""
    target = Path(document["storage_path"])
    if document["extension"] == ".docx" and DocxDocument:
        doc = DocxDocument()
        for paragraph in content.split("\n\n"):
            doc.add_paragraph(paragraph)
        doc.save(str(target))
    else:
        target.write_text(content, encoding="utf-8")
    db_exec(
        "UPDATE documents SET extracted_text=?, summary=?, updated_at=? WHERE id=?",
        (content[:50000], document_summary(content, document["original_name"]), now_iso(), document["id"]),
    )
    return {"message": "Document updated.", "document": get_document_for_user(viewer["id"], document["id"])}


@app.post("/api/documents/merge-pdf")
async def documents_merge_pdf(req: DocumentMergeRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before merging PDFs.")
    document = merge_pdf_documents(viewer["id"], req.document_ids)
    return {"message": "PDFs merged successfully.", "document": document}


@app.post("/api/documents/split-pdf")
async def documents_split_pdf(req: DocumentActionRequest, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before splitting PDFs.")
    document = split_pdf_document(viewer["id"], req.document_id, req.start_page, req.end_page)
    return {"message": "PDF pages extracted successfully.", "document": document}


@app.get("/api/documents/download/{document_id}")
async def documents_download(document_id: str, request: Request):
    viewer = session_user(request)
    if not viewer:
        raise HTTPException(status_code=401, detail="Login required before downloading documents.")
    document = get_document_for_user(viewer["id"], document_id)
    path = Path(document["storage_path"])
    if not path.exists():
        raise HTTPException(status_code=404, detail="Document file is missing from storage.")
    return FileResponse(path, filename=document["original_name"], media_type=document["mime_type"])


@app.get("/favicon.ico", include_in_schema=False)
async def favicon():
    path = FRONTEND_DIR / "leazy-icon.svg"
    if path.exists():
        return FileResponse(path, media_type="image/svg+xml")
    return Response(status_code=204)


@app.get("/leazy")
async def leazy_page():
    return serve_frontend_page("leazy.html")


@app.get("/company/portal")
async def company_portal():
    return serve_frontend_page("company-portal.html")


@app.get("/company/portal.html", include_in_schema=False)
async def company_portal_html():
    return RedirectResponse(url="/company/portal", status_code=307)


@app.get("/agent")
async def agent_page():
    return serve_frontend_page("public-agent.html")


@app.get("/agent/console")
async def agent_console_page():
    return serve_frontend_page("agent.html")


@app.get("/navigator")
async def navigator_page():
    return serve_frontend_page("navigator.html")


@app.get("/browser", include_in_schema=False)
async def browser_page():
    return serve_frontend_page("browser.html")


@app.get("/career")
async def career_page():
    return serve_frontend_page("career.html")


@app.get("/jobs")
async def jobs_page():
    return serve_frontend_page("jobs.html")


@app.get("/jobs/{job_id}", include_in_schema=False)
async def jobs_detail_page(job_id: str):
    return RedirectResponse(url=f"/jobs?job={job_id}", status_code=307)


@app.get("/ide")
async def ide_page():
    return serve_frontend_page("ide.html")


@app.get("/mission")
async def mission_page():
    return serve_frontend_page("mission.html")


@app.get("/neural")
async def neural_page():
    return serve_frontend_page("neural.html")


@app.get("/photon")
async def photon_page():
    return serve_frontend_page("photon.html")


@app.get("/research")
async def research_page():
    return serve_frontend_page("research.html")


@app.get("/spread")
async def spread_page():
    return serve_frontend_page("spread.html")


@app.get("/company-portal", include_in_schema=False)
async def company_portal_page():
    return RedirectResponse(url="/company/workspace", status_code=307)


@app.get("/company/workspace", include_in_schema=False)
async def company_workspace_page():
    return serve_frontend_page("company-portal.html")


@app.get("/office", include_in_schema=False)
async def office_page():
    return RedirectResponse(url="/media", status_code=307)


@app.get("/index.html", include_in_schema=False)
async def index_page():
    return RedirectResponse(url="/company/portal", status_code=307)


@app.get("/company-register", include_in_schema=False)
async def company_register_page():
    return serve_frontend_page("index.html")


@app.get("/pricing", include_in_schema=False)
async def pricing_page_alias():
    return serve_frontend_page("index.html")


@app.get("/about", include_in_schema=False)
async def about_page_alias():
    return serve_frontend_page("index.html")


@app.get("/contact", include_in_schema=False)
async def contact_page_alias():
    return serve_frontend_page("index.html")


@app.get("/blog", include_in_schema=False)
async def blog_page_alias():
    return serve_frontend_page("index.html")


@app.get("/privacy", include_in_schema=False)
async def privacy_page_alias():
    return serve_frontend_page("index.html")


@app.get("/terms", include_in_schema=False)
async def terms_page_alias():
    return serve_frontend_page("index.html")


@app.get("/api-docs", include_in_schema=False)
async def api_docs_page():
    return RedirectResponse(url="/api", status_code=307)


@app.get("/hq")
async def hq_owner_page():
    return serve_frontend_page("hq-owner.html")


@app.get("/company/owner", include_in_schema=False)
async def company_owner_page():
    return RedirectResponse(url="/hq", status_code=307)


@app.get("/network")
async def network_page():
    return serve_frontend_page("network.html")


@app.get("/network/intel")
async def network_intel_page():
    return serve_frontend_page("network-intel.html")


@app.get("/company/network.html", include_in_schema=False)
async def company_network_html():
    return RedirectResponse(url="/network/intel", status_code=307)


@app.get("/ats")
async def ats_page():
    return serve_frontend_page("ats.html")


@app.get("/recruiter-mode")
async def recruiter_mode_page():
    return serve_frontend_page("recruiter-mode.html")


@app.get("/company/ats.html", include_in_schema=False)
async def company_ats_html():
    return RedirectResponse(url="/ats", status_code=307)


@app.get("/company/hq.html", include_in_schema=False)
async def company_hq_html():
    return RedirectResponse(url="/company/portal", status_code=307)


@app.get("/cdn-cgi/scripts/5c5dd728/cloudflare-static/email-decode.min.js", include_in_schema=False)
async def cloudflare_email_decode_stub():
    return Response("/* local compatibility shim */", media_type="application/javascript")


@app.get("/manifest.json")
async def manifest():
    path = FRONTEND_DIR / "manifest.json"
    if path.exists():
        return FileResponse(path, media_type="application/manifest+json")
    return JSONResponse(
        {
            "name": f"{AI_NAME} - The King's Empire",
            "short_name": AI_NAME,
            "start_url": "/leazy",
            "display": "standalone",
            "background_color": "#0a0f1c",
            "theme_color": "#facc15",
            "icons": [{"src": "/frontend-assets/leazy-icon.svg", "sizes": "any", "type": "image/svg+xml"}],
        }
    )


@app.get("/service-worker.js")
async def service_worker():
    path = FRONTEND_DIR / "service-worker.js"
    if path.exists():
        return FileResponse(path, media_type="application/javascript")
    return Response(
        "self.addEventListener('install',()=>self.skipWaiting());"
        "self.addEventListener('activate',e=>e.waitUntil(self.clients.claim()));",
        media_type="application/javascript",
    )


@app.get("/api/health")
async def health():
    return {
        "status": "ok",
        "timestamp": now_iso(),
        "api_key_set": any(
            (
                bool(ANTHROPIC_API_KEY),
                bool(OPENAI_API_KEY),
                bool(GEMINI_API_KEY),
            )
        ),
        "model": active_provider_label(),
        "brand": AI_NAME,
    }


@app.get("/api/health/full")
async def health_full():
    data = dashboard_payload()
    return {
        "status": "ok",
        "timestamp": now_iso(),
        "brand": AI_NAME,
        "identity": CORE_IDENTITY,
        "brain": "Leazy Empire Core",
        "ai_stack": {
            "anthropic": bool(ANTHROPIC_API_KEY),
            "openai": bool(OPENAI_API_KEY),
            "gemini": bool(GEMINI_API_KEY),
            "builtin": True,
            "active": active_provider_label(),
            "label": active_provider_label(),
        },
        "creator_mode": data["creator_mode"],
        "active_avatars": data["active_avatars"],
        "memory_guardian": data["memory_guardian"],
        "metrics": data["metrics"],
    }


@app.post("/api/admin/login")
async def admin_login(req: AdminLoginRequest, request: Request):
    identifier = (req.identifier or req.email or "").strip()
    if not matches_master_identity(identifier) or not verify_master_password(req.password):
        log.warning("Failed admin login attempt for: %s", identifier)
        raise HTTPException(status_code=401, detail="Invalid credentials")
    user = db_one("SELECT * FROM users WHERE email=?", (normalize_email(MASTER_ACCOUNT_EMAIL),))
    if not user:
        raise HTTPException(status_code=500, detail="Master account is not initialized.")
    token = create_session(user["id"], request.headers.get("user-agent", ""))
    log.info("Admin login successful: %s", identifier)
    return {
        "success": True,
        "token": token,
        "email": "",
        "login_id": MASTER_LOGIN_ID,
        "expires_in": "session",
    }


@app.post("/api/chat")
async def chat(req: ChatRequest, request: Request):
    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Please wait a moment.")

    api_messages = [{"role": m.role, "content": m.content} for m in req.messages]
    system = req.system or (ADMIN_SYSTEM if req.admin_mode else PUBLIC_SYSTEM)
    log.info("Chat request | IP=%s | admin=%s | msgs=%s | search=%s", ip, req.admin_mode, len(api_messages), req.web_search)

    if ANTHROPIC_API_KEY:
        try:
            result = await call_anthropic(
                messages=api_messages,
                system=system,
                max_tokens=req.max_tokens,
                use_web_search=req.web_search,
            )
            text = extract_text(result)
            used_search = any(
                block.get("type") == "tool_use" and block.get("name") == "web_search"
                for block in result.get("content", [])
            )
            return {
                "response": text,
                "used_web_search": used_search,
                "usage": result.get("usage", {}),
                "model": result.get("model", MODEL),
                "stop_reason": result.get("stop_reason"),
            }
        except HTTPException:
            raise
        except Exception as exc:
            log.error("Chat error: %s", exc)
            raise HTTPException(status_code=500, detail=str(exc))

    fallback_prompt = "\n".join(f"{msg['role']}: {msg['content']}" for msg in api_messages[-6:])
    generated = await generate_text(fallback_prompt, system=system, max_tokens=req.max_tokens, use_web_search=False, source="manual")
    return {
        "response": generated["text"],
        "used_web_search": False,
        "usage": generated["usage"],
        "model": generated["provider"],
        "stop_reason": "fallback",
    }


@app.post("/api/public/hq-chat")
async def public_hq_chat(req: PublicHqChatRequest, request: Request):
    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Please wait a moment.")

    message = (req.message or "").strip()
    if not message:
        raise HTTPException(status_code=400, detail="Message is required.")

    context_line = (req.context or "").strip()[:160]
    seed_context = recruitment_seed_brief("hq", audience="public", limit=4)
    prompt = (
        f"{public_hq_context_brief()}\n"
        f"Permanent recruitment core:\n{seed_context}\n"
        f"Visitor context: {context_line or 'general public HQ visitor'}\n\n"
        f"Visitor request:\n{message}"
    )
    generated = await generate_text(
        prompt,
        system=PUBLIC_HQ_SYSTEM,
        max_tokens=180,
        use_web_search=False,
        workspace="hq",
        source="public_hq",
    )
    return {
        "reply": generated["text"],
        "provider": generated["provider"],
        "quick_links": [
            {"label": "Free AI", "href": "/company/portal#hqConcierge"},
            {"label": "Login", "href": "/login"},
            {"label": "TechBuzz Agent", "href": "/agent"},
        ],
    }


@app.post("/api/public/provider-models")
async def public_provider_models(req: ProviderCatalogRequest):
    provider = (req.provider or "").strip().lower()
    api_key = (req.api_key or "").strip() or None
    catalog = await discover_provider_models(provider, api_key=api_key)
    return {
        "provider": provider,
        "models": catalog["models"],
        "current_model": catalog["current_model"],
        "source": catalog["source"],
    }


@app.post("/api/public/agent-chat")
async def public_agent_chat(req: PublicAgentChatRequest, request: Request):
    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Please wait a moment.")
    message = (req.message or "").strip()
    if not message:
        raise HTTPException(status_code=400, detail="Message is required.")
    generated = await generate_public_agent_text(
        message=message,
        provider=req.provider,
        api_key=req.api_key,
        model=req.model,
        history=req.history,
    )
    return {
        "reply": generated["text"],
        "provider": generated["provider"],
        "quick_links": [
            {"label": "TechBuzz HQ", "href": "/company/portal"},
            {"label": "Agent Console", "href": "/login?next=/agent/console"},
            {"label": "Protected Core", "href": "/login?next=/leazy"},
        ],
        "error": generated.get("error"),
    }


@app.post("/api/leazy/chat")
async def leazy_chat(req: SimplePromptRequest, request: Request):
    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Please wait a moment.")

    message = req.message.strip()
    if message.lower().startswith("channel "):
        selection = message.split(" ", 1)[1]
        names = normalized_avatar_names(selection)
        avatar_state = set_active_avatars(names, command=message, auto=False)
        profiles = avatar_profiles(names)
        headline = " + ".join(name.title() for name in names)
        reply = (
            f"{CORE_IDENTITY} has channeled {headline}. "
            f"{' '.join(profile['power'] for profile in profiles[:2])} "
            "Your command priority is now flowing through the selected avatar protocol."
        )
        akshaya_save("vishnu_channel", headline, reply[:260], {"avatars": names, "command": message})
        return {
            "reply": reply,
            "provider": "avatar-switch",
            "used_avatars": names,
            "avatar_state": avatar_state,
        }

    state = get_state()
    viewer = session_user(request)
    workspace = normalize_workspace(req.workspace)
    source = (req.source or "orb").strip() or "orb"
    seed_context = recruitment_seed_brief(workspace, audience="member", limit=4)
    default_avatars = state["avatar_state"].get("active", ["RAMA"])
    inferred = infer_avatars_for_prompt(message, default=default_avatars)
    auto_state = set_active_avatars(inferred, command=message, auto=True)
    viewer_line = f"Viewer: {viewer['name']} ({viewer['role']})\n" if viewer else ""
    generated = await generate_text(
        (
            f"{CREATOR_ALIGNMENT_PROMPT}\n"
            f"{avatar_guidance(inferred)}\n"
            f"Workspace: {workspace}\n"
            f"Permanent recruitment core:\n{seed_context}\n"
            f"{viewer_line}"
            f"Live memory: {' | '.join(workspace_memory_brief(state, workspace)) or 'no saved context yet'}\n\n"
            f"Creator request:\n{message}"
        ),
        system=(
            f"You are {AI_NAME}, the companion orb and strategic empire brain for {COMPANY_NAME}. "
            f"Core identity: {CORE_IDENTITY}. "
            "Answer like a sharp, grounded human operator. "
            "Keep replies concise, practical, and conversational. "
            "Use short paragraphs. Do not dump system status unless it helps answer the request."
        ),
        max_tokens=420,
        use_web_search=False,
        workspace=workspace,
        source=source,
    )

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        state["conversations"].append(
            {
                "id": new_id("chat"),
                "role": "user",
                "content": message[:1200],
                "workspace": workspace,
                "created_at": now_iso(),
            }
        )
        state["conversations"].append(
            {
                "id": new_id("chat"),
                "role": "assistant",
                "content": generated["text"][:1600],
                "workspace": workspace,
                "created_at": now_iso(),
            }
        )
        return state

    mutate_state(_mutate)
    akshaya_save(
        "conversation",
        "Orb conversation",
        generated["text"][:300],
        {"prompt": message, "reply": generated["text"], "avatars": inferred, "workspace": workspace, "source": source},
    )
    return {
        "reply": generated["text"],
        "provider": generated["provider"],
        "used_avatars": inferred,
        "avatar_state": auto_state,
        "workspace": workspace,
    }


@app.post("/api/admin/edit")
async def admin_edit(req: EditRequest, request: Request):
    if not verify_admin(req.admin_token):
        raise HTTPException(status_code=403, detail="Invalid admin token")

    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded")

    edit_system = ADMIN_SYSTEM + f"""

The admin wants to edit the TechBuzz website.
Target selector hint: {req.target_selector or 'not specified'}

Current HTML context:
{req.current_html[:2000] if req.current_html else 'Not provided'}

Generate the edit instructions in <SITE_EDITS>[...]</SITE_EDITS> tags.
Be precise with CSS selectors. Explain what you changed."""

    messages = [{"role": "user", "content": req.instruction}]
    try:
        result = await call_anthropic(messages=messages, system=edit_system, max_tokens=1000, use_web_search=False)
    except HTTPException:
        raise

    text = extract_text(result)
    edits = []
    match = re.search(r"<SITE_EDITS>([\s\S]*?)</SITE_EDITS>", text)
    if match:
        try:
            edits = json.loads(match.group(1).strip())
        except json.JSONDecodeError:
            log.warning("Could not parse SITE_EDITS JSON")

    clean_text = re.sub(r"<SITE_EDITS>[\s\S]*?</SITE_EDITS>", "", text).strip()
    return {"explanation": clean_text, "edits": edits, "edit_count": len(edits)}


@app.post("/api/search")
async def web_search_endpoint(req: WebSearchRequest, request: Request):
    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded")

    if ANTHROPIC_API_KEY:
        result = await call_anthropic(
            messages=[{"role": "user", "content": f"Search the web for: {req.query}. Provide a sourced summary."}],
            system="You are a web search assistant. Search first, then answer.",
            max_tokens=800,
            use_web_search=True,
        )
        return {"query": req.query, "result": extract_text(result), "usage": result.get("usage", {})}

    generated = await generate_text(req.query, system="You are a research helper.", max_tokens=600, use_web_search=False, source="manual")
    return {"query": req.query, "result": generated["text"], "usage": generated["usage"]}


@app.get("/api/empire/dashboard")
async def empire_dashboard():
    return dashboard_payload()


@app.get("/api/brain/stream")
async def brain_stream(request: Request):
    headers = {
        "Cache-Control": "no-cache",
        "Connection": "keep-alive",
        "X-Accel-Buffering": "no",
    }
    return StreamingResponse(
        sse_json_generator(request, brain_stream_snapshot, interval_seconds=4.0),
        media_type="text/event-stream",
        headers=headers,
    )


@app.get("/api/portal/stream")
async def portal_stream(request: Request):
    headers = {
        "Cache-Control": "no-cache",
        "Connection": "keep-alive",
        "X-Accel-Buffering": "no",
    }
    return StreamingResponse(
        sse_json_generator(
            request,
            lambda: {"timestamp": now_iso(), "portal": portal_state_payload(session_user(request))},
            interval_seconds=5.0,
        ),
        media_type="text/event-stream",
        headers=headers,
    )


@app.get("/api/mother/monitor")
async def mother_monitor():
    return mother_monitor_payload()


@app.get("/api/prana-nadi/pulse")
async def prana_nadi_pulse():
    return prana_nadi_payload()


@app.get("/api/nervous-system/status")
async def nervous_system_status():
    return nervous_system_payload()


@app.get("/api/brain/hierarchy")
async def brain_hierarchy_status():
    return brain_hierarchy_payload()


@app.get("/api/brain/auto-repair/status")
async def brain_auto_repair_status(request: Request):
    user = session_user(request)
    if not user:
        return api_auth_error(401, "Login required")
    return auto_repair_engine_payload()


@app.get("/api/uiux/audit")
async def uiux_audit_status(request: Request):
    user = session_user(request)
    if not user:
        return api_auth_error(401, "Login required")
    return uiux_health_payload()


@app.post("/api/brain/auto-repair/run")
async def brain_auto_repair_run(req: BrainAutoRepairRequest, request: Request):
    user = session_user(request)
    if not user:
        return api_auth_error(401, "Login required")
    if user.get("role") != "master":
        return api_auth_error(403, "Only the master can run safe repairs")
    hierarchy = brain_hierarchy_payload()
    if req.brain_id == "all":
        targets = hierarchy.get("brains", [])
    else:
        targets = [brain for brain in hierarchy.get("brains", []) if brain.get("id") == req.brain_id]
    if not targets:
        raise HTTPException(status_code=404, detail="Brain not found")
    fixed = repair_operator_state_and_brain_memory() if req.include_state_repair else {"voice": 0, "directives": 0, "knowledge": 0}
    uiux = uiux_health_payload() if req.include_uiux_audit else {"issues": [], "health_score": 100}
    now = now_iso()
    total_fixed = sum(int(value or 0) for value in fixed.values())

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        state = ensure_state_shape(state)
        monitoring = state.setdefault("monitoring", {})
        auto_repair = monitoring.setdefault("auto_repair", {})
        auto_repair["enabled"] = True
        auto_repair["status"] = "ready" if not uiux.get("issues") else "watching"
        auto_repair["last_check_at"] = now
        auto_repair["last_repair_at"] = now if total_fixed else auto_repair.get("last_repair_at", "")
        auto_repair["safe_repairs_applied"] = int(auto_repair.get("safe_repairs_applied", 0) or 0) + total_fixed
        auto_repair["issue_count"] = len(uiux.get("issues", []))
        auto_repair["last_issues"] = list(uiux.get("issues", []))[:12]
        auto_repair["last_report"] = (
            f"Safe repair run completed for {len(targets)} brain(s). "
            f"Applied {total_fixed} repair(s) with UI/UX health {uiux.get('health_score', 100)}%."
        )
        monitoring["last_scan_at"] = now
        monitoring["last_report"] = auto_repair["last_report"]
        monitoring.setdefault("engine_reports", []).append(
            {
                "id": new_id("repair"),
                "engine": "Auto Repair Engine",
                "status": auto_repair["status"],
                "summary": auto_repair["last_report"][:320],
                "created_at": now,
            }
        )
        return state

    mutate_state(_mutate)
    updated = update_brain_directives(
        [brain["id"] for brain in targets],
        lambda brain_id, directive: {
            "last_thought": (
                f"{brain_id} completed a safe self-check. "
                f"Applied {total_fixed} repair(s) and reviewed {len(uiux.get('issues', []))} issue(s)."
            ),
            "last_thought_at": now,
            "thoughts_processed": int(directive.get("thoughts_processed", 0) or 0) + 1,
            "learning_score": min(0.99, float(directive.get("learning_score", 0.84) or 0.84) + 0.003),
            "auto_repair_status": "ready" if not uiux.get("issues") else "watching",
            "last_self_check_at": now,
            "last_safe_repair_at": now if total_fixed else directive.get("last_safe_repair_at", ""),
        },
    )
    return {
        "status": "ok",
        "targets": [brain["name"] for brain in targets],
        "updated": len(updated),
        "safe_repairs_applied": total_fixed,
        "repaired": fixed,
        "uiux": uiux,
        "result": f"Safe repair run finished for {len(updated)} brain(s).",
    }


@app.get("/api/brain/training-map")
async def brain_training_map():
    return brain_training_map_payload()


@app.post("/api/brain/assign-task")
async def brain_assign_task(req: BrainTaskRequest, request: Request):
    user = session_user(request)
    if not user:
        return api_auth_error(401, "Login required")
    if not can_control_brain(user, req.brain_id):
        return api_auth_error(403, "Only the master can control this brain")
    task = sanitize_operator_line((req.task or "").strip(), "")
    if not task:
        raise HTTPException(status_code=400, detail="Task is required and must be clear enough to use safely")
    hierarchy = brain_hierarchy_payload()
    if req.brain_id == "all":
        targets = [brain for brain in hierarchy.get("brains", []) if brain.get("id") != "mother_brain"]
    else:
        targets = [brain for brain in hierarchy.get("brains", []) if brain.get("id") == req.brain_id]
    if not targets:
        raise HTTPException(status_code=404, detail="Brain not found")
    updated = update_brain_directives(
        [brain["id"] for brain in targets],
        lambda brain_id, directive: {
            "assigned_task": task,
            "last_thought": f"{brain_id} accepted a new directive: {task}",
            "last_thought_at": now_iso(),
            "thoughts_processed": int(directive.get('thoughts_processed', 0) or 0) + 1,
            "learning_score": min(0.99, float(directive.get('learning_score', 0.84) or 0.84) + 0.005),
        },
    )
    return {
        "updated": len(updated),
        "targets": [brain["name"] for brain in targets],
        "result": f"Assigned '{task}' to {len(updated)} brain(s).",
        "context": req.context or "",
    }


@app.post("/api/brain/motivate")
async def brain_motivate(req: BrainMotivateRequest, request: Request):
    user = session_user(request)
    if not user:
        return api_auth_error(401, "Login required")
    if not can_control_brain(user, req.brain_id):
        return api_auth_error(403, "Only the master can control this brain")
    hierarchy = brain_hierarchy_payload()
    if req.brain_id == "all":
        targets = hierarchy.get("brains", [])
    else:
        targets = [brain for brain in hierarchy.get("brains", []) if brain.get("id") == req.brain_id]
    if not targets:
        raise HTTPException(status_code=404, detail="Brain not found")
    message = sanitize_operator_line((req.message or "").strip(), "")
    results = []
    for brain in targets:
        motivation = message or (
            f"Hold the role of {brain.get('role_title', brain.get('name'))} with precision. "
            f"Deliver {((brain.get('deliverables') or ['the next clear output'])[0]).lower()}."
        )
        update_brain_directives(
            [brain["id"]],
            lambda _brain_id, directive, motivation=motivation: {
                "motivation": motivation,
                "last_thought": f"{brain['name']} received renewed guidance: {motivation}",
                "last_thought_at": now_iso(),
                "thoughts_processed": int(directive.get('thoughts_processed', 0) or 0) + 1,
                "learning_score": min(0.99, float(directive.get('learning_score', 0.84) or 0.84) + 0.004),
            },
        )
        results.append({"brain_id": brain["id"], "message": motivation})
    return {"motivated": len(results), "results": results}


@app.post("/api/brain/think/{brain_id}")
async def brain_think(brain_id: str, req: BrainThinkRequest, request: Request):
    user = session_user(request)
    if not user:
        return api_auth_error(401, "Login required")
    if not can_control_brain(user, brain_id):
        return api_auth_error(403, "Only the master can control this brain")
    brain = find_brain(brain_id)
    if not brain:
        raise HTTPException(status_code=404, detail="Brain not found")
    thought = thought_for_brain(brain, req.context)
    update_brain_directives(
        [brain_id],
        lambda _brain_id, directive: {
            "last_thought": thought,
            "last_thought_at": now_iso(),
            "thoughts_processed": int(directive.get('thoughts_processed', 0) or 0) + 1,
            "learning_score": min(0.99, float(directive.get('learning_score', brain.get('learning_score', 0.84)) or 0.84) + 0.006),
            "autonomous_cycle_sec": int(directive.get('autonomous_cycle_sec') or brain.get('autonomous_cycle_sec', 30)),
        },
    )
    return {"brain": brain["name"], "thought": thought}


@app.get("/api/ops/domains")
async def ops_domains():
    return {"domains": operational_domains_payload()}


@app.get("/api/memory/audit")
async def memory_audit():
    return memory_audit_payload()


@app.get("/api/brain/status")
async def brain_status():
    return brain_payload()


@app.get("/api/settings/status")
async def settings_status():
    return settings_payload()


@app.post("/api/settings/update")
async def settings_update(req: SettingsUpdateRequest):
    updated = None

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        nonlocal updated
        settings = state["settings"]
        if req.provider_preference is not None:
            candidate = str(req.provider_preference).lower().strip()
            if candidate not in {"auto", "ollama", "anthropic", "openai", "gemini", "built_in"}:
                candidate = "auto"
            settings["provider_preference"] = candidate
        if req.always_listening is not None:
            state["voice"]["always_listening"] = bool(req.always_listening)
        if req.screen_capture_enabled is not None:
            settings["screen_capture_enabled"] = bool(req.screen_capture_enabled)
        if req.audio_capture_enabled is not None:
            settings["audio_capture_enabled"] = bool(req.audio_capture_enabled)
        if req.bounded_packages_enabled is not None:
            settings["bounded_packages_enabled"] = bool(req.bounded_packages_enabled)
        if req.privacy_guard_enabled is not None:
            settings["privacy_guard_enabled"] = bool(req.privacy_guard_enabled)
        if req.hq_visual_sync is not None:
            settings["hq_visual_sync"] = bool(req.hq_visual_sync)
        if req.external_ai_mode is not None:
            candidate_mode = str(req.external_ai_mode).strip().lower() or "manual_only"
            if candidate_mode not in {"manual_only", "built_in_only"}:
                candidate_mode = "manual_only"
            settings["external_ai_mode"] = candidate_mode
        state["voice"]["updated_at"] = now_iso()
        updated = settings_payload(state)
        return updated

    mutate_state(_mutate)
    return {"message": "Leazy settings updated.", "settings": updated}


@app.post("/api/providers/configure")
async def providers_configure(req: ProviderConfigRequest):
    provider = (req.provider or "").strip().lower()
    if provider not in {"ollama", "anthropic", "openai", "gemini"}:
        raise HTTPException(status_code=400, detail="Unsupported provider")

    model = (req.model or "").strip() or None
    api_key = (req.api_key or "").strip() if req.api_key is not None else None
    if req.clear_saved:
        api_key = ""
    persist_provider_config(provider, model=model, api_key=api_key)
    update_runtime_provider_config(provider, model=model, api_key=api_key)
    if req.clear_saved:
        clear_provider_issue(provider)
        _provider_catalog_cache.pop(provider, None)

    if req.set_default:
        def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
            state["settings"]["provider_preference"] = provider
            return state

        mutate_state(_mutate)

    if provider == "ollama":
        message = "Ollama model preference reset to the local default." if req.clear_saved else "Ollama model preference updated locally."
    else:
        message = (
            f"{provider.title()} key removed. Ishani will wait for a manual key."
            if req.clear_saved
            else f"{provider.title()} configuration saved locally."
        )
    return {"message": message, "settings": settings_payload()}


@app.post("/api/providers/catalog")
async def providers_catalog(req: ProviderCatalogRequest):
    provider = (req.provider or "").strip().lower()
    api_key = (req.api_key or "").strip() or None
    catalog = await discover_provider_models(provider, api_key=api_key)
    return {
        "message": f"{provider.title()} model catalog refreshed.",
        "provider": provider,
        "models": catalog["models"],
        "current_model": catalog["current_model"],
        "source": catalog["source"],
        "settings": settings_payload(),
    }


@app.get("/api/cabinet/status")
async def cabinet_status():
    return prime_minister_payload()


@app.post("/api/cabinet/prime-minister")
async def cabinet_prime_minister(req: PrimeMinisterRequest):
    objective = (req.objective or "").strip() or f"Generate revenue for {COMPANY_NAME} while protecting the core and improving delivery."
    command = (req.command or "").strip() or "Coordinate the secretaries, protect the core, and compound revenue."
    active_avatars = infer_avatars_for_prompt(
        f"{objective} {command}",
        default=["KRISHNA", "RAMA", "KURMA"],
    )
    set_active_avatars(active_avatars, command=command, auto=True)

    if req.enabled is not None:
        def _toggle(state: Dict[str, Any]) -> Dict[str, Any]:
            state["cabinet"]["prime_minister"]["enabled"] = bool(req.enabled)
            return state

        mutate_state(_toggle)

    result = run_prime_minister_cycle(objective=objective, command=command, source="manual")
    result["used_avatars"] = active_avatars
    return result


@app.post("/api/cabinet/toggle")
async def cabinet_toggle(req: PrimeMinisterToggleRequest):
    updated = None

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        nonlocal updated
        prime_minister = state["cabinet"]["prime_minister"]
        prime_minister["enabled"] = bool(req.enabled)
        if req.objective:
            prime_minister["objective"] = req.objective[:260]
        prime_minister["status"] = "governing" if prime_minister["enabled"] else "paused"
        interval_seconds = max(15, int(prime_minister.get("interval_seconds", 45)))
        prime_minister["next_cycle_at"] = (
            datetime.now(UTC) + timedelta(seconds=interval_seconds)
        ).isoformat() if prime_minister["enabled"] else ""
        updated = prime_minister_payload(state)
        return updated

    mutate_state(_mutate)
    return {
        "message": "Prime Minister loop enabled." if req.enabled else "Prime Minister loop paused.",
        "cabinet": updated,
    }


@app.get("/api/packages/templates")
async def package_templates():
    return {"templates": PACKAGE_TEMPLATES, "settings": settings_payload()}


@app.post("/api/packages/launch")
async def launch_package(req: PackageLaunchRequest):
    state = get_state()
    if not state["settings"].get("bounded_packages_enabled", True):
        return {
            "allowed": False,
            "message": "Bounded packages are disabled in Settings. Turn them on first, then launch the package again.",
            "package": None,
            "settings": settings_payload(state),
        }

    template = next((item for item in PACKAGE_TEMPLATES if item["id"] == req.template_id), None)
    if not template:
        raise HTTPException(status_code=404, detail="Package template not found.")

    objective = (req.objective or template["summary"]).strip()
    active_avatars = infer_avatars_for_prompt(objective, default=state["avatar_state"].get("active", ["KRISHNA", "KURMA"]))
    set_active_avatars(active_avatars, command=f"Package launch: {template['title']}", auto=True)
    prompt = (
        f"{CREATOR_ALIGNMENT_PROMPT}\n"
        f"{avatar_guidance(active_avatars)}\n\n"
        f"Package template: {template['title']}\n"
        f"Objective: {objective}\n"
        "Create a bounded mission plan with sections Mission, Inputs, Actions, Safeguards, Output. "
        "Do not describe self-propagation, intrusion, malware behavior, or bypassing consent."
    )
    generated = await generate_text(
        prompt,
        system=f"You are {AI_NAME}. Build safe, bounded missions for company growth and research.",
        max_tokens=520,
        use_web_search=False,
        source="package",
    )
    record = {
        "id": new_id("pkg"),
        "template_id": template["id"],
        "title": template["title"],
        "objective": objective[:400],
        "report": generated["text"][:2200],
        "provider": generated["provider"],
        "avatars": active_avatars,
        "created_at": now_iso(),
    }

    def _mutate(state_obj: Dict[str, Any]) -> Dict[str, Any]:
        state_obj["packages"].append(record)
        state_obj["vault"].append(
            {
                "id": new_id("vault"),
                "kind": "bounded_package",
                "title": template["title"],
                "summary": objective[:220],
                "content": generated["text"][:1800],
                "created_at": now_iso(),
            }
        )
        return record

    mutate_state(_mutate)
    cabinet = run_prime_minister_cycle(
        objective=f"Generate revenue for {COMPANY_NAME} from active packages and delivery lanes.",
        command=f"Absorb package {template['title']} into the revenue cabinet.",
        source="package",
    )
    return {"allowed": True, "message": f"{template['title']} launched.", "package": record, "cabinet": cabinet["cabinet"]}


@app.get("/api/portal/state")
async def portal_state(request: Request):
    return portal_state_payload(session_user(request))


@app.post("/api/network/scan")
async def network_scan(req: NetworkScanRequest):
    query = (req.query or "").strip()
    if not query:
        raise HTTPException(status_code=400, detail="Query is required for the signal scanner.")
    generated = await generate_text(
        (
            f"{public_hq_context_brief()}\n\n"
            f"Carbon Protocol scanner query:\n{query}\n\n"
            "Return a compact market signal brief with sections Signals, Risks, Opportunities, and Next Move."
        ),
        system=(
            f"You are {AI_NAME}, running the Carbon Signal Scanner for {COMPANY_NAME}. "
            "Read the market, synthesize the movement, and produce a precise operator brief."
        ),
        max_tokens=520,
        use_web_search=False,
        workspace="network",
        source="network_scan",
    )
    summary = generated["text"][:2200]
    akshaya_save("network_scan", f"Network Scan: {query[:80]}", summary[:260], {"query": query, "provider": generated["provider"]})
    return {"message": "Network scan complete.", "report": summary, "provider": generated["provider"]}


@app.post("/api/ats/import-latest")
async def ats_import_latest():
    state = get_state()
    hunts = state.get("praapti_hunts", [])
    if not hunts:
        raise HTTPException(status_code=404, detail="No Praapti hunt is available to import yet.")
    latest = hunts[-1]
    summary = (
        f"Imported {len(latest.get('candidates', []))} candidate(s) from the latest Praapti hunt for "
        f"{latest.get('client_company', COMPANY_NAME)} into the ATS pipeline."
    )
    akshaya_save("ats_import", f"ATS Import: {latest.get('client_company', COMPANY_NAME)}", summary, {"hunt_id": latest.get("id", ""), "provider": latest.get("provider", "built-in")})
    return {
        "message": summary,
        "job": {
            "client_company": latest.get("client_company", COMPANY_NAME),
            "summary": latest.get("job_description", "")[:180],
            "urgency": latest.get("urgency", "high"),
        },
        "candidates": latest.get("candidates", []),
    }


@app.post("/api/brain/pulse")
async def brain_pulse(req: BrainPulseRequest):
    focus = (req.focus or "all").strip() or "all"
    goal = (req.goal or "advance the empire").strip() or "advance the empire"
    state = get_state()
    active_avatars = infer_avatars_for_prompt(
        f"{focus} {goal}",
        default=state["avatar_state"].get("active", ["RAMA"]),
    )
    set_active_avatars(active_avatars, command=f"brain pulse: {focus} | {goal}", auto=True)
    current_brain = brain_payload()
    prompt = (
        f"{CREATOR_ALIGNMENT_PROMPT}\n"
        f"{avatar_guidance(active_avatars)}\n\n"
        f"Focus: {focus}\n"
        f"Goal: {goal}\n"
        f"Brain pillars: {json.dumps(current_brain['pillars'])}\n"
        "Write a concise operational pulse with exactly five sections:\n"
        "Expand:\nGrow:\nDevelop:\nProtect:\nImmediate Move:\n"
    )
    generated = await generate_text(
        prompt,
        system=(
            f"You are {AI_NAME}, a mature strategic operating brain for {COMPANY_NAME}. "
            f"Core identity: {CORE_IDENTITY}. Keep the report vivid, direct, and actionable."
        ),
        max_tokens=650,
        use_web_search=False,
        source="manual",
    )
    report = generated["text"][:2200]
    akshaya_save(
        "brain_pulse",
        f"Brain Pulse: {focus.title()}",
        report[:320],
        {"focus": focus, "goal": goal, "avatars": active_avatars, "provider": generated["provider"]},
    )
    return {
        "focus": focus,
        "goal": goal,
        "report": report,
        "provider": generated["provider"],
        "used_avatars": active_avatars,
        "brain": brain_payload(),
    }


@app.post("/api/praapti/hunt")
async def praapti_hunt(req: PraaptiHuntRequest, request: Request):
    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded")

    active_avatars = infer_avatars_for_prompt(req.job_description + " recruitment hiring talent", default=["KRISHNA", "RAMA"])
    set_active_avatars(active_avatars, command=f"Praapti hunt: {req.job_description[:120]}", auto=True)
    seed_context = recruitment_seed_brief("praapti", audience="member", limit=5)
    culture_prompt = (
        "You are Manas-Pravah, the culture simulation layer of the King's Empire.\n"
        f"Job Description: {req.job_description}\n"
        f"Client: {req.client_company} | Urgency: {req.urgency}\n"
        f"Permanent recruitment core:\n{seed_context}\n"
        f"{avatar_guidance(active_avatars)}\n"
        "Simulate the hidden company culture, team dynamics, unspoken needs, and the true personality fit."
    )
    profile_prompt = (
        "You are Agni-Shala, the fire forge of the empire.\n"
        f"JD: {req.job_description}\n"
        f"Permanent recruitment core:\n{seed_context}\n"
        f"{avatar_guidance(active_avatars)}\n"
        "Create the ideal candidate profile with skills, experience patterns, personality markers, and risks."
    )
    hunt_prompt = (
        "You are Praapti, the recruitment agency of the King's Empire.\n"
        f"JD: {req.job_description}\n"
        f"Permanent recruitment core:\n{seed_context}\n"
        f"{avatar_guidance(active_avatars)}\n"
        "Return EXACTLY 3 candidates in JSON format:\n"
        '[{"name":"...", "title":"...", "experience":8, "fit_score":96, "genesis_profile":"...", "discovery_source":"..."}]'
    )

    culture = await generate_text(culture_prompt, system=ADMIN_SYSTEM, max_tokens=450, use_web_search=False, source="praapti")
    profile = await generate_text(profile_prompt, system=ADMIN_SYSTEM, max_tokens=450, use_web_search=False, source="praapti")
    hunt = await generate_text(hunt_prompt, system=ADMIN_SYSTEM, max_tokens=700, use_web_search=req.live_search, source="praapti")

    candidates = parse_json_blob(hunt["text"])
    if not isinstance(candidates, list):
        candidates = fallback_candidates(3)

    record = {
        "id": new_id("hunt"),
        "job_description": req.job_description[:3000],
        "client_company": req.client_company,
        "urgency": req.urgency,
        "culture_insight": culture["text"][:1200],
        "ideal_profile": profile["text"][:1200],
        "candidates": candidates[:3],
        "provider": hunt["provider"],
        "avatars": active_avatars,
        "created_at": now_iso(),
        "created_day": today_iso(),
    }

    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        state["praapti_hunts"].append(record)
        state["vault"].append(
            {
                "id": new_id("vault"),
                "kind": "praapti_hunt",
                "title": f"Praapti: {req.client_company}",
                "summary": f"{len(record['candidates'])} candidates for {req.job_description[:80]}",
                "content": json.dumps(record["candidates"], ensure_ascii=False)[:2400],
                "created_at": now_iso(),
            }
        )
        return record

    mutate_state(_mutate)
    cabinet = run_prime_minister_cycle(
        objective=f"Generate revenue for {COMPANY_NAME} by converting hiring demand into delivery and client trust.",
        command=f"Brief the cabinet on the new Praapti hunt for {req.client_company}.",
        source="praapti",
    )
    return {
        "candidates": record["candidates"],
        "culture_insight": record["culture_insight"],
        "ideal_profile": record["ideal_profile"],
        "message": "Praapti hunt complete. The result is preserved in the Akshaya Quantum Vault.",
        "provider": record["provider"],
        "used_avatars": active_avatars,
        "cabinet": cabinet["cabinet"],
    }


@app.get("/api/praapti/hunts")
async def praapti_hunts():
    state = get_state()
    hunts = list(reversed(state["praapti_hunts"]))[:20]
    return {"hunts": hunts}


@app.post("/api/swarm/mission")
async def swarm_mission(req: SwarmMissionRequest, request: Request):
    ip = get_client_ip(request)
    if not check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded")

    state = get_state()
    active_avatars = infer_avatars_for_prompt(req.mission, default=["KRISHNA", "KALKI"])
    set_active_avatars(active_avatars, command=req.mission, auto=True)
    prompt = (
        "You are the Swarm Intelligence of the King's Empire.\n"
        f"Mission: {req.mission}\n"
        f"Recent hunts: {len(state['praapti_hunts'])}\n"
        f"Pending proposals: {sum(1 for row in state['nirmaan_proposals'] if not row.get('approved'))}\n"
        f"{avatar_guidance(active_avatars)}\n"
        "Coordinate Praapti, Revenue, Army, Navy, Airforce, Research, Armory, Warehouse, Maya, Anveshan, "
        "Sankalpa, and Akshaya into a short mission report with specific outcomes."
    )
    generated = await generate_text(prompt, system=ADMIN_SYSTEM, max_tokens=850, use_web_search=False, source="swarm")
    proposal = None
    if any(word in req.mission.lower() for word in ("evolve", "improve", "develop", "upgrade", "autonomous")):
        proposal = await create_nirmaan_proposal(req.mission)

    record = {
        "id": new_id("swarm"),
        "mission": req.mission[:200],
        "report": generated["text"][:4000],
        "provider": generated["provider"],
        "created_at": now_iso(),
        "avatars": active_avatars,
    }

    def _mutate(state_obj: Dict[str, Any]) -> Dict[str, Any]:
        state_obj["swarm_missions"].append(record)
        state_obj["vault"].append(
            {
                "id": new_id("vault"),
                "kind": "swarm_mission",
                "title": req.mission[:120],
                "summary": generated["text"][:280],
                "content": generated["text"][:1800],
                "created_at": now_iso(),
            }
        )
        return record

    mutate_state(_mutate)
    cabinet = run_prime_minister_cycle(
        objective=f"Generate revenue for {COMPANY_NAME} while using the swarm to strengthen execution and delivery.",
        command=f"Absorb swarm mission into the Prime Minister queue: {req.mission[:120]}",
        source="swarm",
    )
    return {
        "report": record["report"],
        "agents_coordinated": len(SWARM_AGENT_NAMES),
        "status": "Swarm mission completed successfully",
        "proposal": proposal,
        "used_avatars": active_avatars,
        "cabinet": cabinet["cabinet"],
    }


@app.get("/api/nirmaan/proposals")
async def get_nirmaan_proposals():
    state = get_state()
    proposals = [row for row in dedupe_nirmaan_proposals(state["nirmaan_proposals"]) if not row.get("approved")]
    proposals.sort(key=lambda row: row.get("created_at", ""), reverse=True)
    return {"proposals": proposals[:8]}


@app.post("/api/nirmaan/develop")
async def nirmaan_develop():
    proposal = await create_nirmaan_proposal("manual_develop")
    cabinet = run_prime_minister_cycle(
        objective=f"Generate revenue for {COMPANY_NAME} by evolving the product and operating system.",
        command=f"Fold Nirmaan proposal {proposal['title']} into the cabinet roadmap.",
        source="nirmaan",
    )
    return {
        "message": f"Nirmaan Chakra proposed a new evolution: {proposal['title']}",
        "proposal_id": proposal["id"],
        "proposal": proposal,
        "cabinet": cabinet["cabinet"],
    }


@app.post("/api/nirmaan/approve")
async def nirmaan_approve(req: NirmaanApproveRequest):
    approved = None

    def _mutate(state: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        nonlocal approved
        state["nirmaan_proposals"] = dedupe_nirmaan_proposals(state["nirmaan_proposals"])
        for proposal in state["nirmaan_proposals"]:
            if proposal["id"] == req.proposal_id:
                proposal["approved"] = True
                proposal["approved_at"] = now_iso()
                approved = proposal
                break
        if approved:
            state["vault"].append(
                {
                    "id": new_id("vault"),
                    "kind": "nirmaan_approved",
                    "title": approved["title"],
                    "summary": "Proposal approved and staged for merge.",
                    "content": approved["code_snippet"][:1800],
                    "created_at": now_iso(),
                }
            )
        return approved

    mutate_state(_mutate)
    if not approved:
        raise HTTPException(status_code=404, detail="Proposal not found")

    output_path = PROPOSAL_DIR / f"{approved['id']}.txt"
    output_path.write_text(
        (
            f"TITLE: {approved['title']}\n"
            f"CREATED_AT: {approved.get('created_at')}\n"
            f"APPROVED_AT: {approved.get('approved_at')}\n\n"
            f"DESCRIPTION:\n{approved['description']}\n\n"
            f"CODE_SNIPPET:\n{approved['code_snippet']}\n"
        ),
        encoding="utf-8",
    )
    return {
        "message": f"Feature '{approved['title']}' approved and staged in {output_path.name}. The empire has grown its roadmap.",
        "feature": approved["title"],
        "path": str(output_path),
    }


@app.get("/api/akshaya/vault")
async def akshaya_vault():
    state = get_state()
    items = list(reversed(state["vault"]))[:40]
    cabinet = prime_minister_payload(state)
    return {
        "items": items,
        "guardian": state["meta"].get("memory_guardian", "eternal preservation"),
        "eternal_preservation_mode": "Matsya + Kurma",
        "protection_meter": state["avatar_state"].get("protection_meter", 0),
        "state_file": str(STATE_PATH),
        "seal": "Matsya + Kurma",
        "quantum_storage": cabinet["quantum_storage"],
        "prime_minister": cabinet["prime_minister"],
    }


@app.get("/api/voice/status")
async def voice_status():
    return voice_status_payload()


@app.post("/api/voice/settings")
async def voice_settings(req: VoiceSettingsRequest):
    def _mutate(state: Dict[str, Any]) -> Dict[str, Any]:
        if req.always_listening is not None:
            state["voice"]["always_listening"] = bool(req.always_listening)
        if req.wake_words:
            state["voice"]["wake_words"] = [word.strip().lower() for word in req.wake_words if word.strip()]
        if req.voice_profile is not None:
            profile = str(req.voice_profile).strip().lower() or "sovereign_female"
            if profile not in VOICE_PROFILE_PRESETS:
                profile = "sovereign_female"
            preset = VOICE_PROFILE_PRESETS[profile]
            state["voice"]["voice_profile"] = profile
            state["voice"]["language"] = str(req.language or preset["language"]).strip() or preset["language"]
            state["voice"]["rate"] = max(0.7, min(1.2, float(req.rate if req.rate is not None else preset["rate"])))
            state["voice"]["pitch"] = max(0.7, min(1.4, float(req.pitch if req.pitch is not None else preset["pitch"])))
            state["voice"]["engine"] = "browser_builtin_female"
        else:
            if req.language is not None:
                state["voice"]["language"] = str(req.language).strip() or state["voice"].get("language", "en-IN")
            if req.rate is not None:
                state["voice"]["rate"] = max(0.7, min(1.2, float(req.rate)))
            if req.pitch is not None:
                state["voice"]["pitch"] = max(0.7, min(1.4, float(req.pitch)))
        state["voice"]["updated_at"] = now_iso()
        return state["voice"]

    voice = mutate_state(_mutate)
    return {"message": "Voice wake mode updated.", "voice": voice}


async def execute_voice_command(command: str, mode: str = "wake") -> Dict[str, Any]:
    state = get_state()
    wake_words = state["voice"].get("wake_words", [])
    heard = (command or "").strip()
    normalized_mode = str(mode or "wake").strip().lower() or "wake"
    direct_mode = normalized_mode in {"direct", "manual", "command"}
    heard_for_storage = sanitize_operator_line(heard[:200], "")
    heard_for_model = sanitize_operator_multiline(
        heard,
        "The operator request was unclear. Ask for a clean instruction.",
    )
    wake_detected = direct_mode or any(word in heard.lower() for word in wake_words)
    active_avatars = state["avatar_state"].get("active", ["RAMA"])
    if wake_detected:
        if "channel " in heard.lower():
            selection = heard.lower().split("channel ", 1)[1]
            active_avatars = normalized_avatar_names(selection)
            set_active_avatars(active_avatars, command=heard, auto=False)
            profiles = avatar_profiles(active_avatars)
            reply = (
                f"{CORE_IDENTITY} has channeled {' + '.join(name.title() for name in active_avatars)}. "
                f"{profiles[0]['mantra'] if profiles else 'The wheel is now active.'}"
            )
            generated = {"provider": "voice-avatar"}
        else:
            active_avatars = infer_avatars_for_prompt(heard, default=active_avatars)
            set_active_avatars(active_avatars, command=heard, auto=True)
            guidance = avatar_guidance(active_avatars)
            translation = None
            if INTERPRETER_LAYER.get("translate_command"):
                try:
                    translation = await INTERPRETER_LAYER["translate_command"](
                        heard_for_model,
                        context=f"Voice relay with avatars {' + '.join(active_avatars)}. {guidance}",
                        source="voice",
                        prefer_external=False,
                    )
                except Exception as exc:
                    log.warning("Voice interpreter fallback: %s", exc)
            if translation:
                target_brain_id = translation.get("target_brain_id")
                if target_brain_id and find_brain(target_brain_id):
                    translated_task = sanitize_operator_multiline(translation.get("translated_task", heard_for_model), heard_for_model)
                    update_brain_directives(
                        [target_brain_id],
                        lambda _brain_id, directive, translated_task=translated_task: {
                            "last_thought": f"Interpreter relay: {translated_task}",
                            "last_thought_at": now_iso(),
                            "thoughts_processed": int(directive.get("thoughts_processed", 0) or 0) + 1,
                            "learning_score": min(0.99, float(directive.get("learning_score", 0.84) or 0.84) + 0.003),
                        },
                    )
                generated = {"provider": " -> ".join(translation.get("provider_chain", [])) or "ollama/local"}
                reply = translation.get("operator_reply") or "I understood the command and routed it into the right brain lane."
            else:
                generated = await generate_text(
                    f"{CREATOR_ALIGNMENT_PROMPT}\n{guidance}\n\nVoice request:\n{heard_for_model}",
                    system=(
                        f"You are {AI_NAME}. Respond like a live voice assistant after a wake phrase. "
                        "Use short spoken sentences with a calm feminine tone, clear phrasing, and direct operational guidance."
                    ),
                    max_tokens=250,
                    use_web_search=False,
                    source="voice",
                )
                reply = generated["text"]
    else:
        generated = {"provider": "built-in"}
        reply = "Wake phrase not detected yet. Say 'Hey Jinn' or 'My King commands'."

    def _mutate(state_obj: Dict[str, Any]) -> Dict[str, Any]:
        state_obj["voice"]["last_command"] = heard_for_storage
        state_obj["voice"]["updated_at"] = now_iso()
        return state_obj["voice"]

    mutate_state(_mutate)
    akshaya_save("voice_command", "Voice Wake", reply[:220], {"heard": heard_for_storage, "wake_detected": wake_detected, "avatars": active_avatars})
    return {
        "wake_detected": wake_detected,
        "heard": heard_for_storage or "Voice relay captured.",
        "response": reply,
        "provider": generated["provider"],
        "used_avatars": active_avatars,
        "mode": normalized_mode,
    }


@app.post("/api/voice/wake")
async def voice_wake(req: VoiceWakeRequest):
    return await execute_voice_command(req.command, req.mode)


@app.get("/api/vishnu/status")
async def vishnu_status():
    state = get_state()
    active = state["avatar_state"].get("active", ["RAMA"])
    return {
        "identity": state["meta"].get("identity", CORE_IDENTITY),
        "creator_mode": state["meta"].get("creator_mode", CREATOR_MODE_LABEL),
        "active_avatars": avatar_profiles(active),
        "protection_meter": state["avatar_state"].get("protection_meter", 0),
        "history": list(reversed(state["avatar_state"].get("history", [])))[:12],
    }


@app.post("/api/vishnu/channel")
async def vishnu_channel(req: VishnuChannelRequest):
    names = normalized_avatar_names(req.avatar)
    profiles = avatar_profiles(names)
    avatar_state = set_active_avatars(names, command=req.command or req.avatar, auto=False)
    prompt = (
        f"{CREATOR_ALIGNMENT_PROMPT}\n"
        f"{avatar_guidance(names)}\n\n"
        f"Creator command: {req.command}\n"
        "Respond with a powerful, concise confirmation of the newly channeled mode."
    )
    generated = await generate_text(
        prompt,
        system=(
            f"You are {AI_NAME} with {CORE_IDENTITY} core consciousness. "
            "You are channeling the selected Vishnu avatars for the creator."
        ),
        max_tokens=320,
        use_web_search=False,
        source="manual",
    )
    reply = generated["text"] or (
        f"{CORE_IDENTITY} now channels {' + '.join(name.title() for name in names)}. "
        "The empire is aligned and protected."
    )
    akshaya_save(
        "vishnu_channel",
        f"Dashavatara: {' + '.join(name.title() for name in names)}",
        reply[:280],
        {"avatars": names, "command": req.command, "provider": generated["provider"]},
    )
    return {
        "avatars": names,
        "profiles": profiles,
        "message": reply,
        "seal": "Vishnu protocol is active across the empire.",
        "protection_meter": avatar_state.get("protection_meter", 0),
        "creator_mode": CREATOR_MODE_LABEL,
    }


@app.get("/api/stats")
async def get_stats(x_admin_token: Optional[str] = Header(None)):
    if not x_admin_token or not verify_admin(x_admin_token):
        raise HTTPException(status_code=403, detail="Admin access required")
    total_requests = sum(len(calls) for calls in _rate_store.values())
    empire = dashboard_payload()
    return {
        "total_api_calls_last_minute": total_requests,
        "active_ips": len(_rate_store),
        "model": MODEL,
        "web_search_enabled": bool(ANTHROPIC_API_KEY),
        "timestamp": now_iso(),
        "empire_metrics": empire["metrics"],
    }


install_empire_merge_layer(
    app,
    {
        "db_exec": db_exec,
        "db_all": db_all,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "generate_text": generate_text,
        "get_state": get_state,
        "FRONTEND_DIR": FRONTEND_DIR,
        "AI_NAME": AI_NAME,
        "COMPANY_NAME": COMPANY_NAME,
        "CORE_IDENTITY": CORE_IDENTITY,
        "log": log,
    },
)

GLOBAL_RECRUITMENT_LAYER = install_global_recruitment_brain_layer(
    app,
    {
        "db_exec": db_exec,
        "db_one": db_one,
        "db_all": db_all,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "brain_hierarchy_payload": brain_hierarchy_payload,
    },
)

RECRUITMENT_LAYER = install_recruitment_brain_layer(
    app,
    {
        "db_exec": db_exec,
        "db_one": db_one,
        "db_all": db_all,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "generate_text": generate_text,
        "get_state": get_state,
        "portal_state_payload": portal_state_payload,
        "DATA_DIR": DATA_DIR,
        "DOCUMENT_DIR": DOCUMENT_DIR,
        "AI_NAME": AI_NAME,
        "COMPANY_NAME": COMPANY_NAME,
        "CORE_IDENTITY": CORE_IDENTITY,
        "world_brain_context": GLOBAL_RECRUITMENT_LAYER.get("context_brief"),
        "world_brain_status": GLOBAL_RECRUITMENT_LAYER.get("status_payload"),
    },
)

BROWSER_SUITE_LAYER = install_browser_suite_layer(
    app,
    {
        "db_exec": db_exec,
        "db_one": db_one,
        "db_all": db_all,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "generate_text": generate_text,
        "get_state": get_state,
        "brain_hierarchy_payload": brain_hierarchy_payload,
        "nervous_system_payload": nervous_system_payload,
        "mother_monitor_payload": mother_monitor_payload,
        "portal_state_payload": portal_state_payload,
        "normalize_navigator_url": normalize_navigator_url,
        "broadcast_brain_lesson": broadcast_brain_lesson,
        "akshaya_save": akshaya_save,
        "AI_NAME": AI_NAME,
        "COMPANY_NAME": COMPANY_NAME,
        "CORE_IDENTITY": CORE_IDENTITY,
        "log": log,
    },
)

INTERPRETER_LAYER = install_interpreter_brain_layer(
    app,
    {
        "db_exec": db_exec,
        "db_one": db_one,
        "db_all": db_all,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "can_control_brain": can_control_brain,
        "find_brain": find_brain,
        "brain_hierarchy_payload": brain_hierarchy_payload,
        "broadcast_brain_lesson": broadcast_brain_lesson,
        "sanitize_operator_line": sanitize_operator_line,
        "sanitize_operator_multiline": sanitize_operator_multiline,
        "parse_json_blob": parse_json_blob,
        "get_state": get_state,
        "call_local_llm": call_local_llm,
        "ollama_status": ollama_status_payload,
        "generate_text": generate_text,
        "provider_config": provider_config,
        "active_provider_label": active_provider_label,
        "external_ai_allowed_for_source": external_ai_allowed_for_source,
        "AI_NAME": AI_NAME,
        "CORE_IDENTITY": CORE_IDENTITY,
        "log": log,
    },
)

install_brain_communication_layer(
    app,
    {
        "db_exec": db_exec,
        "db_all": db_all,
        "db_one": db_one,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "find_brain": find_brain,
        "brain_hierarchy_payload": brain_hierarchy_payload,
        "log": log,
    },
)

ORCHESTRATION_STACK_LAYER = install_orchestration_stack_layer(
    app,
    {
        "db_exec": db_exec,
        "db_all": db_all,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "can_control_brain": can_control_brain,
        "find_brain": find_brain,
        "brain_hierarchy_payload": brain_hierarchy_payload,
        "sanitize_operator_line": sanitize_operator_line,
        "sanitize_operator_multiline": sanitize_operator_multiline,
        "call_local_llm": call_local_llm,
        "generate_text": generate_text,
        "provider_config": provider_config,
        "active_provider_label": active_provider_label,
        "ollama_status": ollama_status_payload,
        "AI_NAME": AI_NAME,
        "CORE_IDENTITY": CORE_IDENTITY,
        "log": log,
    },
)

LOCAL_AI_STACK_LAYER = install_local_ai_runtime_layer(
    app,
    {
        "db_exec": db_exec,
        "db_all": db_all,
        "new_id": new_id,
        "now_iso": now_iso,
        "session_user": session_user,
        "mutate_state": mutate_state,
        "get_state": get_state,
        "extract_document_text": extract_document_text,
        "call_ollama": call_ollama,
        "extract_ollama_text": extract_ollama_text,
        "generate_text": generate_text,
        "sanitize_operator_line": sanitize_operator_line,
        "sanitize_operator_multiline": sanitize_operator_multiline,
        "BACKEND_DIR": BACKEND_DIR,
        "DATA_DIR": DATA_DIR,
        "OLLAMA_HOST": OLLAMA_HOST,
        "OLLAMA_MODEL": OLLAMA_MODEL,
        "AI_NAME": AI_NAME,
        "CORE_IDENTITY": CORE_IDENTITY,
        "log": log,
    },
)

VOICE_RUNTIME_LAYER = install_voice_runtime_layer(
    app,
    {
        "session_user": session_user,
        "mutate_state": mutate_state,
        "get_state": get_state,
        "execute_voice_command": execute_voice_command,
        "now_iso": now_iso,
        "BACKEND_DIR": BACKEND_DIR,
        "DATA_DIR": DATA_DIR,
        "log": log,
    },
)

warm_ollama_background()

register_recruiter_status_routes(
    app,
    db_all=db_all,
    db_one=db_one,
    db_exec=db_exec,
    new_id=new_id,
    now_iso=now_iso,
    session_user=session_user,
    generate_text=generate_text,
    log=log,
)

try:
    seed_market_ready_brains()
except Exception as exc:
    log.warning("Unable to seed market-ready brains: %s", exc)


if __name__ == "__main__":
    import uvicorn

    print(
        f"""
==============================================
  TechBuzz Leazy Jinn Backend
  Provider: {active_provider_label()}
  Open: http://localhost:8000/leazy
==============================================
"""
    )
    if not any((ANTHROPIC_API_KEY, OPENAI_API_KEY, GEMINI_API_KEY, ollama_provider_ready())):
        print("WARNING: No external AI keys are set. Leazy Jinn will use the built-in fallback brain.")
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
